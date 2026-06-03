"""
Build MP-Competitive-Analysis.pptx following the `competitive-analysis` skill in
plugins/vertical-plugins/financial-analysis. 18 slides, insight-driven titles,
2x2 positioning matrix, peer deep dives, moats grid, bull/base/bear, catalysts.

Run:
    python3 build_deck.py
Output:
    MP-Competitive-Analysis.pptx
    charts/  (intermediate PNGs)
"""
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"
CHARTS.mkdir(exist_ok=True)
OUT = HERE / "MP-Competitive-Analysis.pptx"

# Palette per skill (2-3 colors max, muted)
NAVY    = RGBColor(0x17, 0x36, 0x5D)
COPPER  = RGBColor(0xB2, 0x6B, 0x1E)   # rare-earth accent
CHARCOAL= RGBColor(0x40, 0x40, 0x40)
GREY    = RGBColor(0xBF, 0xBF, 0xBF)
LGREY   = RGBColor(0xF2, 0xF2, 0xF2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Times New Roman"

# ============================================================
# CHARTS (matplotlib → PNG)
# ============================================================
plt.rcParams.update({
    "font.family": "serif",
    "font.serif": ["DejaVu Serif"],
    "axes.edgecolor": "#404040",
    "axes.labelcolor": "#404040",
    "xtick.color": "#404040",
    "ytick.color": "#404040",
    "axes.grid": True,
    "grid.color": "#E0E0E0",
    "grid.linestyle": "-",
    "grid.linewidth": 0.5,
})

def chart_ndpr():
    fig, ax = plt.subplots(figsize=(9, 4.5), dpi=150)
    months = ["Jan'26","Feb","Mar","Apr","May'26"]
    spot   = [53, 68, 85, 105, 100]
    floor  = [110]*5
    ax.fill_between(months, spot, color="#B26B1E", alpha=0.18)
    ax.plot(months, spot, color="#B26B1E", linewidth=2.5, marker="o", label="NdPr oxide spot ($/kg)")
    ax.plot(months, floor, color="#2E7D32", linewidth=2, linestyle="--", label="MP / DoD price floor ($110/kg)")
    ax.set_ylabel("USD per kg")
    ax.set_title("NdPr oxide spot price — Jan to May 2026", fontsize=13, color="#17365D", pad=10)
    ax.legend(loc="lower right", frameon=False)
    ax.set_ylim(40, 130)
    fig.tight_layout()
    out = CHARTS / "ndpr.png"
    fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out

def chart_2x2():
    # Axes: x = Revenue scale ($M, log), y = Integration depth (1-6 stages reached)
    fig, ax = plt.subplots(figsize=(9.5, 6), dpi=150)
    # data: (rev, integration_score, name, color, size)
    data = [
        (275,  4.0, "MP Materials",  "#B26B1E", 800),  # mine, crack, sep, magnet planned (3 done + magnet planned = 4)
        (578,  3.5, "Lynas",         "#17365D", 800),  # mine, crack, sep, heavy + (no magnet) = 3.5
        (70,   1.5, "Energy Fuels",  "#7F8C8D", 300),
        (5,    2.0, "USA Rare Earth","#7F8C8D", 300),  # magnet only — partial
        (5910, 6.0, "China N. RE",   "#C0392B", 1300),
        (5850, 0.5, "Albemarle",     "#95A5A6", 700),  # adjacent
    ]
    for rev, integ, name, col, size in data:
        ax.scatter(rev, integ, s=size, color=col, alpha=0.65, edgecolors="#202020", linewidths=1.2)
        ax.annotate(name, (rev, integ), xytext=(10, 8), textcoords="offset points",
                    fontsize=10, color="#202020", weight="bold")
    ax.set_xscale("log")
    ax.set_xlabel("Revenue scale, USD millions (TTM, log)")
    ax.set_ylabel("Vertical integration depth (stages reached)")
    ax.set_yticks([0,1,2,3,4,5,6])
    ax.set_yticklabels(["Adjacent","—","Processing","Sep","Mine+Sep+HRE","Sep + Magnet planned","Fully integrated"])
    ax.set_xlim(2, 12000)
    ax.set_ylim(-0.5, 6.7)
    ax.set_title("Positioning: scale × vertical integration",
                 fontsize=13, color="#17365D", pad=12)
    # quadrant shading
    ax.axhline(3.5, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.axvline(500, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.text(2.5,   6.3, "Integrated juniors\n(rare)",  fontsize=9, color="#7F8C8D")
    ax.text(2500,  6.3, "Integrated leaders",          fontsize=9, color="#7F8C8D")
    ax.text(2.5,   0.0, "Niche / adjacent",            fontsize=9, color="#7F8C8D")
    ax.text(2500,  0.0, "Scaled adjacents",            fontsize=9, color="#7F8C8D")
    fig.tight_layout()
    out = CHARTS / "positioning.png"
    fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out

NDPR_PNG = chart_ndpr()
MATRIX_PNG = chart_2x2()

# ============================================================
# DECK BUILD
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(45720)  # ~0.05in
    tf.vertical_anchor = anchor
    # First paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs = [text]
    else:
        runs = text
    for i, line in enumerate(runs):
        if i == 0:
            r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align
            r = p2.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def slide_title(slide, title, *, size=26):
    # Top color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = COPPER
    bar.line.fill.background()
    # Title text
    add_text(slide, Inches(0.5), Inches(0.20), Inches(12.3), Inches(0.9),
             title, size=size, bold=True, color=NAVY)

def slide_footer(slide, source_text):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             source_text, size=9, italic=True, color=GREY)

def add_table(slide, left, top, width, height, data, *, header=True, col_widths=None,
              first_col_bold=True, font_size=12):
    """data: list of rows, each row is list of strings."""
    rows = len(data); cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            tf = cell.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    if header and r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    elif first_col_bold and c_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = NAVY
                    else:
                        run.font.color.rgb = CHARCOAL
            if header and r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif r_idx % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Emu(45720)
            cell.margin_top = cell.margin_bottom = Emu(36576)
    return tbl_shape

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
# Big navy block top half
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.5))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
# Copper accent stripe
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.12))
stripe.fill.solid(); stripe.fill.fore_color.rgb = COPPER; stripe.line.fill.background()

add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
         "NYSE: MP — Investment Decision Frame", size=18, italic=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.4),
         "MP Materials", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.9),
         "Competitive Landscape & Investment Read", size=28, color=COPPER)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "As of May 2026", size=14, italic=True, color=CHARCOAL)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
         ["Peer set: Lynas · Energy Fuels · USA Rare Earth · China Northern Rare Earth · Albemarle",
          "Built using `competitive-analysis` + `comps-analysis` skills (financial-services repo)",
          "Companion workbook: MP-Comps-Analysis.xlsx"],
         size=13, color=CHARCOAL)

# ============================================================
# SLIDE 2 — EXEC SUMMARY
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Only scaled US rare-earth play; DoD as anchor shareholder rewrites the risk profile")

bullets = [
    ("1. The setup",
     "MP is the only fully scaled US rare-earth miner and the only Western producer with end-to-end mine-to-magnet ambition. "
     "Q1 2026 marked an inflection: revenue $90.6M (+49% YoY), record 917 MT NdPr oxide (+63%), Adj EBITDA flipped positive to $36.6M."),
    ("2. The moat",
     "US Department of Defense holds ~15% via preferred + warrant at $30.03/sh (Jul 2025 deal). The 10-year DoD offtake includes a "
     "$110/kg NdPr price floor — structural downside protection no peer has. MP is the only US public-company peer with this asymmetry."),
    ("3. The risk",
     "Market cap ~$9.8B on TTM revenue $275M = priced for the magnet ramp. Stage III first commercial product is 2028. "
     "China still controls >70% of mining, >80% of processing — NdPr is exposed to Chinese export policy."),
    ("4. The read",
     "Long thesis intact for investors who accept multi-year execution risk in exchange for asymmetric upside on Section 232 / "
     "decoupling tailwinds and a 2028 magnet revenue inflection. Pair-trade vs CNREG offers tighter risk profile."),
]
y = Inches(1.4)
for hdr, body in bullets:
    add_text(s, Inches(0.7), y, Inches(2.3), Inches(1.2), hdr, size=14, bold=True, color=NAVY)
    add_text(s, Inches(3.0), y, Inches(9.6), Inches(1.2), body, size=12, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "Sources: MP Q1 2026 10-Q (mp-20260331); MP DoD partnership press release (Jul 2025); CNBC, S&P Global. See companion .xlsx Notes tab for full source list.")

# ============================================================
# SLIDE 3 — THE QUESTION
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "The question this analysis answers")
add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(1.2),
         '"Can MP convert rare-earth scarcity and US industrial policy into durable cash flow ahead of a competitive peer ramp — '
         'and is today\'s ~$10B market cap a reasonable entry on that thesis?"',
         size=22, italic=True, color=NAVY)

add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.4), "We test that question across five lenses:", size=14, bold=True, color=CHARCOAL)
lenses = [
    "(1) Market — is the NdPr price environment supportive?",
    "(2) Competitive — does MP have a defensible position vs Western and Chinese peers?",
    "(3) Financial — what does the valuation look like vs peer comps?",
    "(4) Strategic — what's the moat and how durable is it?",
    "(5) Decision — what's the cleanest expression of conviction (or wait)?",
]
y = Inches(4.1)
for line in lenses:
    add_text(s, Inches(0.9), y, Inches(12), Inches(0.4), line, size=13, color=CHARCOAL)
    y += Inches(0.45)
slide_footer(s, "Framework adapted from the `competitive-analysis` skill (financial-services/plugins/vertical-plugins/financial-analysis).")

# ============================================================
# SLIDE 4 — MARKET CONTEXT
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Rare-earth demand structurally tight; NdPr +>100% YTD with $110/kg DoD floor under MP")
s.shapes.add_picture(str(NDPR_PNG), Inches(0.5), Inches(1.3), width=Inches(8.5))
add_text(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(0.5), "What the chart says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.3), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• NdPr spot opened 2026 ~$53/kg.",
          "• By end-April hit ~$140/kg in some retail benchmarks; SMM industrial ~$100/kg in May.",
          "• Roughly +100–160% YTD depending on benchmark.",
          "",
          "• MP's $110/kg DoD price floor sits above current SMM spot.",
          "• MP earns the floor on volume covered by the PPA; otherwise sells at spot.",
          "• Lynas, UUUU, USAR are unhedged — they ride spot in both directions.",
          "",
          "• Drivers: China export quotas tightening, accelerating EV magnet demand, US decoupling capex."],
         size=12, color=CHARCOAL)
slide_footer(s, "NdPr prices: SMM, strategicmetalsinvest.com (May 2026). MP–DoD price floor: MP investor relations, Jul 2025.")

# ============================================================
# SLIDE 5 — INDUSTRY ECONOMICS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "China owns >70% of mining and >80% of processing — value pools sit downstream")

# Value chain stages as boxes
stages = ["Mining", "Cracking & Leaching", "Separation (NdPr / HRE)", "Metal / Alloy", "Sintered Magnets", "Downstream apps"]
margins = ["10–20%", "5–10%", "15–25%", "5–10%", "20–30%", "—"]
china_share = ["~70%", "~80%", "~85%", "~90%", "~90%", "n/a"]

x0 = Inches(0.5); y0 = Inches(1.5); w = Inches(2.05); h = Inches(0.8)
for i, st in enumerate(stages):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x0 + Inches(0.05) + i*Inches(2.1), y0, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.color.rgb = NAVY
    tf = box.text_frame; tf.margin_left = tf.margin_right = Emu(45720)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = st; r.font.name=FONT; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=WHITE

# Row labels
add_text(s, Inches(0.5), Inches(2.5), Inches(1.5), Inches(0.4), "Typical margin:", size=11, bold=True, color=CHARCOAL)
add_text(s, Inches(0.5), Inches(2.95), Inches(1.5), Inches(0.4), "China share:", size=11, bold=True, color=CHARCOAL)

# Margin / share row
for i in range(6):
    add_text(s, x0 + Inches(0.05) + i*Inches(2.1), Inches(2.5), w, Inches(0.4),
             margins[i], size=11, color=COPPER, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, x0 + Inches(0.05) + i*Inches(2.1), Inches(2.95), w, Inches(0.4),
             china_share[i], size=11, color=RED, align=PP_ALIGN.CENTER, bold=True)

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5),
         "What this means for investors", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
         ["• The highest-margin stages (Separation, Sintered Magnets) are where strategic value sits — and where China is most dominant.",
          "• MP and USAR are the only Western names targeting Sintered Magnets at scale. Lynas is strong through Separation but does not produce magnets.",
          "• Energy Fuels is a processing pivot from uranium — capacity optionality, but no magnets.",
          "• China Northern Rare Earth is fully integrated AND state-controlled — sets the marginal price for the global market.",
          "• This is why MP's Stage III magnet facility (Independence, TX, commercial 2028) is the thesis-critical milestone."],
         size=12, color=CHARCOAL)
slide_footer(s, "Market shares: Discovery Alert, Farmonaut, SMM (2025/26 industry reports). Margins are typical-industry ranges, not company-specific.")

# ============================================================
# SLIDE 6 — MP PROFILE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "MP profile — Q1'26 inflection: revenue +49%, NdPr +63%, Adj EBITDA flipped to +$36.6M")

# Metrics table
data = [
    ["Metric", "Value", "YoY change", "Source / note"],
    ["Q1 2026 Revenue", "$90.6M", "+49%", "Q1 2026 earnings release"],
    ["Q1 2026 NdPr oxide production", "917 MT", "+63%", "Record quarter"],
    ["Q1 2026 NdPr oxide sales", "1,006 MT", "+117%", "Sales outpaced production via inventory draw"],
    ["Q1 2026 Adj EBITDA", "$36.6M", "+$49.9M", "Driven by $110/kg DoD price floor (PPA, eff. Oct 2025)"],
    ["Magnetics segment revenue", "$21.1M", "+$15.9M", "Magnetic precursor products"],
    ["TTM revenue", "$275M", "n/a", "Sum of trailing four quarters"],
    ["Market cap", "$9.81B", "n/a", "May 15, 2026 (StockAnalysis)"],
    ["52-week range", "$18.64 – $100.25", "n/a", "Stock $61.74 (May 21)"],
    ["FY 2026 consensus revenue growth", "+67%", "n/a", "Analyst aggregate"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.3), data,
          col_widths=[Inches(3.6), Inches(2.4), Inches(1.8), Inches(4.5)], font_size=12)
slide_footer(s, "MP Materials Q1 2026 earnings release; 10-Q FY2026 Q1 (mp-20260331); StockAnalysis.com (May 2026).")

# ============================================================
# SLIDE 7 — DOD DEAL
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "DoD ~15% stake + 10-yr offtake + $110/kg NdPr floor = public–private moat")
data = [
    ["Element", "Detail"],
    ["Form", "DoD purchases $400M of newly-created preferred stock, convertible to common + warrant for additional common"],
    ["Aggregate stake", "~15% of MP common (post-conversion + warrant exercise) as of Jul 9, 2025"],
    ["Conversion price", "$30.03 per share of common stock"],
    ["Close date", "July 11, 2025"],
    ["DoD position", "Becomes MP's LARGEST single shareholder"],
    ["Offtake commitment", "10-year DoD off-take of MP magnet/REE output"],
    ["NdPr price floor (PPA)", "$110/kg on NdPr products covered by the agreement; effective Oct 1, 2025"],
    ["Use of proceeds", "Expand separation & processing AND magnet production capacity"],
    ["Stage III magnet facility", "Construction begins 2026; first commercial product expected 2028"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.3), data,
          col_widths=[Inches(3.6), Inches(8.7)], font_size=12)
slide_footer(s, "Sources: MP Materials press release (Jul 10, 2025); CNBC; Center on Global Energy Policy (Columbia SIPA); Payne Institute explainer.")

# ============================================================
# SLIDE 8 — PEER SET / TIER DIAGRAM
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer set — four strategic groups, two of which MP competes in directly")
groups = [
    ("Tier 1 — Western pure-plays",  ["MP Materials (NYSE: MP) — only scaled US name",
                                       "Lynas Rare Earths (ASX: LYC) — only scaled non-US Western name"], NAVY),
    ("Tier 2 — US juniors / ramping", ["Energy Fuels (NYSE: UUUU) — uranium with REE pivot",
                                       "USA Rare Earth (Nasdaq: USAR) — magnet-focused, $1.6B US gov't investment"], COPPER),
    ("Tier 3 — Chinese majors (context)", ["China Northern Rare Earth (SHA: 600111) — world's largest, state-owned",
                                            "(Shenghe, JL Mag, etc. — not investable for most Western mandates)"], RED),
    ("Tier 4 — Critical-minerals adjacent", ["Albemarle (NYSE: ALB) — lithium leader, not REE",
                                              "Included for broader critical-minerals reference only"], GREY),
]
y = Inches(1.4)
for title, items, col in groups:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(4.2), Inches(0.5), title, size=14, bold=True, color=col)
    for i, it in enumerate(items):
        add_text(s, Inches(0.85), y + Inches(0.45) + i*Inches(0.32), Inches(11.5), Inches(0.3),
                 "• " + it, size=12, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "Peer set per `comps-analysis` 4–6 rule. Excluded: Chinese mid-caps (not investable), early-stage exploration plays without near-term cash flow.")

# ============================================================
# SLIDE 9 — POSITIONING 2x2
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Scale × Integration: MP and Lynas alone occupy the Western 'fully-integrated' quadrant")
s.shapes.add_picture(str(MATRIX_PNG), Inches(0.4), Inches(1.3), width=Inches(8.8))
add_text(s, Inches(9.4), Inches(1.4), Inches(3.7), Inches(0.5), "What the matrix says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.4), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• China Northern RE is the only fully integrated, scaled producer globally — but un-investable for many.",
          "",
          "• MP sits in the top-right Western quadrant, alone with Lynas — and uniquely targets sintered magnets (vs Lynas which does not).",
          "",
          "• USAR is a magnet-only pure-play; small scale.",
          "",
          "• Energy Fuels is processing-only with REE optionality.",
          "",
          "• Albemarle is in the adjacent lithium / bromine space — not a true comp."],
         size=12, color=CHARCOAL)
slide_footer(s, "Integration score derived from public stage disclosures. Bubble size reflects market cap (qualitative).")

# ============================================================
# SLIDE 10 — LYNAS DEEP DIVE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dive: Lynas — Most-profitable Western pure-play, slower up the value chain")
metrics = [
    ["Metric", "Value"],
    ["Market cap", "~US$10.8B (A$15.1B)"],
    ["Q3 FY26 revenue", "A$265M (+115% YoY, +31% QoQ)"],
    ["H1 FY26 revenue / NI", "A$413.7M / A$80.2M"],
    ["FY26 forecast revenue", "~A$1.1B (~2x FY25)"],
    ["FY26 NdPr volume", "~8,800 MT (+35% YoY)"],
    ["FY26 NdPr avg realised", "A$118/kg (+48%)"],
    ["TTM EBITDA / Margin", "US$211M / ~36%"],
    ["TTM P/E", "~209x (forward ~37x)"],
    ["Key asset", "Mt Weld, WA — highest-grade REE deposit globally"],
    ["Processing", "Kuantan, Malaysia; Kalgoorlie cracking; Seadrift, TX (heavy REE, DoD-funded)"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.0), metrics,
          col_widths=[Inches(2.3), Inches(3.9)], font_size=11)

qual = [
    ["Category", "Assessment"],
    ["Business", "Mine → Crack/Leach → Separation, including Heavy REE (Dy/Tb) at US Seadrift facility. Does NOT produce magnets."],
    ["Strengths", "• Highest-grade asset globally\n• Already profitable (EBITDA +$211M TTM)\n• US heavy-REE pivot with DoD support"],
    ["Weaknesses", "• Stops at separation — no magnet revenue\n• Single-mine concentration risk (Mt Weld)\n• Currency exposure (AUD revenue, USD pricing)"],
    ["Strategy", "Volume growth (8,800 MT NdPr in FY26) + US footprint expansion. Will compete with MP in heavy REE."],
]
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.0), qual,
          col_widths=[Inches(1.2), Inches(4.7)], font_size=11)
slide_footer(s, "Lynas H1 FY26 results (Dec 2025); Q3 FY26 update; StockAnalysis.com; S&P Global Market Intelligence (Dec 2025).")

# ============================================================
# SLIDE 11 — JUNIORS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dives: US juniors — capacity optionality, execution risk, long-dated")
uuuu = [
    ["Energy Fuels (UUUU)", ""],
    ["Business", "US uranium producer; pivoting into REE processing at White Mesa Mill (Utah)"],
    ["Phase 2 plan", "6,000 tpa NdPr + 240 tpa Dy + 66 tpa Tb"],
    ["Latest data", "BFS published Jan 15, 2026 — \"lowest-cost NdPr in the world\""],
    ["Strengths", "Existing milling infrastructure; low CAPEX vs greenfield"],
    ["Weaknesses", "REE revenue still pre-commercial; uranium-dominated"],
    ["Read", "Optionality play; thesis hinges on FID and ramp execution"],
]
usar = [
    ["USA Rare Earth (USAR)", ""],
    ["Business", "Mine-to-magnet integrated supply chain (Texas to Oklahoma)"],
    ["Q2 2026 event", "Stillwater, OK magnet line commercial production begins"],
    ["Financing", "$1.5B PIPE closed (Inflection Point anchor); $1.6B US gov't investment"],
    ["Strengths", "Magnet pure-play; gov't backing; first commercial magnets in scope"],
    ["Weaknesses", "Pre-meaningful revenue; small scale; integration not fully proven"],
    ["Read", "Single-thesis bet: can it deliver US magnets at scale on schedule?"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.0), uuuu,
          col_widths=[Inches(1.6), Inches(4.6)], font_size=11)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.0), usar,
          col_widths=[Inches(1.6), Inches(4.3)], font_size=11)
slide_footer(s, "UUUU Phase 2 BFS (Jan 2026); USAR 8-K (Q2 2026 magnet commissioning); Motley Fool (Feb 2026 US$1.6B investment).")

# ============================================================
# SLIDE 12 — CHINESE MAJORS / CNREG
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer context: China Northern Rare Earth — sets the price; un-investable for most Western capital")
cnreg = [
    ["Metric", "Value"],
    ["Global share — mining", "~60–70% controlled by China; CNREG largest single producer"],
    ["Global share — processing", ">80%"],
    ["FY2025 revenue", "¥42.6B (US$5,910M) — up 29% YoY"],
    ["FY2025 net income", "¥2.25B (US$313M) — up 124% YoY"],
    ["H1 2025 NI growth", "+1,882% to +2,014% YoY (off low base, rare-earth price recovery)"],
    ["Q1 2026 net profit guidance", "US$125–131M (+109% to +118% YoY)"],
    ["Parent / ownership", "Baogang Group (state-owned, Inner Mongolia)"],
    ["Strategic role", "Sets global NdPr benchmark price; vertical integration mine → magnet"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(7.0), Inches(5.0), cnreg,
          col_widths=[Inches(2.6), Inches(4.4)], font_size=11)

add_text(s, Inches(7.8), Inches(1.4), Inches(5.0), Inches(0.5), "Why this matters for MP", size=14, bold=True, color=NAVY)
add_text(s, Inches(7.8), Inches(1.95), Inches(5.0), Inches(5.0),
         ["• CNREG's price actions set the marginal NdPr clearing price — MP's economics ride on that benchmark for any volume outside the DoD floor.",
          "",
          "• China can introduce export quotas / restrictions on heavy REE (Dy/Tb) at policy speed — direct upside for MP/Lynas, but execution-risk lever.",
          "",
          "• CNREG profits +124% YoY says: China is winning today even as the West tries to decouple. Decoupling theme is real but slow.",
          "",
          "• Pair-trade idea: long MP / short CNREG offers a cleaner geopolitical-decoupling expression than MP outright."],
         size=11, color=CHARCOAL)
slide_footer(s, "Shanghai Metals Market (SMM); Discovery Alert; Farmonaut (2025/26 industry reports).")

# ============================================================
# SLIDE 13 — COMPARATIVE SCOREBOARD
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Comparative scoreboard — MP leads scale & integration vector; lags near-term margins")
def dots(level):
    # level 1-3
    if level == 3: return "●●● "
    if level == 2: return "●●○ "
    if level == 1: return "●○○ "
    return "○○○ "

scoreboard = [
    ["Dimension", "MP",         "Lynas",      "UUUU",       "USAR",       "CNREG",      "Albemarle"],
    ["Revenue scale (TTM)",
        dots(2) + "$275M", dots(2) + "$578M", dots(1) + "$70M[E]", dots(0) + "Pre-rev", dots(3) + "$5.9B", dots(3) + "$5.85B"],
    ["Vertical integration",
        dots(3) + "Mine→Mag plan", dots(2) + "Mine→Sep+HRE", dots(1) + "Proc", dots(1) + "Magnet only", dots(3) + "Full", dots(0) + "Adjacent"],
    ["TTM profitability",
        dots(1) + "Neg→Pos Adj", dots(3) + "+36% EBITDA", dots(0) + "Neg", dots(0) + "Neg", dots(2) + "+12% EBITDA", dots(3) + "+43% EBITDA"],
    ["Govt backing / floor",
        dots(3) + "DoD 15% + $110 floor", dots(2) + "US/AU grants", dots(1) + "Modest", dots(3) + "$1.6B US inv.", dots(3) + "State-owned", dots(1) + "IRA / DoE"],
    ["China-policy upside",
        dots(3) + "Direct", dots(3) + "Direct", dots(2) + "Modest", dots(3) + "Direct", dots(0) + "Negative", dots(1) + "Indirect"],
    ["Valuation absorb-shock",
        dots(2) + "Floor protects", dots(3) + "Profitable", dots(1) + "Capacity bet", dots(1) + "Optionality", dots(3) + "Cheap", dots(2) + "Improving"],
    ["Investability for US mandate",
        dots(3) + "Yes", dots(2) + "ASX-listed", dots(3) + "Yes", dots(3) + "Yes", dots(0) + "No", dots(3) + "Yes"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), scoreboard,
          col_widths=[Inches(2.6), Inches(1.7), Inches(1.7), Inches(1.5), Inches(1.6), Inches(1.6), Inches(1.6)], font_size=10)
slide_footer(s, "● = Strong  ○ = Weak. Subjective ratings calibrated to data in companion .xlsx Operating Metrics / Valuation / REE-Specific tabs.")

# ============================================================
# SLIDE 14 — MOATS / STRATEGIC SYNTHESIS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Strategic synthesis — Moat: integration + govt partnership · Risk: NdPr cycle + Stage III timing")
moats = [
    ["Moat type",          "MP",      "Lynas",   "UUUU",   "USAR",   "CNREG",  "What it means"],
    ["Network effects",    "Weak",    "Weak",    "Weak",   "Weak",   "Mod",    "Limited — REE is industrial supply chain, not platform"],
    ["Switching costs",    "Strong",  "Mod",     "Weak",   "Mod",    "Strong", "Mag spec-in is sticky; 10-yr offtake hard-wires customers"],
    ["Scale economies",    "Mod",     "Mod",     "Weak",   "Weak",   "Strong", "Cost curve favors integrated scale players (CNREG)"],
    ["Intangibles (govt)", "Strong",  "Mod",     "Weak",   "Mod",    "Strong", "MP–DoD relationship is the single biggest Western intangible"],
    ["Permits / location", "Strong",  "Mod",     "Mod",    "Mod",    "Mod",    "Mountain Pass is the only operating US REE mine"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.0), moats,
          col_widths=[Inches(2.1), Inches(1.2), Inches(1.2), Inches(1.0), Inches(1.0), Inches(1.2), Inches(4.6)], font_size=11)
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4), "Net moat assessment", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
         ["MP's durable advantages: (i) only operating US REE mine, (ii) DoD anchor + price floor, (iii) only Western mine→magnet ambition at scale.",
          "MP's structural vulnerabilities: (i) magnet revenue not until 2028, (ii) NdPr spot still drives marginal volume, (iii) capex-heavy through 2028."],
         size=12, color=CHARCOAL)
slide_footer(s, "Moat framework per `competitive-analysis` skill Step 9. Ratings: Strong / Moderate (Mod) / Weak — calibrated against peer comparison.")

# ============================================================
# SLIDE 15 — BULL / BASE / BEAR
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Range of outcomes anchored by $110/kg floor (downside) and 2028 magnet revenue (upside)")
scenarios = [
    ["Scenario", "Prob.", "Key driver", "Implied direction"],
    ["Bull", "30%", "Section 232 tariffs on REE imports; China heavy-REE export restrictions; Apple/GM expand magnet volume; Stage III on schedule",
        "MP re-rates 50–100% on multi-year magnet revenue visibility; magnet revenue 2028 → 2–3x"],
    ["Base", "50%", "NdPr ranges $80–$150/kg; DoD floor protects covered volume; Stage III breaks ground on time; Lynas/USAR add competitive supply but market absorbs",
        "MP delivers TTM revenue growth >50%; stock tracks NdPr + magnet ramp milestones; ~20–40% upside over 18 months"],
    ["Bear", "20%", "China–US trade détente collapses NdPr; Stage III delays past 2028; USAR/UUUU ramp faster than expected, eroding pricing power",
        "MP re-rates 30–40% lower; trades closer to 4x EV/Sales (peer median) on stalled magnet narrative"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5), scenarios,
          col_widths=[Inches(1.4), Inches(1.0), Inches(5.4), Inches(4.5)], font_size=12)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
         "Signposts to watch (next 3 events that re-rate the scenario weighting)", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         "1) NdPr spot vs $110/kg floor   2) Stage III construction milestone in H2 2026   3) Section 232 / export-control decisions on critical minerals",
         size=12, color=CHARCOAL)
slide_footer(s, "Probabilities are author's judgment, not consensus. For investment use, calibrate to your own conviction and time horizon.")

# ============================================================
# SLIDE 16 — CATALYSTS TIMELINE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Catalysts & signposts — 2026 → 2028")
events = [
    ("Q2 2026",  "USAR",          "Stillwater (OK) sintered NdFeB magnet line commercial production begins", "Competitive supply tick-up"),
    ("H2 2026",  "MP Materials",  "Stage III magnet facility construction commences (Independence, TX)",       "Capex spike; long-dated optionality"),
    ("H2 2026",  "Lynas",         "US heavy-REE (Dy/Tb) plant at Seadrift TX ramps",                            "First non-China heavy-REE supply"),
    ("Q4 2026",  "Energy Fuels",  "Phase 2 REE circuit final investment decision expected",                     "Validates 6,000 tpa NdPr capex"),
    ("Ongoing",  "US Policy",     "Section 232 review / China export-quota decisions",                          "Sets relative price for non-China REE"),
    ("2027",     "MP Materials",  "Apple / GM magnet supply expansion expected",                                "Revenue mix shifts toward magnets"),
    ("2028",     "MP Materials",  "First commercial sintered magnet product (Independence, TX)",                "THESIS-CRITICAL — completes mine-to-magnet"),
]
data = [["When", "Who", "Event", "Impact on MP thesis"]] + [list(e) for e in events]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.3), data,
          col_widths=[Inches(1.2), Inches(2.0), Inches(5.5), Inches(3.6)], font_size=11)
slide_footer(s, "Catalysts compiled from MP investor commentary, peer guidance, and recent media. Track via /catalysts skill in production usage.")

# ============================================================
# SLIDE 17 — DECISION FRAME
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Decision frame — four ways to express the rare-earths thesis")
options = [
    ("Long MP outright (highest conviction)",
     "Triggers: NdPr sustainably > $110/kg floor; Stage III milestones on schedule; new offtake volume from Apple/GM. "
     "Risk: magnet ramp slips past 2028; China relaxes export rules and NdPr collapses. Position sizing reflects multi-year hold.",
     COPPER),
    ("Pair trade: long MP / short CNREG",
     "Triggers: Section 232 tariffs on REE imports; escalating US–China decoupling. "
     "Risk: CNREG profits +124% YoY says China is winning today — short side may be early. Cleaner geopolitical expression than MP outright.",
     NAVY),
    ("Basket via REMX (theme exposure)",
     "REMX YTD ~25–42% — diversified critical-minerals exposure. "
     "Trade-off: significant Chinese name weight inside the ETF; less direct US-domestic-supply exposure than MP single-name.",
     CHARCOAL),
    ("Avoid / wait",
     "Triggers to re-enter: (i) Stage III magnet line achieves first commercial output (~2028), (ii) trailing EBITDA turns sustainably positive, "
     "(iii) NdPr breaks below $80/kg without sufficient floor-protected volume mix. Acknowledges 2026 valuation already prices in significant ramp.",
     GREY),
]
y = Inches(1.4)
for title, body, col in options:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(11.8), Inches(0.4), title, size=14, bold=True, color=col)
    add_text(s, Inches(0.85), y + Inches(0.45), Inches(11.8), Inches(0.85), body, size=11.5, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "This deck is research and decision framing — not investment advice. Cross-check valuation against a primary terminal (CapIQ / FactSet / Bloomberg) before acting.")

# ============================================================
# SLIDE 18 — SOURCES & CAVEATS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Sources, data gaps, and caveats")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "Primary sources", size=14, bold=True, color=NAVY)
sources = [
    "MP Materials Q1 2026 10-Q (mp-20260331) and Q1 2026 earnings release",
    "MP Materials – DoD partnership press release (Jul 10, 2025); CNBC (Pentagon largest shareholder)",
    "Lynas H1 FY26 results; Q3 FY26 trading update; StockAnalysis.com LYC.AX statistics",
    "S&P Global Market Intelligence — Lynas 2026 outlook (Dec 2025)",
    "Energy Fuels Phase 2 BFS press release (Jan 15, 2026); Q1 2026 results",
    "USA Rare Earth 8-K filings (Jan/Feb 2026 PIPE close); magnet-line commissioning press release",
    "China Northern Rare Earth FY25 results — Shanghai Metals Market (SMM); Discovery Alert",
    "Albemarle Q1 2026 earnings; SWOT analysis (Investing.com)",
    "NdPr pricing: SMM, strategicmetalsinvest.com, rare-earth-mining.com (May 2026)",
    "Market context: VanEck REMX ETF, Center on Global Energy Policy (Columbia SIPA)",
]
y = Inches(1.85)
for src in sources:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.3), "• " + src, size=11, color=CHARCOAL)
    y += Inches(0.30)

add_text(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.4),
         "Data gaps and caveats (read before using in a binding decision)", size=14, bold=True, color=NAVY)
y += Inches(0.55)
caveats = [
    "MCP data sources (CapIQ / FactSet / Daloopa) WERE NOT USED — the financial-services plugin connectors were not configured in this environment. Numbers should be cross-checked against a terminal before action.",
    "[E] flags in the companion .xlsx Inputs tab identify figures estimated from public commentary; cell comments document each assumption.",
    "Chinese major financials reported in CNY — conversions use ~7.2 CNY/USD (May 2026 indicative); rerun if FX moves materially.",
    "Forward consensus (P/E forward, FY26 EPS estimates) drawn from public aggregators rather than sell-side feeds. Variance to consensus may exist.",
    "This research is for educational/decision-framing use only and is not investment advice.",
]
for c in caveats:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.5), "• " + c, size=11, color=CHARCOAL)
    y += Inches(0.36)
slide_footer(s, "Built per the `competitive-analysis` and `comps-analysis` skills in financial-services/plugins/vertical-plugins/financial-analysis. Companion file: MP-Comps-Analysis.xlsx.")

# ============================================================
prs.save(OUT)
print(f"Wrote {OUT}")
