"""
Build GOOG-Competitive-Analysis.pptx — investment-decision deck for Alphabet.
18 slides + 3 charts. Pattern follows analysis/AVGO/build_deck.py.
"""
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"; CHARTS.mkdir(exist_ok=True)
OUT = HERE / "GOOG-Competitive-Analysis.pptx"

NAVY = RGBColor(0x17, 0x36, 0x5D)
GBLUE = RGBColor(0x42, 0x85, 0xF4)   # Google blue accent
CHARCOAL = RGBColor(0x40, 0x40, 0x40); GREY = RGBColor(0xBF, 0xBF, 0xBF)
LGREY = RGBColor(0xF2, 0xF2, 0xF2); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B); GREEN = RGBColor(0x2E, 0x7D, 0x32)
FONT = "Times New Roman"

plt.rcParams.update({"font.family":"serif","font.serif":["DejaVu Serif"],
    "axes.edgecolor":"#404040","axes.labelcolor":"#404040","xtick.color":"#404040","ytick.color":"#404040",
    "axes.grid":True,"grid.color":"#E0E0E0","grid.linewidth":0.5})


def chart_capex():
    fig, ax = plt.subplots(figsize=(9, 4.6), dpi=150)
    yrs = ["FY23A", "FY24A", "FY25A", "FY26E", "FY27E", "FY28E", "FY29E", "FY30E"]
    capex = [32.3, 52.5, 92, 185, 200, 180, 160, 150]
    colors = ["#888"]*3 + ["#4285F4"]*5
    bars = ax.bar(yrs, capex, color=colors, edgecolor="#202020", linewidth=0.5)
    for b, v in zip(bars, capex):
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+3, f"${v:.0f}B", ha="center", fontsize=9, fontweight="bold")
    ax.set_ylabel("Capex ($B)")
    ax.set_title("Alphabet capex cycle — $92B (FY25) → $185B (FY26 guide) → taper",
                 fontsize=12.5, color="#17365D", pad=10)
    ax.set_ylim(0, 220)
    fig.tight_layout(); out = CHARTS / "capex.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


def chart_2x2():
    fig, ax = plt.subplots(figsize=(9.5, 6), dpi=150)
    # (rev_$B TTM, growth %, name, color, size~mktcap)
    data = [
        (425, 22, "Alphabet",     "#4285F4", 1300),
        (223, 33, "Meta",         "#17365D", 800),
        (318, 18, "Microsoft",    "#7F8C8D", 1500),
        (743, 17, "Amazon",       "#95A5A6", 1100),
        (47, 16, "Netflix",       "#C0392B", 500),
        (2.96, 12, "Trade Desk",  "#7F8C8D", 300),
        (2.2, 69, "Reddit",       "#2E7D32", 350),
    ]
    for rev, gr, name, col, size in data:
        ax.scatter(rev, gr, s=size, color=col, alpha=0.7, edgecolors="#202020", linewidths=1.1)
        ax.annotate(name, (rev, gr), xytext=(8, 7), textcoords="offset points", fontsize=10, fontweight="bold")
    ax.set_xscale("log")
    ax.set_xlabel("Revenue scale, $B TTM (log)")
    ax.set_ylabel("Revenue growth, MRQ YoY (%)")
    ax.axhline(20, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.axvline(100, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.set_xlim(1, 1500); ax.set_ylim(0, 80)
    ax.set_title("Positioning: scale × growth (bubble = market cap)", fontsize=12.5, color="#17365D", pad=12)
    fig.tight_layout(); out = CHARTS / "positioning.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


def chart_football():
    fig, ax = plt.subplots(figsize=(9.2, 4.6), dpi=150)
    rows = [
        ("DCF — base (WACC/g grid $167-334)", 167, 334),
        ("DCF — bull (Cloud margin >36%, faster taper)", 320, 450),
        ("Comps — mega-cap median fwd P/E ~23x", 280, 380),
        ("Comps — EV/EBITDA mega-cap median", 260, 380),
    ]
    ypos = range(len(rows))
    for y, (lab, lo, hi) in zip(ypos, rows):
        ax.barh(y, hi-lo, left=lo, height=0.5, color="#17365D", alpha=0.75, edgecolor="#202020")
        ax.text(lo-6, y, f"${lo:.0f}", va="center", ha="right", fontsize=9)
        ax.text(hi+6, y, f"${hi:.0f}", va="center", ha="left", fontsize=9)
    ax.set_yticks(list(ypos)); ax.set_yticklabels([r[0] for r in rows], fontsize=9.5)
    ax.axvline(358, color="#C0392B", linewidth=2)
    ax.text(358, len(rows)-0.3, "  Current $358", color="#C0392B", fontsize=10, fontweight="bold")
    ax.axvline(213, color="#2E7D32", linewidth=1.5, linestyle="--")
    ax.text(213, -0.9, "DCF base $213", color="#2E7D32", fontsize=9, ha="center")
    ax.set_xlim(100, 500); ax.set_xlabel("Implied value per share ($)")
    ax.set_title("Valuation football field — market price sits in/above the bull range", fontsize=12.5, color="#17365D", pad=10)
    ax.invert_yaxis()
    fig.tight_layout(); out = CHARTS / "football.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


CAPEX_PNG = chart_capex(); MATRIX_PNG = chart_2x2(); FB_PNG = chart_football()

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(45720); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    runs = [text] if isinstance(text, str) else text
    for i, line in enumerate(runs):
        if i == 0: r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align; r = p2.add_run()
        r.text = line; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def slide_title(slide, title, *, size=24):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = GBLUE; bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95), title, size=size, bold=True, color=NAVY)


def slide_footer(slide, txt):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3), txt, size=9, italic=True, color=GREY)


def add_table(slide, left, top, width, height, data, *, header=True, col_widths=None, font_size=12):
    rows = len(data); cols = len(data[0])
    shp = slide.shapes.add_table(rows, cols, left, top, width, height); tbl = shp.table
    if col_widths:
        for i, w in enumerate(col_widths): tbl.columns[i].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT; run.font.size = Pt(font_size)
                    if header and ri == 0: run.font.bold = True; run.font.color.rgb = WHITE
                    elif ci == 0: run.font.bold = True; run.font.color.rgb = NAVY
                    else: run.font.color.rgb = CHARCOAL
            if header and ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else: cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Emu(45720); cell.margin_top = cell.margin_bottom = Emu(27432)
    return shp


# SLIDE 1 — COVER
s = prs.slides.add_slide(BLANK)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.5))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.12))
stripe.fill.solid(); stripe.fill.fore_color.rgb = GBLUE; stripe.line.fill.background()
add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5), "NASDAQ: GOOG / GOOGL — Investment Decision Frame", size=18, italic=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.4), "Alphabet Inc.", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.9), "Competitive Landscape, Comps & DCF Valuation", size=26, color=GBLUE)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5), "As of June 3, 2026", size=14, italic=True, color=CHARCOAL)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.6),
         ["Mega-cap cohort: META · MSFT · AMZN · NFLX  +  TTD, RDDT (ads-pure memo)",
          "Built with deep-research + competitive-analysis + comps-analysis + 3-statement-model + dcf-model skills",
          "Companion: GOOG-Comps-Analysis.xlsx · GOOG-Model.xlsx"],
         size=13, color=CHARCOAL)

# SLIDE 2 — THESIS
s = prs.slides.add_slide(BLANK)
slide_title(s, "Cheap on fwd P/E (~24x), but $185B FY26 capex compresses FCF — DCF $213 vs $358")
bullets = [
    ("1. The franchise", "Alphabet Q1 CY26: revenue $109.9B (+22%), Cloud $20B (+63%), Services $89.6B (+16%). Cloud op margin reaccelerated to 32.9% (from 17.8% YoY). Cloud backlog $460B (~2x QoQ). $4.35T market cap on ~$425B TTM revenue."),
    ("2. The capex shock", "Capex ramps from $92B (FY25) to $180-190B (FY26 guide) — a ~$95B step. Q1 actual was $22.4B (back-end loaded). Combined with ~$85B June equity raise (incl. Berkshire private placement), this is a structural posture shift from self-funding capital return."),
    ("3. The DCF problem", "Base case ($213/share, -40% vs $358) is heavy because Year-1 unlevered FCF compresses to ~$4B as capex outruns D&A; FCF normalizes to $200B+ by FY30. Even the bull corner (WACC 8.0% / g 4.0%) only reaches $334. The market is pricing aggressive Cloud margin + capex payback."),
    ("4. The read", "GOOG screens cheap on fwd P/E (24x — cheapest mega-cap ex-META) but expensive on intrinsic given the capex cycle and equity raise. Patient long — wait for capex peak signal. DOJ remedies (no Chrome) is a relief; Apple TAC remains the live P&L risk."),
]
y = Inches(1.35)
for h, b in bullets:
    add_text(s, Inches(0.6), y, Inches(2.4), Inches(1.2), h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(3.1), y, Inches(9.6), Inches(1.2), b, size=11.5, color=CHARCOAL)
    y += Inches(1.36)
slide_footer(s, "Sources: Alphabet Q1 CY26 8-K/IR; GOOG-Model.xlsx (validated DCF); company press releases (June 2026 equity raise FWP). Research, not investment advice.")

# SLIDE 3 — QUESTION
s = prs.slides.add_slide(BLANK)
slide_title(s, "The question this analysis answers")
add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(1.2),
         '"Alphabet has the best AI infrastructure franchise (TPU, Cloud +63%) and trades at the cheapest mega-cap forward P/E (~24x), '
         'but is committing ~$185B of FY26 capex and raising ~$85B of equity. What is the intrinsic value through the cycle, and is $358 a reasonable entry?"',
         size=19, italic=True, color=NAVY)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4), "We test that across five lenses:", size=14, bold=True, color=CHARCOAL)
for i, ln in enumerate([
    "(1) Market — Cloud demand, AI capex sustainability, ads resilience.",
    "(2) Competitive — vs MSFT/Azure, AMZN/AWS, META on ads, AI infra winners.",
    "(3) Financial (comps) — how cheap is GOOG vs the mega-cap median?",
    "(4) Intrinsic (DCF) — what is fair value through the capex cycle?",
    "(5) Decision — bull/base/bear and the cleanest expression.",
]):
    add_text(s, Inches(0.9), Inches(4.2)+i*Inches(0.46), Inches(12), Inches(0.4), ln, size=13, color=CHARCOAL)

# SLIDE 4 — CAPEX CHART
s = prs.slides.add_slide(BLANK)
slide_title(s, "The defining number for FY26-27: capex doubles — funding via equity, not just cash")
s.shapes.add_picture(str(CAPEX_PNG), Inches(0.5), Inches(1.35), width=Inches(8.4))
add_text(s, Inches(9.1), Inches(1.4), Inches(4.0), Inches(0.5), "What this means", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.1), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• FY25 capex ~$92B (already 6x 2022). FY26 guide $180-190B — a ~$95B step.",
          "",
          "• Q1 CY26 only $22.4B — implies $50B+/qtr in H2 to hit the guide range.",
          "",
          "• ~$85B June 2026 equity raise (Berkshire pp + ATM + common) funds this — first major issuance in years.",
          "",
          "• My base assumes capex tapers from FY28 ($180B) to FY30 ($150B). Reasonable but a real risk if Cloud demand exceeds capacity.",
          "",
          "• D&A lag means FCF crushed in FY26, normalizing by FY28."],
         size=10.5, color=CHARCOAL)
slide_footer(s, "Sources: Alphabet 10-Ks FY23-25; Q1 CY26 8-K; June 2026 FWP (Berkshire pp); CNBC.")

# SLIDE 5 — SEGMENT SPLIT
s = prs.slides.add_slide(BLANK)
slide_title(s, "Three businesses: Services (~$343B cash cow), Cloud (re-accel to +63%), Other Bets (Waymo)")
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
    ["Segment (FY25)", "Revenue", "YoY", "Op income", "Op margin", "Q1 26 trend"],
    ["Google Services", "$342.7B", "~+12%", "~$154B (45%)", "~45%", "+16% YoY"],
    ["  of which Search", "~$220B [E]", "~+13%", "—", "—", "Resilient w/ AI Overviews"],
    ["  of which YouTube Ads", "~$36B [E]", "+15-18%", "—", "—", "Ad reaccel"],
    ["Google Cloud", "$58.7B", "+36%", "$13.9B", "24%", "+63% rev / 32.9% margin"],
    ["Other Bets", "$1.5B", "~flat", "($7.5B)", "(488%)", "Waymo $16B raise @ $126B val"],
    ["Total Alphabet", "$402.8B", "+15%", "~$129B", "~32%", "Q1: $109.9B (+22%)"],
], col_widths=[Inches(2.7), Inches(1.5), Inches(1.4), Inches(2.0), Inches(1.5), Inches(3.2)], font_size=10.5)
add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
         ["• Services is the cash engine (~45% op margin), and the AI Overviews integration is preserving Search query share.",
          "• Cloud is the marginal growth story — margin expansion from 5% (FY23) to 24% (FY25) to 32.9% (Q1 26) is the leverage that justifies the capex.",
          "• Other Bets is a long-dated option — Waymo $16B raise in Feb 2026 at ~$126B implied valuation marks the position.",
          "• Anthropic + Broadcom 3.5GW TPU deal (April 2026) extends GOOG's third-party AI accelerator franchise — a Cloud-backlog anchor."],
         size=12, color=CHARCOAL)
slide_footer(s, "Sources: Alphabet 10-K FY25; Q1 CY26 8-K segment disclosures; February 2026 Waymo funding round (TSG / press); April 2026 Anthropic-Google-Broadcom announcement.")

# SLIDE 6 — PROFILE
s = prs.slides.add_slide(BLANK)
slide_title(s, "Alphabet profile — Q1 CY26 prints + balance-sheet restructuring underway")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0), [
    ["Metric", "Value", "Note"],
    ["Q1 CY26 revenue", "$109.9B (+22%)", "Beat consensus; +19% cc"],
    ["Cloud Q1 26 revenue / margin", "$20.0B / 32.9%", "Margin up from 17.8% YoY (huge)"],
    ["Cloud backlog (Q1 26)", "$460B", "~2x QoQ; ~50% recognized in 24 mo"],
    ["Q1 CY26 net income / EPS", "$62.6B / $5.11", "+81% YoY (incl. Anthropic gains)"],
    ["Q1 capex / annual guide", "$22.4B / $180-190B", "Back-end loaded ramp"],
    ["June 2026 equity raise (announced)", "~$84.75B", "Berkshire pp + $40B ATM + $40B common"],
    ["Net cash (Q1 26)", "~$80B", "Cash+sec $126.8B – LT debt $46.5B (after Q1 $31B senior notes)"],
    ["FY25 buybacks / div", "$45.7B / ~$10B", "Buybacks down from ~$62B FY24 as capex absorbs"],
    ["Market cap / share", "~$4.35T / ~$358", "12.12B diluted shares"],
    ["P/E TTM / Fwd", "~32x / ~24x", "Cheapest mega-cap fwd P/E ex-META"],
], col_widths=[Inches(4.5), Inches(3.5), Inches(4.3)], font_size=11)
slide_footer(s, "Sources: Alphabet Q1 CY26 8-K/IR; FWP for June 2026 equity raise; StockAnalysis market data.")

# SLIDE 7 — INDUSTRY (CLOUD + ADS)
s = prs.slides.add_slide(BLANK)
slide_title(s, "Cloud + ads landscape — GOOG is the #3 cloud catching #2, and the AI infra premium is real")
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.6), [
    ["Dimension", "Alphabet", "Microsoft", "Amazon"],
    ["Cloud revenue Q1 26", "$20.0B (+63%)", "Azure +40% (within IC $34.7B)", "AWS $37.6B (+28%)"],
    ["Cloud op margin", "32.9% (Q1 26)", "39.7% (IC segment)", "37.7% (AWS)"],
    ["Cloud backlog", "$460B", "~$298B (RPO; Q3 FY26)", "~$200B+ [E]"],
    ["AI angle", "TPU + Anthropic 3.5GW", "OpenAI 27% stake + Azure", "Bedrock + Anthropic $8B deal"],
    ["FY26 capex", "$180-190B", "~$190B (calendar)", "~$200B+ (Q1 alone $44B)"],
    ["AI revenue run-rate", "Cloud + Search AI surfaces", "$37B AI ARR (+123%)", "Bedrock + ads"],
], col_widths=[Inches(3.0), Inches(3.3), Inches(3.0), Inches(3.0)], font_size=11)
add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4),
         ["• All three are racing to spend ~$200B+/yr on AI infra; the differentiator is monetization economics on each $ of capex.",
          "• Google's TPU + Anthropic relationship is the cleanest case for owning the inference workload at lower unit cost than NVDA GPUs.",
          "• Microsoft has the OpenAI moat (and a $228B implied stake). Amazon has scale + cash. Alphabet has cost (TPU) + breadth (Search + YouTube + Cloud)."],
         size=12, color=CHARCOAL)
slide_footer(s, "Sources: MSFT Q3 FY26 IR; AMZN Q1 CY26 IR; Alphabet Q1 CY26 IR; SemiAnalysis (TPU v7 Ironwood).")

# SLIDE 8 — MOATS
s = prs.slides.add_slide(BLANK)
slide_title(s, "Alphabet moat — Search distribution + TPU cost + YouTube data + Anthropic anchor + Waymo option")
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.3), [
    ["Layer", "Detail"],
    ["Search distribution", "~90% global query share; AI Overviews preserving query stream; Chrome 60%+ browser share — DOJ DID NOT order divestiture (Sept 2025 ruling)"],
    ["TPU + Anthropic anchor", "Ironwood TPU v7 shipping 2026; Anthropic 3.5GW commit (April 2026); ~1M TPU chips to Anthropic — third-party validation"],
    ["YouTube + first-party data", "~2B+ users; richest video + interest graph; ad business +15-18%"],
    ["Cloud backlog", "$460B (~2x QoQ); RPO recognized over 24 mo provides revenue visibility"],
    ["Waymo (option value)", "$16B Feb 2026 raise at ~$126B implied valuation; Alphabet's majority owner"],
    ["Cash + debt access", "Pre-raise net cash ~$80B + investment-grade credit + capacity to raise (proven June 2026)"],
    ["GDM / Gemini", "Frontier AI model (Pro 1.5+); Search integration; consumer Gemini subs growing"],
], col_widths=[Inches(3.4), Inches(8.9)], font_size=11)
slide_footer(s, "Sources: Alphabet IR; DOJ remedies ruling (Sept 2025, Hughes Hubbard summary); Waymo funding round (Feb 2026); Anthropic-Broadcom-Google deal (April 2026).")

# SLIDE 9 — PEER SET
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer set — mega-cap cohort for the platform comparison; ads-pure for the Search/YouTube lens")
groups = [
    ("Mega-cap cohort (main stat set)", ["Meta (META) — cheapest fwd P/E (~18x); Reality Labs $4B/qtr drag", "Microsoft (MSFT) — Azure + OpenAI moat; ~$190B FY26 capex", "Amazon (AMZN) — AWS reaccel +28%; ads $70B+ TTM", "Netflix (NFLX) — streaming premium ~26x P/E"], NAVY),
    ("Ads-pure memo set (Search/YouTube reference)", ["The Trade Desk (TTD) — programmatic; decelerating (+12%)", "Reddit (RDDT) — hyper-growth ad surface (+69%); data-licensing"], GBLUE),
]
y = Inches(1.5)
for t, items, col in groups:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(2.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(8), Inches(0.5), t, size=14, bold=True, color=col)
    for i, it in enumerate(items):
        add_text(s, Inches(0.85), y+Inches(0.45)+i*Inches(0.45), Inches(11.5), Inches(0.4), "• " + it, size=12, color=CHARCOAL)
    y += Inches(2.6)
slide_footer(s, "Stats over the 4 mega-cap peers; TTD + RDDT are memo (you don't median them into a mega-cap set).")

# SLIDE 10 — POSITIONING
s = prs.slides.add_slide(BLANK)
slide_title(s, "Positioning: META has the smallest scale + fastest growth; GOOG and AMZN cluster")
s.shapes.add_picture(str(MATRIX_PNG), Inches(0.4), Inches(1.3), width=Inches(8.8))
add_text(s, Inches(9.4), Inches(1.4), Inches(3.7), Inches(0.5), "What this says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.4), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• META is unusual: smaller revenue base (~$223B) but fastest growth (+33%); justifies its lower fwd P/E being a buying opportunity if the AI ads thesis works.",
          "",
          "• MSFT and AMZN are the scale benchmarks (~$300-700B revenue) at +17-18% growth.",
          "",
          "• GOOG sits between — $425B at +22% — a sweet spot on scale × growth.",
          "",
          "• NFLX is a smaller premium streaming play; RDDT is the hyper-growth outlier (+69%)."],
         size=11, color=CHARCOAL)
slide_footer(s, "Bubble = market cap. Source: GOOG-Comps-Analysis.xlsx.")

# SLIDE 11 — META DEEP DIVE
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dive: Meta — cheapest mega-cap on fwd P/E (~18x); Reality Labs the structural drag")
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.0), [
    ["Metric", "Value"],
    ["Market cap", "~$1.55T"],
    ["TTM revenue", "~$223B"],
    ["Q1 26 growth", "+33% YoY"],
    ["Gross margin", "~82%"],
    ["TTM EBITDA / margin", "~$110B / ~51%"],
    ["TTM net income", "~$71B"],
    ["P/E TTM / Fwd", "~22x / ~18x"],
    ["EV/EBITDA / EV/Rev", "~12x / ~7x"],
    ["Reality Labs (Q1 26)", "Loss $4.03B on $402M rev (cumulative >$90B)"],
], col_widths=[Inches(2.5), Inches(3.7)], font_size=11)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.0), [
    ["Category", "Assessment"],
    ["Business", "Family of Apps (FB, IG, WhatsApp, Threads) + Reality Labs (Quest, Ray-Ban Meta, AR)."],
    ["Strengths", "• Cheapest mega-cap fwd P/E\n• Family of Apps op margin >40%\n• Llama open-source moat\n• 3.5B+ DAUs"],
    ["Weaknesses", "• RL drag persistent\n• Apple ATT exposure ongoing\n• Capex $125-145B FY26 (raised) — margin pressure"],
    ["Read-through to GOOG", "META multiple compression bounded GOOG's downside on ads; META's AI-spend mix more focused on inference (similar to GOOG's TPU bet)"],
], col_widths=[Inches(1.4), Inches(4.5)], font_size=10.5)
slide_footer(s, "Sources: Meta Q1 CY26 8-K; multiples.vc; StockAnalysis. RL: 10-Q.")

# SLIDE 12 — MSFT + AMZN
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dives: Microsoft (the cloud benchmark) + Amazon (the FCF whale)")
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.2), [
    ["Microsoft (MSFT)", ""],
    ["Market cap / TTM rev", "~$4.72T / ~$318B"],
    ["Fiscal Q3 FY26 growth", "+18% YoY"],
    ["Azure growth / Intelligent Cloud margin", "+40% / 39.7%"],
    ["AI ARR / Copilot seats", "$37B (+123%) / >20M"],
    ["P/E TTM / Fwd", "~26x / ~23x"],
    ["OpenAI stake", "~27% at ~$228B (vs $13B invested = ~17x MoIC)"],
    ["FY26 capex", "~$190B (calendar)"],
    ["Read", "The benchmark for cloud + AI. Owns the OpenAI relationship economics."],
], col_widths=[Inches(2.5), Inches(3.7)], font_size=10.5)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.2), [
    ["Amazon (AMZN)", ""],
    ["Market cap / TTM rev", "~$2.7T / ~$743B"],
    ["AWS Q1 26", "$37.6B (+28%) — fastest in 15Q; $14.2B op income (37.7%)"],
    ["NA retail / Intl op margin", "9.0% / 3.6%"],
    ["Ads Q1 26 / TTM", "$17.2B (+24%) / $70B+"],
    ["P/E TTM / Fwd", "~30x / ~28x"],
    ["FY25 net debt", "~$66B (incl. lease)"],
    ["Q1 capex", "$44B (annualized ~$175B+)"],
    ["Read", "AWS reaccel is the standout; retail margin lifting. Highest fwd P/E in the cohort."],
], col_widths=[Inches(2.5), Inches(3.4)], font_size=10.5)
slide_footer(s, "Sources: MSFT Q3 FY26 IR; AMZN Q1 CY26 IR; Bloomberg (OpenAI stake May 2026).")

# SLIDE 13 — SCOREBOARD
s = prs.slides.add_slide(BLANK)
slide_title(s, "Comparative scoreboard — GOOG leads on Cloud growth + multiple value; trails on near-term FCF")
def dots(n): return {3: "●●● ", 2: "●●○ ", 1: "●○○ "}.get(n, "○○○ ")
add_table(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.4), [
    ["Dimension", "GOOG", "META", "MSFT", "AMZN", "NFLX"],
    ["Revenue scale", dots(2)+"$425B", dots(1)+"$223B", dots(2)+"$318B", dots(3)+"$743B", dots(1)+"$47B"],
    ["Growth (MRQ)", dots(2)+"+22%", dots(3)+"+33%", dots(1)+"+18%", dots(1)+"+17%", dots(1)+"+16%"],
    ["Cloud growth", dots(3)+"+63%", dots(0)+"n/a", dots(2)+"Azure +40%", dots(2)+"AWS +28%", dots(0)+"n/a"],
    ["AI infra moat", dots(3)+"TPU + Anthropic", dots(2)+"Llama", dots(3)+"OpenAI 27%", dots(2)+"Bedrock", dots(0)+"limited"],
    ["FCF margin", dots(1)+"compressed FY26", dots(2)+"~25%", dots(2)+"~32%", dots(2)+"~10%", dots(3)+"~25%"],
    ["Balance sheet", dots(3)+"net cash $80B", dots(3)+"net cash $40B", dots(2)+"slight cash", dots(1)+"net debt $66B", dots(1)+"net debt $4B"],
    ["Value (fwd P/E)", dots(2)+"~24x", dots(3)+"~18x", dots(2)+"~23x", dots(1)+"~30x", dots(1)+"~26x"],
], col_widths=[Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)], font_size=10.5)
slide_footer(s, "● = stronger; ○ = weaker. Higher fwd P/E = weaker value score. Source: GOOG-Comps-Analysis.xlsx.")

# SLIDE 14 — VALUATION I (FOOTBALL FIELD)
s = prs.slides.add_slide(BLANK)
slide_title(s, "Valuation I — DCF $213 vs $358; comps $260-380; only the bull case touches the price")
s.shapes.add_picture(str(FB_PNG), Inches(0.5), Inches(1.35), width=Inches(8.4))
add_text(s, Inches(9.1), Inches(1.4), Inches(4.0), Inches(0.5), "What this shows", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.1), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• DCF base ($213) sits well below $358; full grid $167-334.",
          "",
          "• Bull DCF (Cloud margin >36%, faster capex taper) reaches ~$320-450 — overlaps the current price.",
          "",
          "• Comps (mega-cap median fwd P/E ~23x × FY27 EPS ~$13) imply ~$300-380.",
          "",
          "• Conclusion: $358 is priced for the bull capex-payback case."],
         size=11, color=CHARCOAL)
slide_footer(s, "Source: GOOG-Model.xlsx (DCF); GOOG-Comps-Analysis.xlsx (multiples).")

# SLIDE 15 — VALUATION II (DCF + SENSITIVITY)
s = prs.slides.add_slide(BLANK)
slide_title(s, "Valuation II — base DCF $213; sensitivity grid $167-334; only bull WACC/g exceeds price")
add_table(s, Inches(0.5), Inches(1.4), Inches(5.6), Inches(4.6), [
    ["DCF bridge (base)", "$M / $"],
    ["PV of explicit FCF (FY26-30)", "343,877"],
    ["PV of terminal value", "2,162,879"],
    ["Enterprise value", "2,506,757"],
    ["(+) Net cash", "80,300"],
    ["Equity value", "2,587,057"],
    ["÷ Diluted shares (M)", "12,120"],
    ["Implied value / share", "$213.45"],
    ["Current price", "$357.73"],
    ["Upside / (downside)", "(40.3%)"],
], col_widths=[Inches(3.4), Inches(2.2)], font_size=11)
add_text(s, Inches(6.4), Inches(1.4), Inches(6.4), Inches(0.4), "Sensitivity — implied $/share (WACC × g)", size=12, bold=True, color=NAVY)
add_table(s, Inches(6.4), Inches(1.85), Inches(6.5), Inches(3.2), [
    ["WACC ↓ / g →", "2.0%", "2.5%", "3.0%", "3.5%", "4.0%"],
    ["8.0%", "231", "250", "272", "300", "334"],
    ["8.5%", "211", "227", "245", "267", "294"],
    ["9.0% (base)", "195", "208", "223", "241", "263"],
    ["9.5%", "180", "191", "204", "219", "237"],
    ["10.0%", "167", "177", "188", "201", "215"],
], col_widths=[Inches(1.7), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95)], font_size=10.5)
add_text(s, Inches(6.4), Inches(5.2), Inches(6.5), Inches(1.7),
         ["• TV is ~86% of EV — extremely WACC/g-sensitive (Cloud growth + capex cycle).",
          "• FY26 unlevered FCF compresses to ~$4B as capex outruns D&A; normalizes to $200B+ by FY30.",
          "• Bull (Cloud margin >36%, faster taper) overlays cleanly with ~$320-450 range."],
         size=11, color=CHARCOAL)
slide_footer(s, "Source: GOOG-Model.xlsx (validated, 0 formula errors, 267 formulas).")

# SLIDE 16 — BULL/BASE/BEAR
s = prs.slides.add_slide(BLANK)
slide_title(s, "Bull / base / bear — anchored to Cloud margin trajectory and capex efficiency")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.8), [
    ["Scenario", "Prob.", "Key drivers", "DCF value"],
    ["Bull", "30%", "Cloud margin >36% by FY28; capex tapers FY27; Waymo achieves commercial monetization; DOJ remedies have no economic impact",
     "~$320-450 — overlaps the current $358 price"],
    ["Base", "45%", "Cloud margin 30%→36% by FY30; capex stays elevated through FY28; Apple TAC modestly impaired; equity raise dilutive but not repeated",
     "~$213 (~40% below)"],
    ["Bear", "25%", "Cloud margin plateaus at 32%; capex stays at $200B+ as Anthropic + sovereign deals demand more; Apple TAC eliminated; DOJ remedies bite Search",
     "~$167-188 (~50% below)"],
], col_widths=[Inches(1.3), Inches(0.9), Inches(5.6), Inches(4.5)], font_size=11)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4), "Signposts that move the weighting", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.9),
         "1) Q2/Q3 CY26 capex run-rate vs the $180-190B guide  2) Cloud op margin trajectory  3) DOJ appeals / Apple default-payment economics  4) Waymo IPO / direct monetization",
         size=12, color=CHARCOAL)
slide_footer(s, "Probabilities are author's judgment. Bull DCF requires both capex and margin assumptions to swing favorably.")

# SLIDE 17 — DECISION
s = prs.slides.add_slide(BLANK)
slide_title(s, "Decision frame & sources")
add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4), "Four ways to express a view", size=14, bold=True, color=NAVY)
opts = [
    ("Patient long — wait for the capex peak signal", "Best franchise of the mega-caps but at $358 you're paying for successful payback. Buy on Q2/Q3 capex disappointments or any Cloud growth deceleration scare.", NAVY),
    ("Relative-value pair: long GOOG / short AMZN", "Same AI infra exposure; GOOG cheaper on fwd P/E (24x vs 30x) with stronger Cloud growth (+63% vs +28%); cleaner geopolitical / antitrust profile post-Sept 2025 ruling.", GBLUE),
    ("Theme play: long META (highest growth, cheapest P/E)", "If you want mega-cap AI capex theme with most multiple expansion potential, META at ~18x fwd P/E + 33% growth is the best risk-reward.", CHARCOAL),
    ("Trim / avoid here", "If you believe capex doesn't pay back fast enough OR DOJ remedies bite, current price already discounts the bull case.", GREY),
]
y = Inches(1.7)
for t, b, col in opts:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(11.8), Inches(0.35), t, size=13, bold=True, color=col)
    add_text(s, Inches(0.85), y+Inches(0.38), Inches(11.8), Inches(0.6), b, size=11, color=CHARCOAL)
    y += Inches(1.02)
add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.35), "Sources & caveats", size=12, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.9),
         ["Alphabet Q1 CY26 8-K/IR; FY23-25 10-Ks; June 2026 equity raise FWP (Berkshire pp); peer 10-Q/IR; CNBC/Bloomberg/Reuters for DOJ + Waymo + Anthropic deal.",
          "MCP terminals NOT configured. [E] flags in workbooks. Other Bets revenue is residual. Research only — NOT investment advice."],
         size=10, italic=True, color=CHARCOAL)
slide_footer(s, "Built per: deep-research + competitive-analysis + comps-analysis + 3-statement-model + dcf-model. Companions: GOOG-Model.xlsx, GOOG-Comps-Analysis.xlsx.")

# SLIDE 18 — SOURCES (lite)
s = prs.slides.add_slide(BLANK)
slide_title(s, "Primary sources & data-gap caveats")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4), "Primary sources", size=14, bold=True, color=NAVY)
srcs = [
    "Alphabet Q1 CY26 8-K & 10-Q (Apr 2026); FY25 10-K (Feb 2026); FY24 10-K; FY23 10-K",
    "Alphabet FWP filings (June 2026 equity raise — Berkshire Hathaway private placement, ATM, common)",
    "DOJ v. Google search remedies ruling (Sept 2025); appeals (Jan 2026 GOOG, DOJ cross-appeal)",
    "Anthropic-Google-Broadcom 3.5GW TPU expansion (April 2026); SemiAnalysis on Ironwood TPU v7",
    "Waymo $16B funding round (Feb 2026) at ~$126B implied valuation (TSG, press)",
    "Peer 10-Qs / IR releases: META Q1 26, MSFT Q3 FY26, AMZN Q1 26, NFLX Q1 26, TTD Q1 26, RDDT Q1 26",
    "Market data: StockAnalysis, companiesmarketcap, GuruFocus, Macrotrends, FinanceCharts (~June 1-3, 2026)",
]
y = Inches(1.85)
for src in srcs:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.3), "• " + src, size=11, color=CHARCOAL); y += Inches(0.30)
add_text(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.4), "Caveats", size=14, bold=True, color=NAVY)
y += Inches(0.55)
caveats = [
    "MCP terminal connectors (CapIQ/FactSet/Daloopa) NOT configured — public sources used and [E]-flagged in workbooks.",
    "GOOG TTM NI inflated by Anthropic-related fair-value gains; conservative GAAP basis used in model.",
    "DCF assumes capex peaks FY27 and tapers; if Anthropic/sovereign deals demand more, FCF normalization delayed.",
    "Other Bets revenue (~$1.5B) is largely Waymo; not broken out separately.",
    "This is research / decision-framing only — NOT investment advice.",
]
for c in caveats:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.5), "• " + c, size=11, color=CHARCOAL); y += Inches(0.36)
slide_footer(s, "Built per competitive-analysis + comps-analysis + 3-statement-model + dcf-model + audit-xls + ib-check-deck skills.")

prs.save(OUT)
print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
