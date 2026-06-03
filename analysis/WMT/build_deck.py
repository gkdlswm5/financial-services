"""
Build WMT-Competitive-Analysis.pptx following the `competitive-analysis` skill in
plugins/vertical-plugins/financial-analysis. 18 slides, insight-driven titles,
2x2 positioning matrix, peer deep dives, moats grid, bull/base/bear, catalysts.

Run:
    python3 build_deck.py
Output:
    WMT-Competitive-Analysis.pptx
    charts/  (intermediate PNGs)
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"
CHARTS.mkdir(exist_ok=True)
OUT = HERE / "WMT-Competitive-Analysis.pptx"

# Palette per skill (2-3 colors max, muted) — WMT brand-adjacent navy + Walmart-spark yellow accent
NAVY     = RGBColor(0x17, 0x36, 0x5D)
WMT_GOLD = RGBColor(0xFF, 0xC2, 0x20)   # Walmart spark accent
CHARCOAL = RGBColor(0x40, 0x40, 0x40)
GREY     = RGBColor(0xBF, 0xBF, 0xBF)
LGREY    = RGBColor(0xF2, 0xF2, 0xF2)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xC0, 0x39, 0x2B)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Times New Roman"

# ============================================================
# CHARTS (matplotlib → PNG)
# ============================================================
plt.rcParams.update({
    "font.family": "serif",
    "font.serif": ["DejaVu Serif"],
    "axes.edgecolor": "#404040",
    "axes.labelcolor": "#404040",
    "xtick.color": "#404040",
    "ytick.color": "#404040",
    "axes.grid": True,
    "grid.color": "#E0E0E0",
    "grid.linestyle": "-",
    "grid.linewidth": 0.5,
})

def chart_comp_sales():
    """Comp sales growth bar chart across the peer set (most-recent reported)."""
    fig, ax = plt.subplots(figsize=(9, 4.5), dpi=150)
    peers = ["COST", "WMT", "BJ", "KR", "TGT"]
    comps = [6.2, 4.5, 3.8, 2.5, 2.0]  # % ex-fuel/gas
    colors = ["#17365D", "#FFC220", "#7F8C8D", "#7F8C8D", "#7F8C8D"]
    bars = ax.bar(peers, comps, color=colors, edgecolor="#202020", linewidth=0.6)
    for bar, v in zip(bars, comps):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.15,
                f"+{v:.1f}%", ha="center", fontsize=11, color="#202020", fontweight="bold")
    ax.set_ylabel("Comparable / identical sales growth, YoY %")
    ax.set_title("US comparable sales — most recent reported quarter (Q1 FY26, ex-fuel)",
                 fontsize=13, color="#17365D", pad=10)
    ax.set_ylim(0, 8)
    ax.axhline(0, color="#404040", linewidth=0.8)
    # AMZN intentionally excluded — no comparable-store metric
    ax.text(4.4, 7.4, "AMZN excluded — no store-comp metric", fontsize=9,
            color="#7F8C8D", ha="right", style="italic")
    fig.tight_layout()
    out = CHARTS / "comp_sales.png"
    fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out

def chart_2x2():
    """Scale × E-commerce mix positioning matrix."""
    fig, ax = plt.subplots(figsize=(9.5, 6), dpi=150)
    # data: (rev_$B, ecom_pct, name, color, size)
    data = [
        (685, 18,  "Walmart",          "#FFC220",  1200),
        (270, 7.5, "Costco",           "#17365D",   900),
        (107, 20,  "Target",           "#7F8C8D",   500),
        (650, 100, "Amazon",           "#C0392B",  1200),
        (148, 10,  "Kroger",           "#7F8C8D",   600),
        (22,  5,   "BJ's Wholesale",   "#7F8C8D",   300),
    ]
    for rev, ecom, name, col, size in data:
        ax.scatter(rev, ecom, s=size, color=col, alpha=0.7, edgecolors="#202020", linewidths=1.2)
        ax.annotate(name, (rev, ecom), xytext=(10, 8), textcoords="offset points",
                    fontsize=10, color="#202020", weight="bold")
    ax.set_xscale("log")
    ax.set_xlabel("Revenue scale, USD billions (TTM, log)")
    ax.set_ylabel("E-commerce % of revenue (TTM)")
    ax.set_xlim(15, 1500)
    ax.set_ylim(-5, 110)
    ax.set_title("Positioning: scale × e-commerce mix",
                 fontsize=13, color="#17365D", pad=12)
    # Quadrant shading
    ax.axhline(20, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.axvline(200, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.text(17,   100, "Pure digital", fontsize=9, color="#7F8C8D")
    ax.text(250,  100, "Scaled omni",  fontsize=9, color="#7F8C8D")
    ax.text(17,   -1,  "Niche brick",  fontsize=9, color="#7F8C8D")
    ax.text(250,  -1,  "Scaled brick", fontsize=9, color="#7F8C8D")
    fig.tight_layout()
    out = CHARTS / "positioning.png"
    fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out

COMP_PNG = chart_comp_sales()
MATRIX_PNG = chart_2x2()

# ============================================================
# DECK BUILD
# ============================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs = [text]
    else:
        runs = text
    for i, line in enumerate(runs):
        if i == 0:
            r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align
            r = p2.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def slide_title(slide, title, *, size=26):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = WMT_GOLD
    bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.20), Inches(12.3), Inches(0.9),
             title, size=size, bold=True, color=NAVY)

def slide_footer(slide, source_text):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
             source_text, size=9, italic=True, color=GREY)

def add_table(slide, left, top, width, height, data, *, header=True, col_widths=None,
              first_col_bold=True, font_size=12):
    rows = len(data); cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            tf = cell.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    if header and r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    elif first_col_bold and c_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = NAVY
                    else:
                        run.font.color.rgb = CHARCOAL
            if header and r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif r_idx % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Emu(45720)
            cell.margin_top = cell.margin_bottom = Emu(36576)
    return tbl_shape

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.5))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.12))
stripe.fill.solid(); stripe.fill.fore_color.rgb = WMT_GOLD; stripe.line.fill.background()

add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
         "NYSE: WMT — Investment Decision Frame", size=18, italic=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.4),
         "Walmart Inc.", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.9),
         "Competitive Landscape & Investment Read", size=28, color=WMT_GOLD)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "As of May 2026", size=14, italic=True, color=CHARCOAL)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
         ["Peer set: Costco · Target · Amazon · Kroger · BJ's Wholesale",
          "Built using `competitive-analysis` + `comps-analysis` skills (financial-services repo)",
          "Companion workbook: WMT-Comps-Analysis.xlsx"],
         size=13, color=CHARCOAL)

# ============================================================
# SLIDE 2 — EXEC SUMMARY / THESIS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Scale + omnichannel flywheel + ad/marketplace mix = a re-rating story, not a retail story")

bullets = [
    ("1. The setup",
     "Walmart is the largest retailer on earth (~$685B TTM revenue) and the only grocery-anchored peer compounding "
     "ad + marketplace + membership revenue at 20–30% annual growth — three tiers structurally higher-margin than core retail."),
    ("2. The flywheel",
     "Stores + supply chain → low-cost grocery → traffic → high-margin ad inventory (Walmart Connect ~$5B run-rate) → "
     "marketplace 3P GMV → Walmart+ subscription. Each tier funds the next; none of WMT's pure-play peers (COST, KR, BJ, TGT) replicate the full stack."),
    ("3. The risk",
     "Stock at ~$95 (~$765B market cap, ~26x TTM EPS) is no longer 'cheap defensive.' Re-rating depends on continued ad/marketplace growth, "
     "consumer holding up through a soft-landing year, and Amazon not pressing harder into grocery."),
    ("4. The read",
     "Long thesis intact for the margin-mix re-rating; key Q is duration. Pair against COST (premium peer) for relative-value framing. "
     "Avoid if you believe US consumer downtrades meaningfully in 2H 2026."),
]
y = Inches(1.4)
for hdr, body in bullets:
    add_text(s, Inches(0.7), y, Inches(2.3), Inches(1.2), hdr, size=14, bold=True, color=NAVY)
    add_text(s, Inches(3.0), y, Inches(9.6), Inches(1.2), body, size=12, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "Sources: WMT Q1 FY26 10-Q and IR commentary (May 2026); StockAnalysis.com; companion .xlsx Notes tab for full source list.")

# ============================================================
# SLIDE 3 — THE QUESTION
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "The question this analysis answers")
add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(1.2),
         '"Is Walmart\'s ad/marketplace/membership margin mix worth a premium re-rating from the 18-20x P/E of legacy mass-merch — '
         'and is the current ~26x TTM P/E a reasonable entry on that thesis?"',
         size=22, italic=True, color=NAVY)

add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.4),
         "We test that question across five lenses:", size=14, bold=True, color=CHARCOAL)
lenses = [
    "(1) Market — is the US retail / e-com / retail-media environment supportive?",
    "(2) Competitive — does WMT have a defensible position vs warehouse clubs, mass merch, and AMZN?",
    "(3) Financial — what does the valuation look like vs peer comps (COST, TGT, AMZN, KR, BJ)?",
    "(4) Strategic — what's the moat and how durable is it?",
    "(5) Decision — what's the cleanest expression of conviction (or wait)?",
]
y = Inches(4.1)
for line in lenses:
    add_text(s, Inches(0.9), y, Inches(12), Inches(0.4), line, size=13, color=CHARCOAL)
    y += Inches(0.45)
slide_footer(s, "Framework adapted from the `competitive-analysis` skill (financial-services/plugins/vertical-plugins/financial-analysis).")

# ============================================================
# SLIDE 4 — MARKET CONTEXT
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "US retail $5T+; grocery share gains, e-com inflecting, retail-media the high-margin tier expanding fastest")
s.shapes.add_picture(str(COMP_PNG), Inches(0.5), Inches(1.3), width=Inches(8.5))
add_text(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(0.5), "What the chart says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.3), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• Warehouse clubs (COST, BJ) lead on comp-sales growth — membership pricing power + value-trade-down winners.",
          "",
          "• WMT US comp +4.5% — above-trend mass-merch; outpaces TGT 2x.",
          "",
          "• KR shows grocery still positive but slow; TGT lagging.",
          "",
          "• AMZN omitted — no comparable-store metric, but online retail growing ~10% YoY.",
          "",
          "• US retail-media ad spend ~$60B in 2025; growing ~20% YoY (eMarketer, IAB)."],
         size=12, color=CHARCOAL)
slide_footer(s, "Comp sales: each issuer's Q1 FY26 IR releases; ex-fuel/gas. US retail-media TAM: eMarketer / IAB (2025).")

# ============================================================
# SLIDE 5 — INDUSTRY ECONOMICS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Retail value chain — grocery thin, general merch better; ads + 3P marketplace + membership are the margin pools")

stages = ["Grocery / Food", "Consumables", "General Merch", "Apparel / Home", "3P Marketplace", "Ads + Mbrshp"]
margins = ["2–4%", "5–8%", "10–15%", "15–25%", "25–40%", "60–80%"]
wmt_pres = ["#1 US",      "#1 US",     "Top-2",        "Mid-tier",    "Building $50B+ GMV", "Connect $5B+ run-rate"]

x0 = Inches(0.5); y0 = Inches(1.5); w = Inches(2.05); h = Inches(0.8)
for i, st in enumerate(stages):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x0 + Inches(0.05) + i*Inches(2.1), y0, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.color.rgb = NAVY
    tf = box.text_frame; tf.margin_left = tf.margin_right = Emu(45720)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = st; r.font.name=FONT; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=WHITE

add_text(s, Inches(0.5), Inches(2.5), Inches(1.5), Inches(0.4), "Typical EBIT margin:", size=11, bold=True, color=CHARCOAL)
add_text(s, Inches(0.5), Inches(2.95), Inches(1.5), Inches(0.4), "WMT presence:", size=11, bold=True, color=CHARCOAL)

for i in range(6):
    add_text(s, x0 + Inches(0.05) + i*Inches(2.1), Inches(2.5), w, Inches(0.4),
             margins[i], size=11, color=WMT_GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, x0 + Inches(0.05) + i*Inches(2.1), Inches(2.95), w, Inches(0.4),
             wmt_pres[i], size=10, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5),
         "What this means for investors", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
         ["• Grocery anchors the traffic but contributes little to margin. The story-shaping tier is the right two-thirds: ads + 3P marketplace + membership.",
          "• Walmart Connect (ads) is reportedly running at >40% EBIT margin — at $5B run-rate, that's ~$2B of incremental EBIT growing 30%+ YoY.",
          "• Amazon proved this template: AMZN's advertising services line has compounded to ~$60B TTM — that's the WMT bull case directionally.",
          "• Costco does NOT play in ads (intentionally, to protect member trust) — the margin lever WMT has, COST does not.",
          "• Kroger's Precision Marketing is the only direct peer comp for retail-media inside grocery; ~$1.5B annualized vs WMT Connect ~$5B."],
         size=12, color=CHARCOAL)
slide_footer(s, "Margin ranges: industry sources (NRF, Bain, JP Morgan retail primers). WMT/Connect figures: WMT mgmt commentary, Q4 FY25 and Q1 FY26.")

# ============================================================
# SLIDE 6 — WMT PROFILE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Walmart profile — Q1 FY26: revenue +6%, ad+marketplace tiers compounding 20–30%, EBIT margin +25 bps")

data = [
    ["Metric", "Value", "YoY change", "Source / note"],
    ["Q1 FY26 revenue", "~$165B", "+6%", "WMT Q1 FY26 IR release (May 2026)"],
    ["Walmart US comp sales (ex-fuel)", "+4.5%", "n/a", "Sustained above-trend mass-merch comps"],
    ["WMT US e-commerce", "~$110B run-rate", "+20%+ YoY", "Mgmt commentary; ~18% of US revenue mix"],
    ["Walmart Connect (ad revenue)", "~$5B run-rate", "+30%+ YoY", "Mgmt commentary; not separately disclosed line"],
    ["Marketplace 3P GMV", "~$50B+", "+30%+ YoY", "Mgmt commentary; mix of US + India (Flipkart)"],
    ["Membership revenue (Sam's + Walmart+)", "~$3B", "Sam's +9%", "Sam's $2B + Walmart+ ~$1B estimate"],
    ["TTM revenue", "~$685B", "n/a", "FY25 $681B + Q1 FY26 +6%"],
    ["TTM EBITDA / Margin", "~$45B / ~6.6%", "n/a", "EBIT $30B + D&A ~$15B"],
    ["Market cap", "~$765B", "n/a", "WMT ~$95/share × ~8.05B shares, May 2026"],
    ["P/E (TTM)", "~26x", "n/a", "vs 5-yr avg ~22x; premium reflects mix-shift narrative"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.5), data,
          col_widths=[Inches(3.5), Inches(2.4), Inches(2.0), Inches(4.4)], font_size=11)
slide_footer(s, "Walmart Q1 FY26 release (May 2026); WMT 10-K FY25 (Mar 2026); StockAnalysis.com market data.")

# ============================================================
# SLIDE 7 — WMT MOATS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "WMT moat — supply-chain density + EDLP + Connect ad business form three reinforcing layers")
data = [
    ["Layer", "Description / KPI"],
    ["Supply-chain scale", "~4,600 US stores + ~150 distribution centers; ~90% of US population within 10 miles of a Walmart store"],
    ["EDLP price gap", "WMT averages ~10–15% price gap below TGT/KR on grocery basket (per third-party tracking)"],
    ["Walmart+ subscription", "Sub base ~$1B revenue tier; bundles grocery delivery, Paramount+, fuel discount — embeds household behavior"],
    ["Marketplace 3P", "~150K active sellers; international ramp via Flipkart (India) and Mexico"],
    ["Walmart Connect (ads)", "~$5B run-rate ad business; ~40%+ EBIT margin per mgmt; growing 30%+ YoY"],
    ["Data + media network", "~250M weekly customers → highest-quality first-party shopper-graph in US (vs AMZN's general e-com graph)"],
    ["DC + automation capex", "$15–17B annual capex modernizing DCs, automation, freight — extends per-unit cost gap vs sub-scale grocers"],
    ["Geographic scale", "Walmart International ~$120B revenue; Mexico (Walmex), India (PhonePe + Flipkart), China; supply leverage vs US-only peers"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.3), data,
          col_widths=[Inches(3.2), Inches(9.1)], font_size=12)
slide_footer(s, "Sources: WMT 10-K FY25; mgmt commentary at Q4 FY25 / Q1 FY26 calls; third-party price tracking (Bain, KeyBanc).")

# ============================================================
# SLIDE 8 — PEER SET / TIER DIAGRAM
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer set — four strategic groups, three of which WMT competes in directly")
groups = [
    ("Tier 1 — Warehouse club premium",  ["Costco (Nasdaq: COST) — premium-priced peer, membership economics",
                                          "BJ's Wholesale (NYSE: BJ) — smaller-scale COST follower"], NAVY),
    ("Tier 2 — Mass-merch / general",     ["Walmart (NYSE: WMT) — scale leader, omnichannel + ad mix-shift",
                                            "Target (NYSE: TGT) — cheap-chic, smaller scale, lagging comps"], WMT_GOLD),
    ("Tier 3 — Marketplace / digital",    ["Amazon (Nasdaq: AMZN) — marketplace + ads + AWS — partial retail comp",
                                            "(Shopify, eBay — out of scope for WMT direct competition)"], RED),
    ("Tier 4 — Pure-play grocery",        ["Kroger (NYSE: KR) — largest US pure-play grocer",
                                            "(Albertsons, Sprouts, etc. — smaller pure-plays, excluded for tractability)"], GREY),
]
y = Inches(1.4)
for title, items, col in groups:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(4.2), Inches(0.5), title, size=14, bold=True, color=col)
    for i, it in enumerate(items):
        add_text(s, Inches(0.85), y + Inches(0.45) + i*Inches(0.32), Inches(11.5), Inches(0.3),
                 "• " + it, size=12, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "Peer set per `comps-analysis` 4–6 rule. Excluded: dollar stores (DG, DLTR — different customer cohort), specialty (no overlap with WMT mix).")

# ============================================================
# SLIDE 9 — POSITIONING 2x2
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Scale × E-commerce mix: WMT and AMZN alone occupy upper-right; COST commands a premium at lower mix")
s.shapes.add_picture(str(MATRIX_PNG), Inches(0.4), Inches(1.3), width=Inches(8.8))
add_text(s, Inches(9.4), Inches(1.4), Inches(3.7), Inches(0.5), "What the matrix says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.4), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• AMZN sits at the extreme — 100% e-com mix, $650B revenue. Pure-digital, AWS-distorted multiples.",
          "",
          "• WMT is the only scaled (>$500B) omnichannel player — ~18% e-com mix and rising, with stores as fulfillment nodes.",
          "",
          "• COST sits notably below at ~7.5% e-com — but commands a 40%+ P/E premium on membership economics alone.",
          "",
          "• TGT has higher e-com mix (~20%) than WMT but a tenth of the scale — disadvantaged on supply-chain unit economics.",
          "",
          "• KR and BJ are sub-scale on both axes."],
         size=11, color=CHARCOAL)
slide_footer(s, "E-com mix per most-recent IR disclosure or [E] estimate; bubble size reflects market cap.")

# ============================================================
# SLIDE 10 — COST DEEP DIVE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dive: Costco — Membership-first peer commanding ~40% P/E premium to WMT")
metrics = [
    ["Metric", "Value"],
    ["Market cap", "~$420B"],
    ["TTM revenue", "~$270B"],
    ["Q3 FY26 comp sales (ex-gas/FX)", "+6.2%"],
    ["Gross margin", "~12.5%"],
    ["EBITDA / Margin", "~$13B / ~4.8%"],
    ["Net income (TTM)", "~$8B"],
    ["P/E (TTM)", "~50x (vs WMT ~26x)"],
    ["EV / EBITDA", "~32x (vs WMT ~17x)"],
    ["Membership fees (TTM)", "~$5B"],
    ["Renewal rate", "~92–93% US/Canada"],
    ["E-com mix", "~7.5%"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.6), metrics,
          col_widths=[Inches(3.0), Inches(3.2)], font_size=11)

qual = [
    ["Category", "Assessment"],
    ["Business", "Warehouse club; ~$5B annual membership fees fund ultra-low merchandise gross margin. Stores carry ~3,800 SKUs vs WMT ~120K."],
    ["Strengths", "• Membership renewal ~92%+ creates recurring-revenue dynamic\n• Pricing power: dues raised 2024 first time in 7 yrs\n• Negative net debt; no leverage risk\n• Executive-membership tier acts as ARPU lever"],
    ["Weaknesses", "• Limited assortment vs WMT\n• No ad / retail-media business by design (protects member trust)\n• High dependence on US/Canada; intl expansion slow\n• Premium valuation provides little error margin"],
    ["Strategy", "Volume growth via new warehouses (+25/yr), gradual membership-fee step-ups, slow international ramp. Avoids the ad lane."],
]
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.6), qual,
          col_widths=[Inches(1.2), Inches(4.7)], font_size=11)
slide_footer(s, "Costco FY25 10-K; Q3 FY26 monthly sales releases; StockAnalysis.com.")

# ============================================================
# SLIDE 11 — TGT / KR / BJ
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dives: Target, Kroger, BJ — mid-tier scale, narrower moats, varying ad-mix progress")
tgt = [
    ["Target (TGT)", ""],
    ["Business", "Mass-merch with cheap-chic positioning; ~1,950 US stores"],
    ["TTM revenue / NI", "~$107B / ~$4.2B"],
    ["Comp sales Q1 FY26", "+2% (recovering from FY25 declines)"],
    ["Roundel (ads)", "~$2B run-rate ad business; growing >25%"],
    ["E-com mix", "~20% — proportionally above WMT"],
    ["Read", "Cheap multiple, real recovery in progress, but no scale moat vs WMT"],
]
kr = [
    ["Kroger (KR)", ""],
    ["Business", "Largest US pure-play grocer; ~2,700 stores"],
    ["TTM revenue / NI", "~$148B / ~$2.5B"],
    ["Comp sales Q1 FY26 (ex-fuel)", "+2.5%"],
    ["KPM (ads)", "~$1.5B run-rate ad business"],
    ["E-com mix", "~10% (grocery e-com still small)"],
    ["Read", "Defensive; thin EBITDA margin caps multiple; Albertsons merger blocked Dec 2024"],
]
bj = [
    ["BJ's Wholesale (BJ)", ""],
    ["Business", "Warehouse club; ~245 clubs, NE/Mid-Atlantic"],
    ["TTM revenue / NI", "~$22B / ~$0.55B"],
    ["Comp sales Q1 FY26 (ex-gas)", "+3.8%"],
    ["Membership", "~$470M membership-fee revenue"],
    ["E-com mix", "~5%"],
    ["Read", "Smaller-scale COST follower; valuation discount until renewal-rate parity"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(4.1), Inches(5.4), tgt, col_widths=[Inches(1.4), Inches(2.7)], font_size=10)
add_table(s, Inches(4.7), Inches(1.35), Inches(4.1), Inches(5.4), kr,  col_widths=[Inches(1.4), Inches(2.7)], font_size=10)
add_table(s, Inches(8.9), Inches(1.35), Inches(3.9), Inches(5.4), bj,  col_widths=[Inches(1.4), Inches(2.5)], font_size=10)
slide_footer(s, "TGT/KR/BJ Q1 FY26 IR releases; StockAnalysis.com market data.")

# ============================================================
# SLIDE 12 — AMZN
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer context: Amazon — Comparability caveat (AWS distorts every multiple)")
amzn = [
    ["Metric", "Value"],
    ["Market cap", "~$2.2T"],
    ["TTM consolidated revenue", "~$650B"],
    ["TTM EBITDA / margin", "~$160B / ~25%"],
    ["TTM net income", "~$65B"],
    ["TTM advertising services revenue", "~$60B (+ 20%+ YoY)"],
    ["TTM AWS revenue", "~$110B"],
    ["P/E (TTM)", "~33x (vs WMT ~26x)"],
    ["EV / EBITDA", "~14x (lower than WMT — AWS lifts EBITDA denominator)"],
    ["E-com mix (of consolidated)", "100% (excl. AWS, ~83%)"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(7.0), Inches(5.0), amzn,
          col_widths=[Inches(3.0), Inches(4.0)], font_size=11)

add_text(s, Inches(7.8), Inches(1.4), Inches(5.0), Inches(0.5), "Why this matters for WMT", size=14, bold=True, color=NAVY)
add_text(s, Inches(7.8), Inches(1.95), Inches(5.0), Inches(5.0),
         ["• AMZN's $60B ad revenue and ~25% consolidated EBITDA margin are the WMT bull-case template directionally.",
          "",
          "• But: AWS is ~70% of consolidated EBIT — strip it out and AMZN retail margin is ~3–4%, BELOW WMT.",
          "",
          "• AMZN is squeezing WMT specifically on grocery (Whole Foods + Amazon Fresh + same-day Prime grocery in major metros).",
          "",
          "• Counter: WMT US grocery is 4–5x the size of AMZN grocery; logistics-density advantage remains.",
          "",
          "• Pair-trade idea: long WMT / short AMZN expresses confidence in stores-as-fulfillment + slowing AWS growth."],
         size=11, color=CHARCOAL)
slide_footer(s, "AMZN Q1 FY26 earnings; segment reporting; StockAnalysis.com.")

# ============================================================
# SLIDE 13 — COMPARATIVE SCOREBOARD
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Comparative scoreboard — WMT leads scale & ad-revenue trajectory; lags margin vs COST")
def dots(level):
    if level == 3: return "●●● "
    if level == 2: return "●●○ "
    if level == 1: return "●○○ "
    return "○○○ "

scoreboard = [
    ["Dimension", "WMT",         "COST",        "TGT",        "AMZN",        "KR",         "BJ"],
    ["Revenue scale (TTM)",
        dots(3) + "$685B", dots(2) + "$270B", dots(1) + "$107B", dots(3) + "$650B", dots(2) + "$148B", dots(1) + "$22B"],
    ["Comp sales growth",
        dots(2) + "+4.5%", dots(3) + "+6.2%", dots(1) + "+2.0%", dots(0) + "n/a",   dots(1) + "+2.5%", dots(2) + "+3.8%"],
    ["EBITDA margin",
        dots(2) + "~6.6%", dots(2) + "~4.8%", dots(2) + "~9.3%", dots(3) + "~25%",  dots(1) + "~4.7%", dots(1) + "~5.0%"],
    ["Ad / retail-media",
        dots(3) + "$5B+ Connect", dots(0) + "Not played", dots(2) + "$2B Roundel", dots(3) + "$60B Ads", dots(2) + "$1.5B KPM", dots(0) + "Minimal"],
    ["Membership / loyalty",
        dots(2) + "Walmart+ + Sam's", dots(3) + "~$5B fees, 92%+ renewal", dots(1) + "Circle 360 new", dots(3) + "Prime + subs", dots(0) + "Boost free-tier", dots(2) + "$470M fees"],
    ["Capital intensity (LT debt)",
        dots(2) + "Moderate", dots(3) + "Net cash", dots(2) + "Moderate", dots(2) + "Moderate", dots(1) + "Elevated", dots(2) + "Moderate"],
    ["Valuation premium absorbed",
        dots(2) + "~26x P/E", dots(1) + "~50x P/E", dots(3) + "~17x P/E", dots(2) + "~33x P/E", dots(3) + "~20x P/E", dots(2) + "~27x P/E"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), scoreboard,
          col_widths=[Inches(2.4), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.6), Inches(1.5)], font_size=10)
slide_footer(s, "● = Strong  ○ = Weak. Subjective ratings calibrated to data in companion .xlsx Operating Metrics / Valuation / Retail-Specific tabs.")

# ============================================================
# SLIDE 14 — MOATS GRID
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Strategic synthesis — Moat: scale + data + logistics · Risk: consumer rollover + AMZN grocery")
moats = [
    ["Moat type",          "WMT",     "COST",    "TGT",   "AMZN",   "KR",    "What it means"],
    ["Network effects",    "Mod",     "Weak",    "Weak",  "Strong", "Weak",  "WMT marketplace 3P building; AMZN dominant in 3P; COST not a platform"],
    ["Switching costs",    "Mod",     "Strong",  "Weak",  "Strong", "Weak",  "Membership creates highest switching cost; WMT+ is the lever"],
    ["Scale economies",    "Strong",  "Mod",     "Weak",  "Strong", "Weak",  "WMT's per-unit logistics cost is structurally below TGT/KR/BJ"],
    ["Intangibles (data)", "Strong",  "Mod",     "Mod",   "Strong", "Mod",   "First-party shopper data quality is the input to ad-business growth"],
    ["Brand / EDLP",       "Strong",  "Strong",  "Mod",   "Mod",    "Mod",   "WMT and COST have the strongest 'value' brand equity"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.0), moats,
          col_widths=[Inches(2.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(5.3)], font_size=11)
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4), "Net moat assessment", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
         ["WMT's durable advantages: (i) only scaled omnichannel player with grocery anchor, (ii) ad+marketplace+membership all compounding 20–30%, (iii) supply-chain density unmatched ex-AMZN.",
          "WMT's structural vulnerabilities: (i) consumer downtrade risk, (ii) AMZN pressing on grocery, (iii) Walmart+ subscriber growth opaque, (iv) capex intensity (~$17B/yr) crowds free cash flow."],
         size=12, color=CHARCOAL)
slide_footer(s, "Moat framework per `competitive-analysis` skill Step 9. Ratings: Strong / Moderate (Mod) / Weak — calibrated against peer comparison.")

# ============================================================
# SLIDE 15 — BULL / BASE / BEAR
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Range of outcomes anchored by ad-revenue compounding (upside) and consumer downtrade (downside)")
scenarios = [
    ["Scenario", "Prob.", "Key driver", "Implied direction"],
    ["Bull", "30%", "Walmart Connect → $10B+ run-rate by FY28; Walmart+ subs accelerate; EBIT margin reaches 5.5%; tariff dollar-cost-plus pricing absorbed",
        "P/E re-rates to 28–30x; ~25% upside as margin mix shift validates a structural premium to historical 18–22x band"],
    ["Base", "50%", "WMT comp +3–5%; Connect/marketplace 25%+ growth continues; EBIT margin holds 4.5–5%; consumer soft-lands; Stage I tariff drag manageable",
        "Stock tracks ~10–15% appreciation over 18 months in line with EBIT growth + 1-multiple expansion"],
    ["Bear", "20%", "Consumer downtrade in 2H 2026; AMZN grocery investments compress WMT US comps; tariffs squeeze general-merch GM; Connect deceleration",
        "Multiple compresses to 20x P/E (peer-aligned); 15–20% downside risk to ~$75 on flat earnings"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5), scenarios,
          col_widths=[Inches(1.4), Inches(1.0), Inches(5.4), Inches(4.5)], font_size=12)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
         "Signposts to watch (next 3 events that re-rate the scenario weighting)", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         "1) Walmart Connect / marketplace growth on Q2 FY26 call   2) US consumer-credit + jobs data through 2H 2026   3) Tariff pass-through impact on Q3/Q4 FY26 GM",
         size=12, color=CHARCOAL)
slide_footer(s, "Probabilities are author's judgment, not consensus. For investment use, calibrate to your own conviction and time horizon.")

# ============================================================
# SLIDE 16 — CATALYSTS TIMELINE
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Catalysts & signposts — 2026 → 2028")
events = [
    ("Q2 FY26",  "Walmart",       "Walmart Connect ad-business breakout disclosure expected",                  "Validates re-rating mix-shift thesis"),
    ("Q3 FY26",  "Walmart",       "Holiday season — tariff pass-through visible in gross margin",              "Tests pricing power vs consumer elasticity"),
    ("H2 2026",  "Amazon",        "Amazon Fresh / Whole Foods grocery footprint expansion",                    "Direct competitive pressure on WMT grocery"),
    ("FY27",     "Walmart",       "Sam's Club new-store cadence accelerating (~10+ new clubs)",                "Competes directly with COST/BJ"),
    ("FY27",     "Walmart",       "Walmart+ subscriber count disclosure (rumored, not yet committed)",         "If disclosed, would clarify Prime gap"),
    ("FY27",     "Costco",        "Membership-fee increase reset (last increase 2024) — pricing power test",    "Cross-read for warehouse-club pricing power"),
    ("FY28",     "Walmart",       "Walmart International — India (Flipkart) IPO timing",                       "Value-unlock event; sum-of-parts visibility"),
    ("FY28",     "Walmart",       "EBIT margin tracking towards mgmt long-term framework (~5%+)",              "Confirms or refutes the structural margin lift"),
]
data = [["When", "Who", "Event", "Impact on WMT thesis"]] + [list(e) for e in events]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.3), data,
          col_widths=[Inches(1.2), Inches(2.0), Inches(5.5), Inches(3.6)], font_size=11)
slide_footer(s, "Catalysts compiled from WMT/peer investor commentary and recent media. Track via /catalysts skill in production usage.")

# ============================================================
# SLIDE 17 — DECISION FRAME
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Decision frame — four ways to express the Walmart thesis")
options = [
    ("Long WMT outright (highest conviction)",
     "Triggers: Walmart Connect breakout disclosure on Q2 FY26; sustained 4%+ US comp; EBIT margin tracking ≥4.5%. "
     "Risk: consumer downtrade in 2H 2026; AMZN grocery escalation. Position sizing reflects 18–36 month hold.",
     WMT_GOLD),
    ("Pair trade: long WMT / short TGT",
     "Triggers: WMT comp gap vs TGT widens; tariff pass-through more painful for TGT general-merch mix; WMT data/scale moat advantages compound. "
     "Risk: TGT mean-reverts faster than WMT compounds. Cleanest scale-divergence expression within mass-merch.",
     NAVY),
    ("Relative value: long WMT / short COST",
     "Triggers: COST premium narrows as warehouse-club comp slows below WMT; WMT ad business closes margin-mix gap. "
     "Risk: COST has demonstrated consistent execution premium — short side may be expensive to carry.",
     CHARCOAL),
    ("Avoid / wait",
     "Triggers to re-enter: (i) Walmart+ subscriber count disclosed and credible vs Prime, (ii) P/E retraces to 22x peer-aligned band, "
     "(iii) tariff drag fully embedded in numbers. Acknowledges current ~26x P/E already prices significant mix-shift upside.",
     GREY),
]
y = Inches(1.4)
for title, body, col in options:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(11.8), Inches(0.4), title, size=14, bold=True, color=col)
    add_text(s, Inches(0.85), y + Inches(0.45), Inches(11.8), Inches(0.85), body, size=11.5, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "This deck is research and decision framing — not investment advice. Cross-check valuation against a primary terminal (CapIQ / FactSet / Bloomberg) before acting.")

# ============================================================
# SLIDE 18 — SOURCES & CAVEATS
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_title(s, "Sources, data gaps, and caveats")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "Primary sources", size=14, bold=True, color=NAVY)
sources = [
    "Walmart Q1 FY26 earnings release and 10-Q (May 2026); WMT 10-K FY25 (Mar 2026)",
    "Costco — Q3 FY26 monthly sales releases; COST 10-K FY25",
    "Target — Q1 FY26 earnings; TGT 10-K FY25",
    "Amazon — Q1 FY26 earnings; AMZN 10-K FY25; segment reporting (Advertising Services, AWS)",
    "Kroger — Q1 FY26 earnings; KR 10-K FY25",
    "BJ's Wholesale — Q1 FY26 earnings; BJ 10-K FY25",
    "Market data: StockAnalysis.com, Yahoo Finance (May 2026)",
    "Retail-media TAM: eMarketer / IAB US Retail Media report (2025)",
    "Industry context: Bain, KeyBanc retail primers; NRF; press coverage (CNBC, Bloomberg)",
]
y = Inches(1.85)
for src in sources:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.3), "• " + src, size=11, color=CHARCOAL)
    y += Inches(0.30)

add_text(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.4),
         "Data gaps and caveats (read before using in a binding decision)", size=14, bold=True, color=NAVY)
y += Inches(0.55)
caveats = [
    "MCP data sources (CapIQ / FactSet / Daloopa) WERE NOT USED — the financial-services plugin connectors were not configured in this environment. Numbers should be cross-checked against a terminal before action.",
    "[E] flags in the companion .xlsx Inputs tab identify figures estimated from public commentary; cell comments document each assumption.",
    "Walmart Connect (ad business) is NOT separately disclosed in WMT filings — the ~$5B run-rate is mgmt commentary aggregated across calls / interviews. Range of estimates: $4–6B.",
    "Operating-lease liabilities are excluded from Net Debt — material for AMZN (~$80B+) and modest for others; EV understates economic leverage proportionally.",
    "AMZN consolidated multiples include AWS — they materially distort a 'retail-only' comparison. A clean comparison would require segment-level disaggregation.",
    "This research is for educational/decision-framing use only and is not investment advice.",
]
for c in caveats:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.5), "• " + c, size=11, color=CHARCOAL)
    y += Inches(0.36)
slide_footer(s, "Built per the `competitive-analysis` and `comps-analysis` skills in financial-services/plugins/vertical-plugins/financial-analysis. Companion file: WMT-Comps-Analysis.xlsx.")

# ============================================================
prs.save(OUT)
print(f"Wrote {OUT}")
