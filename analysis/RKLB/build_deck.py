"""
Build RKLB-Competitive-Analysis.pptx — investment-decision deck for Rocket Lab.
14-slide tight version (small-cap-growth adaptation per scope).
"""
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"; CHARTS.mkdir(exist_ok=True)
OUT = HERE / "RKLB-Competitive-Analysis.pptx"

NAVY = RGBColor(0x17, 0x36, 0x5D); ORANGE = RGBColor(0xF2, 0x6A, 0x21)  # RKLB orange
CHARCOAL = RGBColor(0x40, 0x40, 0x40); GREY = RGBColor(0xBF, 0xBF, 0xBF)
LGREY = RGBColor(0xF2, 0xF2, 0xF2); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B); GREEN = RGBColor(0x2E, 0x7D, 0x32)
FONT = "Times New Roman"

plt.rcParams.update({"font.family":"serif","font.serif":["DejaVu Serif"],
    "axes.edgecolor":"#404040","axes.labelcolor":"#404040","xtick.color":"#404040","ytick.color":"#404040",
    "axes.grid":True,"grid.color":"#E0E0E0","grid.linewidth":0.5})


def chart_evrev():
    fig, ax = plt.subplots(figsize=(9, 4.6), dpi=150)
    peers = ["IRDM", "LMT", "RDW", "PL", "SpaceX*", "RKLB", "ASTS"]
    vals = [7.5, 1.85, 9.3, 49, 95, 110, 410]
    colors = ["#7F8C8D", "#7F8C8D", "#7F8C8D", "#7F8C8D", "#2E7D32", "#F26A21", "#7F8C8D"]
    bars = ax.barh(peers, vals, color=colors, edgecolor="#202020", linewidth=0.5)
    for b, v in zip(bars, vals):
        ax.text(b.get_width()+8, b.get_y()+b.get_height()/2, f"{v:.0f}x", va="center", fontsize=10, fontweight="bold")
    ax.set_xlabel("EV / TTM Revenue (x)")
    ax.set_title("EV/Revenue — RKLB sits in the upper-growth zone (only ASTS higher)", fontsize=12.5, color="#17365D", pad=10)
    ax.set_xlim(0, 480)
    fig.tight_layout(); out = CHARTS / "evrev.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


def chart_neutron():
    fig, ax = plt.subplots(figsize=(9, 4.6), dpi=150)
    yrs = ["FY25A", "FY26E", "FY27E", "FY28E", "FY29E", "FY30E", "FY31E"]
    elec = [21, 22, 28, 30, 32, 32, 32]
    neut = [0, 1, 5, 12, 16, 18, 20]
    x = range(len(yrs))
    ax.bar(x, elec, color="#7F8C8D", label="Electron + HASTE", edgecolor="#202020", linewidth=0.4)
    ax.bar(x, neut, bottom=elec, color="#F26A21", label="Neutron", edgecolor="#202020", linewidth=0.4)
    for i, (e, n) in enumerate(zip(elec, neut)):
        ax.text(i, e+n+1, f"{e+n}", ha="center", fontsize=9, fontweight="bold")
    ax.set_xticks(list(x)); ax.set_xticklabels(yrs)
    ax.set_ylabel("Launches per year")
    ax.set_title("Launch cadence — Neutron debut Q4 2026; ~20/yr at mature cadence (vs SpaceX Falcon 9 ~140)",
                 fontsize=12, color="#17365D", pad=10)
    ax.legend(loc="upper left", fontsize=10)
    ax.set_ylim(0, 60)
    fig.tight_layout(); out = CHARTS / "neutron.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


EVR_PNG = chart_evrev(); NEUT_PNG = chart_neutron()

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(45720); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    runs = [text] if isinstance(text, str) else text
    for i, line in enumerate(runs):
        if i == 0: r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align; r = p2.add_run()
        r.text = line; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def slide_title(slide, title, *, size=22):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95), title, size=size, bold=True, color=NAVY)


def slide_footer(slide, txt):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3), txt, size=9, italic=True, color=GREY)


def add_table(slide, left, top, width, height, data, *, col_widths=None, font_size=12):
    rows = len(data); cols = len(data[0])
    shp = slide.shapes.add_table(rows, cols, left, top, width, height); tbl = shp.table
    if col_widths:
        for i, w in enumerate(col_widths): tbl.columns[i].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT; run.font.size = Pt(font_size)
                    if ri == 0: run.font.bold = True; run.font.color.rgb = WHITE
                    elif ci == 0: run.font.bold = True; run.font.color.rgb = NAVY
                    else: run.font.color.rgb = CHARCOAL
            if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else: cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Emu(45720); cell.margin_top = cell.margin_bottom = Emu(27432)
    return shp


# SLIDE 1 — COVER
s = prs.slides.add_slide(BLANK)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.5))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.12))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ORANGE; stripe.line.fill.background()
add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5), "NASDAQ: RKLB — Investment Decision Frame", size=18, italic=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.4), "Rocket Lab USA", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.9), "Competitive Landscape, Comps & DCF (small-cap-growth)", size=24, color=ORANGE)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5), "As of June 3, 2026", size=14, italic=True, color=CHARCOAL)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.6),
         ["Peers: ASTS · PL · IRDM · RDW  +  LMT (old-space anchor)  +  SpaceX (private memo)",
          "Built with deep-research + competitive-analysis + comps-analysis + 3-statement-model + dcf-model skills",
          "Companions: RKLB-Comps-Analysis.xlsx · RKLB-Model.xlsx"],
         size=13, color=CHARCOAL)

# SLIDE 2 — THESIS / HEADLINE FINDING
s = prs.slides.add_slide(BLANK)
slide_title(s, "Best-in-class small-cap space pure-play — but at ~$123 you're paying for SpaceX-lite execution")
bullets = [
    ("1. The franchise", "RKLB Q1 CY26: $200M (+63%), GAAP GM 38.2%, $2.22B backlog (+108% YoY), $1.48B cash. Bifurcated: Launch Services (Electron 21/yr) + Space Systems (~67% of rev, incl. Mynaric optical comms closed April 2026)."),
    ("2. The Neutron option", "Q4 2026 NET first launch (slipped from mid-2026 after Q4 2025 tank-test failure). 5+ Neutron contracts booked Q1 26 alone; 5 more in the May 'mystery customer' 8-launch deal. Target $50-55M/launch at ~50% mature GM. 5 in FY27 → 18-20/yr at mature cadence."),
    ("3. The valuation problem", "Stock has rocketed to $123 (~$75B mkt cap, ~110x TTM rev, ~6x from $20-30 in early 2025). Analyst consensus PT $103.91 — 16% BELOW spot. Our base-case DCF: $9.35/share (-92% vs $123). The market is pricing SpaceX-lite optionality beyond FY31."),
    ("4. The read", "Wonderful business, broken price for fundamentals-driven entry. Long for momentum/optionality, NOT for intrinsic value. Watch Neutron Q4 26 test stand cadence — slip = mean-reversion catalyst. Pair vs RDW for cleaner risk (same theme, 9x vs 110x EV/Rev)."),
]
y = Inches(1.35)
for h, b in bullets:
    add_text(s, Inches(0.6), y, Inches(2.4), Inches(1.2), h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(3.1), y, Inches(9.6), Inches(1.2), b, size=11.5, color=CHARCOAL)
    y += Inches(1.36)
slide_footer(s, "Sources: RKLB Q1 CY26 IR (May 7, 2026); StockAnalysis (June 2 close); RKLB-Model.xlsx. Research, NOT investment advice.")

# SLIDE 3 — THE QUESTION
s = prs.slides.add_slide(BLANK)
slide_title(s, "The question this analysis answers")
add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(1.2),
         '"Rocket Lab is the highest-quality small-cap space pure-play with backlog up +108% YoY and Neutron debut due Q4 2026 — but the stock has 6x\'d to $123 and the analyst consensus PT is BELOW spot. '
         'What is the intrinsic value through the Neutron ramp, and how should that frame an entry?"',
         size=18, italic=True, color=NAVY)
add_text(s, Inches(0.7), Inches(3.9), Inches(12), Inches(0.4), "Five lenses:", size=14, bold=True, color=CHARCOAL)
for i, ln in enumerate([
    "(1) Market — space launch + space infra demand; sovereign + commercial mix.",
    "(2) Competitive — vs SpaceX (private, dominant) + ASTS, PL, IRDM, RDW.",
    "(3) Financial (comps) — EV/Revenue the cleaner gauge in a pre-profit cohort.",
    "(4) Intrinsic (DCF) — small-cap-growth DCF with wide bands; conservative case.",
    "(5) Decision — express the theme without paying ~110x trailing revenue.",
]):
    add_text(s, Inches(0.9), Inches(4.4)+i*Inches(0.46), Inches(12), Inches(0.4), ln, size=13, color=CHARCOAL)

# SLIDE 4 — MARKET
s = prs.slides.add_slide(BLANK)
slide_title(s, "Space market is consolidating around SpaceX + a handful of pure-plays; sovereign demand is the tailwind")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.6), [
    ["Segment", "TAM / state", "Key dynamic", "RKLB position"],
    ["Launch (small/medium)", "~$15-20B/yr, growing", "SpaceX dominates Falcon 9 (~140/yr); ride-share commoditizing small payloads", "Electron leader in small (~$10M); Neutron targets medium ($50M)"],
    ["Space Systems / Components", "~$25B+ (sat manufacturing, subsystems)", "Verticalization (SpaceX in-house); growing 3rd-party supply", "RKLB Photon, solar arrays (SolAero), star trackers, Mynaric optical comms"],
    ["Satellite Communications", "~$50B (incl. Starlink dominance)", "Starlink consolidating LEO comms (>10M subs); IRDM serves narrowband", "Not direct (RKLB is supplier, not operator)"],
    ["US Gov't / Defense", "~$25B+ (DoD + NASA + NRO)", "SDA Tranche 2 + USAF/AFRL Rocket Cargo + classified launches", "$1.3B SDA T2 Beta contract; AFRL Rocket Cargo NET 2026"],
    ["Total addressable", "$100B+ commercial+gov't 2030E", "SpaceX captures launch + comms; opportunity in components + non-SpaceX launch", "Top 5 customers = 49% revenue (heavy USG)"],
], col_widths=[Inches(3.0), Inches(2.7), Inches(3.5), Inches(3.1)], font_size=11)
slide_footer(s, "Sources: Industry estimates; SpaceX S-1 (April 2026); SDA + AFRL contract announcements; RKLB 10-K customer concentration disclosure.")

# SLIDE 5 — RKLB PROFILE
s = prs.slides.add_slide(BLANK)
slide_title(s, "Rocket Lab profile — Q1 26 record, $2.22B backlog, $1.48B cash, Neutron Q4 NET")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0), [
    ["Metric", "Value", "Note"],
    ["Q1 CY26 revenue", "$200.3M (+63.5% YoY)", "Beat consensus by ~$20M; record quarter"],
    ["  Launch Services", "$63.7M (+78.9%)", "6 Electron+HASTE missions"],
    ["  Space Systems", "$136.7M (+57.2%)", "~68% of revenue; Photon + Mynaric"],
    ["Q1 26 gross margin", "38.2%", "Improving from FY25 GAAP ~30%"],
    ["Q1 26 adj. EBITDA / NI", "-$11.8M / -$45M", "Narrowing losses; -$0.07 EPS"],
    ["Backlog (Q1 26)", "$2.22B (+108% YoY, +20.2% QoQ)", "~36% deliverable within 12 months"],
    ["Cash + marketable securities", "$1.48B", "+$2B total liquidity; FY25 ATM raised $1.146B"],
    ["FY26 H1 (Q1 actual + Q2 guide mid)", "~$432M ($200M + $232.5M)", "Q2 guide +12% above $208M consensus"],
    ["Mynaric (closed April 14, 2026)", "$155.3M (cash + 2.28M shares)", "Optical inter-satellite laser comms; Rocket Lab Europe"],
    ["Market cap / share / 52-wk high", "~$75B / $123.32 / $150.23 (May 27)", "Up ~5-6x from $20-30 in early 2025"],
], col_widths=[Inches(3.7), Inches(3.6), Inches(5.0)], font_size=10.5)
slide_footer(s, "Sources: RKLB Q1 CY26 8-K / 10-Q (May 7, 2026); FY25 10-K; press releases on Mynaric close.")

# SLIDE 6 — NEUTRON RAMP CHART
s = prs.slides.add_slide(BLANK)
slide_title(s, "Neutron is the thesis — Q4 2026 debut, ~20 launches/yr at mature cadence")
s.shapes.add_picture(str(NEUT_PNG), Inches(0.5), Inches(1.35), width=Inches(8.4))
add_text(s, Inches(9.1), Inches(1.4), Inches(4.0), Inches(0.5), "Why this matters", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.1), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• Neutron list price $50-55M target — vs SpaceX Falcon 9 ~$67M.",
          "",
          "• 5+ Neutron contracts booked Q1 26 (record); 5 more in May 'mystery customer' 8-launch deal.",
          "",
          "• Mature ~20/yr (vs SpaceX ~140 Falcon 9): RKLB targets 15% of medium-lift market = ~$1B/yr at scale.",
          "",
          "• Q4 2026 first launch will be EXPENDABLE (first stage). Reuse target ~3-5 flights from late 2027+.",
          "",
          "• Engine qualification complete; integrated systems test next."],
         size=10.5, color=CHARCOAL)
slide_footer(s, "Sources: Spaceflight Now (May 2026 'mystery customer' deal); RKLB Q1 CY26 IR; engine test commentary.")

# SLIDE 7 — PEER COMPARISON CHART
s = prs.slides.add_slide(BLANK)
slide_title(s, "EV/Revenue — RKLB at ~110x sits in the upper-growth zone (ASTS only higher; SpaceX ~95x)")
s.shapes.add_picture(str(EVR_PNG), Inches(0.5), Inches(1.35), width=Inches(8.4))
add_text(s, Inches(9.1), Inches(1.4), Inches(4.0), Inches(0.5), "Reading", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.1), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• IRDM (7.5x) = mature cash-flow anchor; RKLB will land here if growth stalls.",
          "",
          "• RDW (9x) = closest growth comp; pre-profit; same theme, much cheaper.",
          "",
          "• PL (49x) = subscription EO premium.",
          "",
          "• SpaceX (95x) = the upper bound for a credible operator.",
          "",
          "• RKLB (110x) PRICES AS IF it's better than SpaceX on terminal value.",
          "",
          "• ASTS (410x) is pre-revenue option — not a usable anchor."],
         size=10.5, color=CHARCOAL)
slide_footer(s, "Source: RKLB-Comps-Analysis.xlsx. SpaceX EV/Rev derived from $1.75T IPO target / $18.7B FY25 revenue.")

# SLIDE 8 — PEER DEEP DIVE
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep-dive: what each peer tells you about RKLB's valuation")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ["Peer", "Mkt cap", "TTM rev / growth", "EV/Rev", "What it tells you about RKLB"],
    ["IRDM", "~$5.25B", "$887M / +2%", "~7.5x", "Mature-state anchor. If RKLB stops growing → P/S compresses ~15x."],
    ["PL", "~$15.4B", "$308M / +26%", "~49x", "First adj. EBITDA+ year. Subscription EO model — RKLB doesn't have."],
    ["RDW", "~$4.1B", "$432M / +58%", "~9x", "Closest growth comp (space infra, pre-profit). Same theme, 1/12 the multiple."],
    ["ASTS", "~$35B", "$85M / pre-rev", "~410x", "Pre-revenue option. Not a usable anchor."],
    ["LMT", "~$121B", "$75.1B / flat", "~1.85x", "Old-space proxy. Space segment ~$13B/yr at ~10% op margin = the floor."],
    ["SpaceX*", "$1.75T target", "$18.7B / +43%", "~95x", "PRIVATE benchmark (S-1 filed April 2026). The upper bound for a credible operator."],
], col_widths=[Inches(1.2), Inches(1.3), Inches(2.4), Inches(1.2), Inches(6.2)], font_size=11)
slide_footer(s, "*SpaceX is private; multiple computed from $1.75T IPO target (Bloomberg) and FY25 S-1 disclosed revenue.")

# SLIDE 9 — VALUATION (FOOTBALL FIELD)
s = prs.slides.add_slide(BLANK)
slide_title(s, "Valuation — the gap between DCF ($9) and market ($123) is the SpaceX-lite optionality")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.4), [
    ["Method", "Implied $/share range", "vs $123.32"],
    ["DCF — base (WACC 13% / g 3.5%)", "~$9", "−92%"],
    ["DCF — sensitivity grid (WACC 11-15%)", "$8 – $14", "Entire grid below price"],
    ["DCF — extended bull (10% terminal growth)", "~$25-40 [E]", "−70% to −80%"],
    ["EV/Revenue — RDW comp (~9x × $679M)", "~$10", "−92%"],
    ["EV/Revenue — PL comp (~49x × $679M)", "~$54", "−56%"],
    ["EV/Revenue — SpaceX comp (~95x × $679M)", "~$104", "−16% (≈ analyst consensus PT $103.91)"],
    ["Analyst consensus PT (S&P, 18 analysts)", "$103.91", "−15.7% — Street already skeptical"],
], col_widths=[Inches(5.4), Inches(3.5), Inches(3.4)], font_size=11)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8),
         ["• Only the SpaceX-comp lens ($104) gets near the current price — and that's against a still-private benchmark.",
          "• DCF is structurally bearish for pre-profit names: heavy in TV, conservative on post-FY31 growth.",
          "• Reality check: at $123 you're paying for ~10%+ annual growth through 2040+. That's SpaceX-lite execution priced in."],
         size=11, color=CHARCOAL)
slide_footer(s, "Source: RKLB-Comps-Analysis.xlsx + RKLB-Model.xlsx. Implied $/share = method × shares ($679M TTM / 608M shares).")

# SLIDE 10 — DCF DETAIL
s = prs.slides.add_slide(BLANK)
slide_title(s, "DCF detail — Neutron commercial revenue FY27+; first profitable year FY28; FY31 uFCF $642M")
add_table(s, Inches(0.5), Inches(1.4), Inches(5.6), Inches(4.6), [
    ["DCF bridge (base)", "$M / $"],
    ["PV of explicit FCF (FY27-31)", "609"],
    ["Terminal value (Gordon)", "6,779"],
    ["PV of terminal value", "3,632"],
    ["Enterprise value", "4,240"],
    ["(+) Net cash", "1,442"],
    ["Equity value", "5,682"],
    ["÷ Diluted shares (M)", "608"],
    ["Implied value / share", "$9.35"],
    ["Current price", "$123.32"],
    ["Upside / (downside)", "(92.4%)"],
], col_widths=[Inches(3.4), Inches(2.2)], font_size=11)
add_text(s, Inches(6.4), Inches(1.4), Inches(6.4), Inches(0.4), "Sensitivity — implied $/share (WACC × g)", size=12, bold=True, color=NAVY)
add_table(s, Inches(6.4), Inches(1.85), Inches(6.5), Inches(3.0), [
    ["WACC ↓ / g →", "2.5%", "3.0%", "3.5%", "4.0%", "4.5%"],
    ["11.0%", "11", "12", "12", "13", "14"],
    ["12.0%", "10", "10", "11", "11", "12"],
    ["13.0% (base)", "9", "9", "10", "10", "10"],
    ["14.0%", "8", "8", "9", "9", "9"],
    ["15.0%", "8", "8", "8", "8", "9"],
], col_widths=[Inches(1.7), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95)], font_size=10.5)
add_text(s, Inches(6.4), Inches(5.0), Inches(6.5), Inches(1.9),
         ["• Wider WACC band (11-15%) reflects small-cap + pre-profit risk; beta 1.8.",
          "• TV is 86% of EV — every $ of post-FY31 growth is amplified.",
          "• To bridge to $123 you'd need terminal growth ~10%+ for many years (SpaceX-lite path)."],
         size=11, color=CHARCOAL)
slide_footer(s, "Source: RKLB-Model.xlsx (validated, 225 formulas, 0 errors).")

# SLIDE 11 — BULL/BASE/BEAR
s = prs.slides.add_slide(BLANK)
slide_title(s, "Bull / base / bear — anchored to Neutron execution + sovereign demand")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.8), [
    ["Scenario", "Prob.", "Key drivers", "Implied value"],
    ["Bull (SpaceX-lite)", "20%", "Neutron mature 30+ launches/yr; sovereign demand drives $5B+ revenue by FY30; 25%+ mature margin",
     "~$80-130 (matches/beats price; extended bull DCF)"],
    ["Base (Neutron ramps as planned)", "55%", "Neutron 18-20/yr at FY31; 15% mature op margin; Space Systems compounds 25%/yr",
     "~$9-15 (DCF); ~$50-100 on EV/Rev comps"],
    ["Bear (Neutron slips/fails)", "25%", "Neutron debut delayed past 2027 or first 3-5 flights have RUD; SDA T2 underwhelms; ATM dilution restarts",
     "~$30-50 (mean-reversion to RDW comp); meaningful downside"],
], col_widths=[Inches(1.6), Inches(0.9), Inches(5.6), Inches(4.2)], font_size=11)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4), "Signposts", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.9),
         "1) Neutron integrated systems test (H2 2026)  2) First Neutron launch (Q4 2026 NET)  3) Additional Neutron contract bookings  4) ATM raise activity (signals cash burn)  5) DoD NSSL Lane 1 acceptance",
         size=12, color=CHARCOAL)
slide_footer(s, "Probabilities are author's judgment. RKLB IR not provided full-year guidance; consensus $900M-$1.0B FY26 (range).")

# SLIDE 12 — DECISION FRAME
s = prs.slides.add_slide(BLANK)
slide_title(s, "Decision frame — four ways to express (or pass on) the Rocket Lab thesis")
opts = [
    ("Pass at $123 — wait for a pullback", "DCF says $9-15; even the SpaceX-lite optionality only justifies ~$100-130. Analyst PT $104 below spot. Wait for Neutron miss/scare or general space-stock correction.", NAVY),
    ("Pair trade: long RDW / short RKLB", "Same growth theme (space infra, pre-profit, US gov't-dependent); RDW at 9x vs RKLB at 110x EV/Rev — relative-value short.", ORANGE),
    ("Express the theme via SpaceX IPO", "When SpaceX IPOs (potentially June 2026 onward), get the dominant operator at ~95x EV/Rev — vs RKLB at 110x without the dominance.", CHARCOAL),
    ("Hold a small position — story-stock optionality", "If you owned RKLB lower, hold a portion for Neutron event-cadence + sovereign tailwinds. Trim into strength; don't add at $123.", GREY),
]
y = Inches(1.5)
for t, b, col in opts:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(11.8), Inches(0.35), t, size=13, bold=True, color=col)
    add_text(s, Inches(0.85), y+Inches(0.38), Inches(11.8), Inches(0.75), b, size=11, color=CHARCOAL)
    y += Inches(1.22)
slide_footer(s, "Research / decision-framing only — NOT investment advice. Cross-check against a primary terminal before acting.")

# SLIDE 13 — RISKS
s = prs.slides.add_slide(BLANK)
slide_title(s, "Key risks — single product (Neutron), customer concentration, valuation")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.4), [
    ["Risk", "Severity", "Detail"],
    ["Neutron execution", "HIGH", "Q4 2026 NET (slipped from mid-2026); first 1-3 flights expendable. RUD on inaugural = stock -30%+."],
    ["Customer concentration", "HIGH", "Top 5 customers = 49% FY25 revenue; heavy USG (SDA, USAF, NASA, NRO, DARPA). Continuing-resolution risk."],
    ["Valuation", "HIGH", "~110x TTM EV/Rev; analyst PT $104 (-16%); embeds ~10%+ growth through 2040+."],
    ["Cash burn / dilution", "MEDIUM", "FY25 ATM raised $1.146B; Q1 26 cash $1.48B — adequate but Neutron build + Mynaric integration consume."],
    ["SpaceX competition", "MEDIUM", "Falcon 9 ~$67M can undercut Neutron's ~$50-55M target; Starship adds heavy-lift competition."],
    ["First Electron failure", "LOW-MED", "21-for-21 in FY25; any future RUD impacts customer trust + insurance."],
    ["Mynaric integration", "LOW-MED", "First international ops (Germany); FX, ITAR/EAR complexity, Mynaric prior cash burn."],
], col_widths=[Inches(2.6), Inches(1.3), Inches(8.4)], font_size=11)
slide_footer(s, "Sources: RKLB FY25 10-K risk factors; Bleecker Street short thesis; Space Insider; press coverage of Neutron delay (Spaceflight Now).")

# SLIDE 14 — SOURCES
s = prs.slides.add_slide(BLANK)
slide_title(s, "Primary sources & data-gap caveats")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4), "Primary sources", size=14, bold=True, color=NAVY)
srcs = [
    "Rocket Lab Q1 CY26 8-K / 10-Q (May 7, 2026); FY25 10-K (Feb 26, 2026); FY24 10-K (Feb 2025)",
    "Rocket Lab IR press releases — Mynaric close (April 14, 2026), 'mystery customer' 8-launch deal (May 7, 2026)",
    "Spaceflight Now / SpaceNews / NASASpaceFlight — Neutron timeline coverage; tank-test failure (Nov 2025); Q4 2026 NET",
    "Peer Q1 26 8-Ks/IR: ASTS (May 11), PL (Mar 19 FY26), IRDM (Apr), RDW (May 8), LMT (Apr)",
    "SpaceX S-1 (April 2026 confidential; first financials disclosed); Bloomberg / Reuters / Satellite Today coverage",
    "Stock data: StockAnalysis, Public.com, Macrotrends, Yahoo Finance (~June 1-3, 2026)",
    "Analyst consensus: S&P Global poll of 18 analysts (PT $103.91, mean Buy)",
]
y = Inches(1.85)
for src in srcs:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.3), "• " + src, size=11, color=CHARCOAL); y += Inches(0.30)
add_text(s, Inches(0.5), y + Inches(0.1), Inches(12.3), Inches(0.4), "Caveats", size=14, bold=True, color=NAVY)
y += Inches(0.55)
caveats = [
    "MCP terminal connectors NOT configured. All data from SEC filings + reputable secondary sources, [E]-flagged.",
    "Pre-profit DCF: wide WACC bands (11-15%); model is structurally bearish for story stocks. Use EV/Rev as the cross-check.",
    "Neutron pricing/cadence are management commentary, not contracts — execution-dependent.",
    "Mynaric (NOT MDA Space) was the 2026 acquisition; integration in Q2-Q3 26 may produce one-time charges.",
    "SpaceX figures are S-1 disclosed but the company is still private; IPO valuation could land $780B-$2T+ range.",
    "Research / decision-framing only — NOT investment advice.",
]
for c in caveats:
    add_text(s, Inches(0.7), y, Inches(12.3), Inches(0.5), "• " + c, size=11, color=CHARCOAL); y += Inches(0.36)
slide_footer(s, "Built per: deep-research + competitive-analysis + comps-analysis + 3-statement-model + dcf-model + audit-xls + ib-check-deck.")

prs.save(OUT)
print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
