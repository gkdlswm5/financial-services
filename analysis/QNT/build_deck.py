"""
Build QNT-Competitive-Analysis.pptx — Quantinuum (NASDAQ: QNT) competitive
analysis deck, following the `competitive-analysis` skill.
"""
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"
CHARTS.mkdir(exist_ok=True)
OUT = HERE / "QNT-Competitive-Analysis.pptx"

NAVY = RGBColor(0x17, 0x36, 0x5D)
MIDBLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHTBLUE = RGBColor(0xD9, 0xE1, 0xF2)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHTGREY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x86, 0x40)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
ORANGE = RGBColor(0xED, 0x7D, 0x31)

FONT = "Times New Roman"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=NAVY, indent=0,
                first_bold=False, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        p.line_spacing = line_spacing
        text = item if isinstance(item, str) else item[0]
        bold = first_bold and i == 0
        if isinstance(item, tuple):
            bold = item[1] if len(item) > 1 else False
        r = p.add_run()
        r.text = "• " + text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_header(slide, title, subtitle=None, num=None):
    add_rect(slide, 0, 0, 13.333, 0.6, fill=NAVY)
    add_text(slide, 0.3, 0.05, 11, 0.5, title, size=22, bold=True,
             color=WHITE, font=FONT)
    if num:
        add_text(slide, 12.3, 0.05, 1, 0.5, num, size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_text(slide, 0.3, 0.7, 12.7, 0.4, subtitle, size=12,
                 color=GREY, font=FONT)


def add_footer(slide, n):
    add_rect(slide, 0, 7.2, 13.333, 0.3, fill=LIGHTGREY)
    add_text(slide, 0.3, 7.22, 6, 0.25,
             "Quantinuum (QNT) — Competitive Analysis | Jun 4, 2026",
             size=9, color=GREY)
    add_text(slide, 12.5, 7.22, 0.7, 0.25, f"{n}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, data, col_widths_in, row_h=0.32,
              header_fill=NAVY, header_color=WHITE, body_size=10,
              header_size=10, zebra=True, first_col_bold=False):
    """data[0] = headers, data[1:] = rows. Each cell is a string or (str, color)."""
    n_rows = len(data); n_cols = len(data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                       Inches(sum(col_widths_in)), Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for i, w in enumerate(col_widths_in):
        tbl.columns[i].width = Inches(w)
    for ri in range(n_rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            val = data[ri][ci]
            fg = None
            if isinstance(val, tuple):
                val, fg = val
            cell.text = str(val)
            tf = cell.text_frame
            tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
            tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(header_size if ri == 0 else body_size)
                run.font.bold = (ri == 0) or (first_col_bold and ci == 0)
                if ri == 0:
                    run.font.color.rgb = header_color
                elif fg is not None:
                    run.font.color.rgb = fg
                else:
                    run.font.color.rgb = NAVY
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif zebra and ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHTGREY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return tbl_shape


# =====================================================================
# Generate charts upfront
# =====================================================================
plt.rcParams.update({
    "font.family": "serif", "font.serif": ["Times New Roman", "DejaVu Serif"],
    "axes.edgecolor": "#17365D", "axes.labelcolor": "#17365D",
    "xtick.color": "#17365D", "ytick.color": "#17365D",
})


def chart_revenue_ramp():
    years = ["FY24A", "FY25A", "FY26E", "FY27E", "FY28E", "FY29E",
             "FY30E", "FY31E", "FY32E", "FY33E", "FY34E", "FY35E"]
    riken = [14.5,18.5,4,2,1,0.5,0,0,0,0,0,0]
    helios = [4,5,22,75,180,350,580,850,1180,1550,1980,2480]
    saas = [2.5,4,8,16,32,60,100,160,240,340,460,600]
    gov = [2,3.4,6,12,20,30,45,65,90,120,160,200]
    fig, ax = plt.subplots(figsize=(8.5, 4.3), dpi=160)
    bottom = np.zeros(len(years))
    for label, vals, col in [
        ("RIKEN (winding down)", riken, "#C0392B"),
        ("Helios cloud + on-prem", helios, "#17365D"),
        ("Quantum Origin + InQuanto SaaS", saas, "#4472C4"),
        ("Government / DARPA / CHIPS", gov, "#2E8640"),
    ]:
        ax.bar(years, vals, bottom=bottom, label=label, color=col,
               edgecolor="white", linewidth=0.5)
        bottom += np.array(vals)
    ax.set_title("QNT revenue build — 4 streams ($M)", fontsize=12, color="#17365D")
    ax.legend(loc="upper left", fontsize=8, frameon=False)
    ax.set_ylabel("Revenue ($M)", fontsize=9)
    ax.tick_params(axis='x', rotation=45, labelsize=8)
    ax.tick_params(axis='y', labelsize=8)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.axvspan(-0.5, 1.5, alpha=0.08, color="grey")
    ax.text(0.5, ax.get_ylim()[1] * 0.95, "Reported", fontsize=8,
            color="#595959", ha="center")
    plt.tight_layout()
    out = CHARTS / "revenue_ramp.png"
    plt.savefig(out, bbox_inches="tight"); plt.close()
    return str(out)


def chart_football_field():
    methods = [
        ("IONQ-anchored\n(~172x EV/FY25)", 45, 50, 55),
        ("Peer median\n(~383x)", 50, 56, 62),
        ("Traditional DCF\n(WACC 15.85%)", 18, 25, 38),
        ("Reverse-engineered\nat $50B TAM, 12%", 55, 68, 85),
        ("QBTS-anchored\n(~1023x EV/FY25)", 110, 140, 165),
        ("Bull case\n(McKinsey high)", 75, 105, 135),
    ]
    fig, ax = plt.subplots(figsize=(9, 4.5), dpi=160)
    y = np.arange(len(methods))
    for i, (label, lo, mid, hi) in enumerate(methods):
        col = "#17365D" if "DCF" not in label else "#C0392B"
        ax.barh(i, hi - lo, left=lo, color=col, alpha=0.65, height=0.55)
        ax.plot([mid], [i], "D", color="white", markersize=8,
                markeredgecolor=col, markeredgewidth=1.5)
        ax.text(hi + 2, i, f"${mid}", fontsize=9, va="center", color=col)
        ax.text(lo - 2, i, f"${lo}", fontsize=8, va="center",
                ha="right", color="#595959")
    ax.axvline(68, color="#ED7D31", linestyle="--", linewidth=1.5)
    ax.text(68, len(methods) - 0.3, "$68 (open)", fontsize=9,
            color="#ED7D31", ha="center")
    ax.set_yticks(y); ax.set_yticklabels([m[0] for m in methods], fontsize=9)
    ax.set_xlabel("Implied price per QNT share ($)", fontsize=10)
    ax.set_title("Football field — QNT implied valuation", fontsize=12,
                 color="#17365D")
    ax.set_xlim(0, 200)
    ax.invert_yaxis()
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    plt.tight_layout()
    out = CHARTS / "football_field.png"
    plt.savefig(out, bbox_inches="tight"); plt.close()
    return str(out)


def chart_modality_2x2():
    fig, ax = plt.subplots(figsize=(7.5, 5.5), dpi=160)
    # x = revenue scale (log proxy), y = technical maturity (1-10)
    data = [
        # (name, x, y, color, size)
        ("QNT (Helios)", 31, 8.5, "#17365D", 600),
        ("IONQ (Forte/Tempo)", 130, 7.5, "#4472C4", 700),
        ("RGTI (108Q)", 7, 6.5, "#ED7D31", 400),
        ("QBTS (Annealer)", 25, 5.5, "#999999", 500),
        ("PsiQuantum (priv.)", 3, 6.0, "#888888", 300),
        ("IBM Q (research)", 0.5, 8.0, "#888888", 250),
        ("Google QAI (research)", 0.5, 9.0, "#888888", 250),
        ("Atom Computing (priv.)", 5, 5.5, "#888888", 250),
    ]
    for name, x, y, col, sz in data:
        ax.scatter(x, y, s=sz, c=col, alpha=0.7, edgecolors="white",
                   linewidths=1.5)
        ax.text(x, y - 0.45, name, fontsize=8, ha="center", color="#17365D")
    ax.axhline(7, color="#595959", linewidth=0.5, alpha=0.3)
    ax.axvline(30, color="#595959", linewidth=0.5, alpha=0.3)
    ax.set_xscale("log")
    ax.set_xlim(0.3, 300)
    ax.set_ylim(4.5, 10)
    ax.set_xlabel("FY25 revenue ($M, log scale)", fontsize=10)
    ax.set_ylabel("Technical maturity / gate fidelity → FTQC ", fontsize=10)
    ax.set_title("2x2 positioning — quantum computing competitors", fontsize=12,
                 color="#17365D")
    ax.text(0.4, 9.5, "Lab leaders\n(no commercial yet)", fontsize=8,
            color="#595959", style="italic")
    ax.text(80, 9.5, "Commercial + technical leaders", fontsize=8,
            color="#595959", style="italic")
    ax.text(0.4, 5.0, "Behind the curve", fontsize=8, color="#595959",
            style="italic")
    ax.text(80, 5.0, "Commercial only", fontsize=8, color="#595959",
            style="italic")
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    plt.tight_layout()
    out = CHARTS / "modality_2x2.png"
    plt.savefig(out, bbox_inches="tight"); plt.close()
    return str(out)


def chart_tam():
    years = [2025, 2028, 2030, 2035, 2040]
    mckinsey_low = [1, 4, 12, 43, 75]
    mckinsey_high = [1.5, 5, 18, 71, 130]
    bcg = [1, 5, 20, 55, 170]
    fig, ax = plt.subplots(figsize=(8, 4), dpi=160)
    ax.fill_between(years, mckinsey_low, mckinsey_high, color="#4472C4",
                    alpha=0.30, label="McKinsey range ($B QC)")
    ax.plot(years, bcg, "o-", color="#C0392B", label="BCG (HW + SW)",
            linewidth=2)
    ax.plot(years, mckinsey_low, "o--", color="#4472C4", linewidth=1)
    ax.plot(years, mckinsey_high, "o--", color="#4472C4", linewidth=1)
    ax.set_title("Quantum computing TAM ($B, 2025-2040)", fontsize=12, color="#17365D")
    ax.set_xlabel("Year", fontsize=10)
    ax.set_ylabel("Annual market size ($B)", fontsize=10)
    ax.legend(fontsize=9, frameon=False)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.grid(True, alpha=0.2)
    plt.tight_layout()
    out = CHARTS / "tam.png"
    plt.savefig(out, bbox_inches="tight"); plt.close()
    return str(out)


def chart_concentration_cliff():
    periods = ["FY2024", "FY2025", "Q1 2025", "Q1 2026"]
    riken_pct = [63, 60, 90, 7]
    fig, ax = plt.subplots(figsize=(8, 4), dpi=160)
    bars = ax.bar(periods, riken_pct, color=["#C0392B" if p > 30 else "#2E8640"
                                              for p in riken_pct],
                  edgecolor="white", linewidth=0.5)
    for bar, p in zip(bars, riken_pct):
        ax.text(bar.get_x() + bar.get_width() / 2, p + 2, f"{p}%",
                ha="center", fontsize=10, fontweight="bold", color="#17365D")
    ax.set_ylabel("RIKEN as % of revenue", fontsize=10)
    ax.set_title("RIKEN concentration cliff — single-customer revenue cliff in Q1'26",
                 fontsize=12, color="#17365D")
    ax.axhline(50, color="grey", linestyle=":", linewidth=1, alpha=0.5)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.set_ylim(0, 100)
    plt.tight_layout()
    out = CHARTS / "riken_cliff.png"
    plt.savefig(out, bbox_inches="tight"); plt.close()
    return str(out)


revenue_chart = chart_revenue_ramp()
ff_chart = chart_football_field()
modality_chart = chart_modality_2x2()
tam_chart = chart_tam()
riken_chart = chart_concentration_cliff()

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = blank()
add_rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
add_text(s, 0.6, 1.0, 12, 1.0, "Quantinuum Inc.", size=44, bold=True,
         color=WHITE)
add_text(s, 0.6, 2.0, 12, 0.6, "NASDAQ: QNT  |  Competitive Analysis, Comps & Dual-DCF Valuation",
         size=18, color=LIGHTBLUE)
add_text(s, 0.6, 4.5, 12, 0.5,
         "Priced $60 (Jun 3, 2026) above $53–55 range  |  Trading $68  |  Post-money ~$15.6B",
         size=15, color=WHITE)
add_text(s, 0.6, 5.0, 12, 0.5,
         "Honeywell parent ~48% voting  |  Pre-profit  |  Apollo (FTQC) target 2029",
         size=13, color=LIGHTBLUE)
add_text(s, 0.6, 6.5, 12, 0.4,
         "Investment workup using 7-skill financial-analysis pipeline",
         size=11, color=LIGHTBLUE)
add_text(s, 0.6, 6.85, 12, 0.4,
         "Decision-framing only. Not investment advice.",
         size=10, color=LIGHTBLUE)
add_text(s, 0.6, 7.15, 12, 0.3,
         "June 4, 2026",
         size=10, color=LIGHTBLUE)

# =====================================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# =====================================================================
s = blank()
add_header(s, "Executive summary — what does $68 mean?", num="2")

add_rect(s, 0.3, 0.95, 12.7, 1.05, fill=LIGHTBLUE)
add_text(s, 0.5, 1.05, 12.3, 0.5,
         "$68 ≈ 2.75× our traditional DCF and ≈ 2.34× the FY35 revenue we model in base case.",
         size=16, bold=True, color=NAVY)
add_text(s, 0.5, 1.5, 12.3, 0.4,
         "Stock prices a 15.4% share of the 2035 quantum-computing TAM. Achievable, but the bar is high.",
         size=13, color=GREY)

# Two-column comparison: Framework 1 vs Framework 2
add_rect(s, 0.3, 2.2, 6.3, 4.5, fill=WHITE, line=NAVY)
add_rect(s, 0.3, 2.2, 6.3, 0.5, fill=NAVY)
add_text(s, 0.5, 2.25, 6, 0.45, "Framework 1 — Traditional DCF",
         size=14, bold=True, color=WHITE)
add_bullets(s, 0.5, 2.85, 6, 3.8, [
    ("WACC: 15.85% (rf 4.3% + beta 2.1 × ERP 5.5%)", True),
    "Terminal growth: 4.0%",
    "12-year horizon (FY24A → FY35E)",
    "FCF crossover: FY2029E (Apollo year)",
    "FY2035E FCF: $1,307M (10-yr CAGR ~75%)",
    "PV of explicit FCF: $1,372M",
    "PV of terminal value: $2,835M",
    "Enterprise value: $4,207M",
    "+ Cash $2,080M = Equity $6,287M",
    ("→ Intrinsic $24.76 / share", True),
    ("→ Stock at +175% premium", True),
], size=11)

add_rect(s, 6.8, 2.2, 6.2, 4.5, fill=WHITE, line=NAVY)
add_rect(s, 6.8, 2.2, 6.2, 0.5, fill=ORANGE)
add_text(s, 7.0, 2.25, 5.9, 0.45, "Framework 2 — Reverse-engineered to $68",
         size=14, bold=True, color=WHITE)
add_bullets(s, 7.0, 2.85, 5.9, 3.8, [
    ("$68 × 253.9M shares = $17.3B mkt cap", True),
    "Less $2.1B cash → $15.2B EV",
    "Assumed FY35 mature EV/Rev: 8.0×",
    "Solve: FY35E rev needed = $7,680M",
    "Model base case FY35E rev: $3,280M",
    "Implied gap: +134% above model",
    "Implied 2035 share of $50B QC TAM: 15.4%",
    "10-yr revenue CAGR required: ~75%",
    "Required: full Apollo (FTQC) success",
    ("→ 'Optionality is mostly priced in'", True),
    ("→ Asymmetry skews DOWN at $68", True),
], size=11)

add_rect(s, 0.3, 6.85, 12.7, 0.3, fill=LIGHTGREY)
add_text(s, 0.5, 6.87, 12.5, 0.25,
         "Decision call: HOLD with downside skew. Buy zone <$45 (peer-median footing). Refresh after Q2'26 print.",
         size=11, bold=True, color=NAVY)

add_footer(s, 2)

# =====================================================================
# SLIDE 3 — THE IPO EVENT
# =====================================================================
s = blank()
add_header(s, "The IPO event — June 3, 2026", num="3",
           subtitle="Priced above range, upsized, ~20x oversubscribed — first major dedicated quantum IPO at scale")

add_rect(s, 0.3, 1.2, 6.3, 5.6, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.3, 6, 0.4, "IPO mechanics", size=14, bold=True)
add_table(s, 0.5, 1.8, [
    ["Item", "Value"],
    ["Initial range (filed)", "$45–50"],
    ["Revised range (May 26)", "$53–55"],
    ["Final priced (Jun 3)", ("$60.00 — above range", GREEN)],
    ["Shares offered (final)", "28,000,000 (upsized from 21M)"],
    ["Gross proceeds", "$1.68B"],
    ["Underwriters", "JPM, MS (joint book) +\nJefferies, Evercore ISI"],
    ["Oversubscription", ("~20× (reported)", GREEN)],
    ["Day-1 open (Jun 4)", ("$68.00 (+13.3% vs IPO)", GREEN)],
    ["Post-money mkt cap (FD)", "~$15.6B"],
    ["Public float", ("~8.3% (very thin)", RED)],
], col_widths_in=[2.4, 3.6], row_h=0.30)

add_rect(s, 6.8, 1.2, 6.2, 5.6, fill=WHITE, line=NAVY)
add_text(s, 7.0, 1.3, 6, 0.4, "Cap-table post-IPO (voting)", size=14, bold=True)
add_bullets(s, 7.0, 1.75, 5.9, 4.5, [
    ("Honeywell: ~48.1% voting power", True),
    "  - Plus Transaction Committee veto on M&A > $10M",
    "  - Plus operating supplier to QNT ('H2 powered by Honeywell')",
    ("Cambridge Quantum Holdings: ~32.5% voting", True),
    "  - Aggregates founder Ilyas Khan (~23%), IBM, JSR",
    ("Combined HON + CQH: ~82% voting", True),
    "Public free float: ~8.3% economic, similar voting",
    "",
    ("Use of proceeds: working capital, R&D,", True),
    "  general corporate purposes. No specific earmark for",
    "  Helios/Sol/Apollo named in S-1.",
    "",
    ("Implied multiple at $60: ~485× EV/FY25 rev", True),
    ("At $68 open: ~491× — peer median ~383× (FY25 basis)", True),
], size=11)

add_footer(s, 3)

# =====================================================================
# SLIDE 4 — COMPANY PROFILE & HONEYWELL PARENT
# =====================================================================
s = blank()
add_header(s, "Company profile — and the Honeywell parent angle", num="4",
           subtitle="QNT is the merged Honeywell Quantum Solutions + Cambridge Quantum (Nov 2021)")

add_rect(s, 0.3, 1.2, 6.3, 5.6, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.3, 6, 0.4, "Quantinuum Inc. (NASDAQ: QNT)", size=14, bold=True)
add_bullets(s, 0.5, 1.8, 5.9, 5.0, [
    ("CEO: Ilyas Khan (founder, ex-Cambridge Quantum)", True),
    "HQ: Cambridge UK + Broomfield CO",
    "Modality: trapped-ion (QCCD architecture)",
    "Tech generations: H1 → H2 → Helios → Sol → Apollo",
    "",
    ("Funding history", True),
    "  Nov '21: Honeywell-Cambridge merger ($300M Honeywell)",
    "  Jan '24: $300M raise @ $5B pre-money (JPMC led)",
    "  Sep '25: $600M raise @ $10B pre-money (NVIDIA +)",
    "  Jun '26: IPO at $60 → post-money ~$15.6B",
    "  Cumulative raised: $2.5B+ over 5 years",
    "",
    ("Financials (S-1)", True),
    "  FY24 revenue $23.0M | net loss $144.1M",
    "  FY25 revenue $30.9M (+34%) | net loss $192.6M",
    "  Q1'26 revenue $5.2M | net loss $136.6M",
    "  Pre-IPO cash $677M → post-IPO ~$2.1B (~4-year runway)",
], size=11)

add_rect(s, 6.8, 1.2, 6.2, 5.6, fill=WHITE, line=NAVY)
add_text(s, 7.0, 1.3, 5.9, 0.4, "Honeywell parent angle (HON)", size=14, bold=True)
add_bullets(s, 7.0, 1.8, 5.9, 5.0, [
    ("Pre-IPO ownership: ~54%", True),
    ("Post-IPO ownership: ~48.1% voting", True),
    "  (sold down slightly through offering)",
    "",
    ("Value to HON shareholders", True),
    "  QNT stake at $68 = ~$7.5B (122M shares × $68)",
    "  HON market cap ~$138B",
    "  → QNT = ~5.4% of HON value",
    "  Modest positive catalyst, NOT transformative",
    "",
    ("Asymmetries for HON shareholders", True),
    "  Upside: monetization optionality (secondary offering)",
    "  Upside: showcase R&D culture",
    "  Risk: ongoing related-party supplier obligations",
    "  Risk: QNT brand reputation flows back to HON",
    "",
    ("HON multi-industrial peers", True),
    "  GE Aerospace, RTX, EMR, ITW",
    "  HON SOTP: QNT carve-out adds ~3-5% to NAV [E]",
], size=11)

add_footer(s, 4)

# =====================================================================
# SLIDE 5 — TAM
# =====================================================================
s = blank()
add_header(s, "Market — quantum computing TAM", num="5",
           subtitle="McKinsey $43-71B by 2035; BCG sees $90-170B HW+SW by 2040")
s.shapes.add_picture(tam_chart, Inches(0.3), Inches(1.1),
                     width=Inches(8.2), height=Inches(5.8))
add_rect(s, 8.8, 1.1, 4.2, 5.8, fill=LIGHTBLUE)
add_text(s, 8.95, 1.2, 4, 0.4, "Why we use $50B", size=13, bold=True)
add_bullets(s, 8.95, 1.6, 4, 5.2, [
    "McKinsey QTM 2026: $43-71B QC by 2035",
    "We use the midpoint ($50B) for the reverse-engineered DCF",
    "",
    ("Splits within $50B:", True),
    "  ~40% hardware ($20B)",
    "  ~35% cloud + services ($17B)",
    "  ~25% software / apps ($12B)",
    "",
    ("Real wild card:", True),
    "  Total economic value $2.7T by 2035",
    "  (downstream usage, not vendor revenue)",
    "  Risk of mistaking 'value created' for",
    "  'value captured' — vendor margin is",
    "  the tractable thesis.",
], size=10)

add_footer(s, 5)

# =====================================================================
# SLIDE 6 — MODALITY LANDSCAPE & 2x2
# =====================================================================
s = blank()
add_header(s, "Competitive landscape — modality + 2x2 positioning", num="6",
           subtitle="Trapped-ion (QNT, IONQ) currently leads on fidelity; superconducting (IBM, Google, RGTI) on scale")
s.shapes.add_picture(modality_chart, Inches(0.3), Inches(1.0),
                     width=Inches(7.5), height=Inches(6.0))

add_rect(s, 8.0, 1.0, 5.0, 6.0, fill=LIGHTGREY)
add_text(s, 8.2, 1.1, 4.8, 0.4, "Modality scorecard", size=13, bold=True)
add_bullets(s, 8.2, 1.5, 4.8, 5.5, [
    ("Trapped-ion (QNT, IONQ)", True),
    "  + Highest gate fidelity in industry",
    "  + Long coherence times",
    "  - Slower gate speeds; scaling unproven beyond ~100q",
    "",
    ("Superconducting (IBM, Google, RGTI)", True),
    "  + Largest qubit counts (Condor 1121q)",
    "  + Fast gates; semis-style scale-up",
    "  - Lower fidelity, shorter coherence",
    "",
    ("Neutral atoms (Atom, QuEra, Pasqal)", True),
    "  + Excellent scaling potential",
    "  + Reconfigurable arrays",
    "  - Earlier commercial maturity",
    "",
    ("Photonic (PsiQuantum, Xanadu)", True),
    "  + Room-temperature operation",
    "  + Excellent for networking",
    "  - Loss problem, gate speed",
    "",
    ("Annealing (QBTS)", True),
    "  + Real customers, optimization niche",
    "  - Not universal gate-model — different category",
], size=10)

add_footer(s, 6)

# =====================================================================
# SLIDE 7 — PEER SET
# =====================================================================
s = blank()
add_header(s, "Peer set — 3 public quantum pure-plays + HON memo", num="7",
           subtitle="As of Jun 3-4, 2026 close")
add_table(s, 0.3, 1.1, [
    ["Ticker", "Modality", "Market cap", "EV", "FY25 rev", "Q1'26 YoY", "EV/FY25", "EV/TTM"],
    [("QNT", NAVY), "Trapped-ion", "$17.3B", "$15.2B", "$30.9M", "—", "~491×", "~466×"],
    ["IONQ", "Trapped-ion", "$25.5B", "$22.4B", "$130.0M", ("+755%", GREEN), "~172×", "~120×"],
    ["RGTI", "Superconducting", "$7.8B", "$7.3B", "$7.1M", ("+199%", GREEN), "~1,023×", "~726×"],
    ["QBTS", "Annealing", "$10.0B", "$9.4B", "$24.6M", ("−81%", RED), "~383×", "~754×"],
    ["HON (memo)", "Conglomerate", "$138B", "$149B", "$38,500M", "+4%", "~3.9×", "~3.9×"],
], col_widths_in=[1.4, 1.7, 1.3, 1.3, 1.3, 1.4, 1.3, 1.4], row_h=0.34, first_col_bold=True)

add_text(s, 0.3, 3.7, 12.7, 0.4,
         "Key observations", size=14, bold=True)
add_bullets(s, 0.3, 4.1, 12.7, 3.0, [
    ("IONQ is the direct comp (trapped-ion, similar TAM exposure) — but has 6× QNT's revenue and 7.5× YoY growth.", True),
    ("QNT trades at a +28% premium to peer median on FY25 revenue (~491× vs ~383×) — but at a −36% DISCOUNT on TTM (~466× vs ~726×).", True),
    ("The 'rich vs peers' narrative is wrong vs RGTI/QBTS, correct only vs IONQ.", True),
    "RGTI revenue is so small ($7M) that any multiple looks distorted — peer-set median has high noise.",
    "QBTS is annealing, not gate-model — not a pure technology comp; trading on Quantum Circuits acquisition optionality.",
    ("HON memo: at $7.5B (~$68), QNT stake is ~5.4% of HON market cap. Modest catalyst.", True),
], size=11)

add_footer(s, 7)

# =====================================================================
# SLIDE 8 — IONQ DEEP DIVE
# =====================================================================
s = blank()
add_header(s, "IONQ — the trapped-ion direct comp", num="8",
           subtitle="More mature commercial revenue; lower multiple; faster growth — the credibility benchmark")

add_rect(s, 0.3, 1.2, 6.3, 5.6, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.3, 6, 0.4, "Financial profile", size=14, bold=True)
add_table(s, 0.5, 1.85, [
    ["Metric", "IONQ"],
    ["Market cap", "$25.5B"],
    ["Enterprise value", "$22.4B"],
    ["FY25 revenue", "$130.0M (+202%)"],
    ["Q1'26 revenue", "$64.7M (+755% YoY)"],
    ["TTM revenue", "$187M"],
    ["Cash & investments", "$3,100M (Q1'26)"],
    ["Q1'26 Adj. EBITDA", "($96.8M)"],
    ["EV/FY25 revenue", "~172×"],
    ["EV/TTM revenue", "~120×"],
    ["Latest system", "Forte / Tempo (5th gen)"],
    ["Customers", "AWS, AstraZeneca, NVIDIA,\nAFRL, ORNL, DARPA (HARQ)"],
], col_widths_in=[2.5, 3.5], row_h=0.30)

add_rect(s, 6.8, 1.2, 6.2, 5.6, fill=WHITE, line=NAVY)
add_text(s, 7.0, 1.3, 6, 0.4, "QNT vs IONQ comparison", size=14, bold=True)
add_bullets(s, 7.0, 1.8, 5.9, 5.0, [
    ("IONQ trades at LESS than half QNT's EV/Rev:", True),
    "  IONQ: ~172× FY25 vs QNT ~491×",
    "  IONQ: ~120× TTM vs QNT ~466×",
    "",
    ("…despite better fundamentals:", True),
    "  6× the revenue base ($130M vs $31M)",
    "  Far higher Q1 growth rate (+755% vs +34% trailing)",
    "  AWS distribution + NVIDIA partnership",
    "  More diversified customer base",
    "",
    ("…and superior tech credentials:", True),
    "  99.99% 2Q fidelity (lab) — best in industry",
    "  #AQ 64 achieved (commercial benchmark)",
    "  DARPA HARQ heterogeneous architecture",
    "",
    ("The IONQ multiple is the 'sanity anchor'", True),
    "  At IONQ's 172× FY25 mult, QNT implied ~$50-55",
    "  Suggests QNT IPO range $53-55 was 'fair vs IONQ'",
    "  Trading at $68 = ~30% premium to that anchor",
], size=11)

add_footer(s, 8)

# =====================================================================
# SLIDE 9 — RGTI + QBTS quick takes
# =====================================================================
s = blank()
add_header(s, "RGTI and QBTS — small-revenue peers (use with care)", num="9",
           subtitle="Tiny denominators distort multiples; both are noise more than signal for QNT pricing")

add_rect(s, 0.3, 1.2, 6.3, 5.6, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.3, 6, 0.4, "Rigetti Computing (RGTI)", size=14, bold=True)
add_bullets(s, 0.5, 1.8, 5.9, 5.0, [
    "Superconducting modality",
    "FY25 revenue: $7.1M (extremely small)",
    "Q1'26 revenue: $4.4M (+199% YoY)",
    "Market cap $7.8B; EV $7.3B",
    "EV / FY25 rev: ~1,023× (very distorted by low denominator)",
    "EV / TTM rev: ~726×",
    "",
    ("Latest milestone:", True),
    "  Cepheus-1-108Q general availability (Apr 2026)",
    "  12 chiplets × 9 qubits each",
    "  99.1% median 2Q fidelity (target 99.5% 2H'26)",
    "",
    ("Roadmap (deferred to 2H'26):", True),
    "  Formal update expected late 2026",
    "  Targeting 'quantum advantage' in ~3 years",
    "",
    "Verdict: RGTI multiple is unusable for sizing QNT; tech is roughly 6-12 months behind QNT Helios",
], size=11)

add_rect(s, 6.8, 1.2, 6.2, 5.6, fill=WHITE, line=NAVY)
add_text(s, 7.0, 1.3, 6, 0.4, "D-Wave Quantum (QBTS)", size=14, bold=True)
add_bullets(s, 7.0, 1.8, 5.9, 5.0, [
    "Quantum annealing modality (NOT universal gate-model)",
    "FY25 revenue: $24.6M (+179% YoY)",
    "Q1'26 revenue: $2.9M (−81% YoY vs Q1'25 system sale)",
    "Market cap $10.0B; EV $9.4B",
    "EV / FY25 rev: ~383×",
    "EV / TTM rev: ~754× (also distorted)",
    "",
    ("Latest system: Advantage2 (May 2025)", True),
    "  4,400+ qubits, Zephyr topology",
    "  20-way connectivity",
    "  2× coherence vs Advantage1",
    "",
    ("Strategic move:", True),
    "  Acquired Quantum Circuits ($250M cash)",
    "  Adds gate-model optionality",
    "  Diversifies modality bet",
    "",
    "Verdict: annealing is a different category from QNT's gate-model. Multiple comp is misleading; technology comp irrelevant.",
], size=11)

add_footer(s, 9)

# =====================================================================
# SLIDE 10 — TECHNOLOGY SCOREBOARD
# =====================================================================
s = blank()
add_header(s, "Technology scoreboard — gate-model quantum computers", num="10",
           subtitle="By 2Q gate fidelity, qubit count, and FTQC roadmap year")

add_table(s, 0.3, 1.1, [
    ["", "QNT (Helios)", "IONQ (Tempo)", "RGTI (Cepheus)", "IBM (Quantum Sys 2)", "Google (Willow)"],
    ["Modality", "Trapped-ion", "Trapped-ion", "Superconducting", "Superconducting", "Superconducting"],
    ["Physical qubits", "98", "64+ AQ", "108", "1,121 (Condor)", "105 (Willow)"],
    ["2Q gate fidelity", ("99.921%", GREEN), "99.99% (lab)", "99.1% (target 99.5%)", "~99.5% (Heron r2)", ("99.4%", GREEN)],
    ["Connectivity", ("All-to-all", GREEN), "All-to-all", "Limited (chiplet)", "Heavy hex", "2D grid"],
    ["Quantum Volume", ("2^25", GREEN), "n/d", "n/d", "1024+", "n/d"],
    ["Logical qubits demoed", ("12 (w/ MSFT)", GREEN), "0 (lab)", "1", "1", ("Below threshold", GREEN)],
    ["FTQC target year", "2029 (Apollo)", "2028 functional", "Unclear", "Beyond 2030", "Beyond 2030"],
    ["Latest commercial system", "Helios (Nov 2025)", "Forte/Tempo (2026)", "Cepheus-108Q (Apr'26)", "Quantum System 2", "Research only"],
    ["Public / Private", "Public (QNT)", "Public (IONQ)", "Public (RGTI)", "IBM segment", "Alphabet segment"],
], col_widths_in=[2.3, 1.85, 1.85, 1.85, 1.85, 1.85], row_h=0.34, first_col_bold=True, body_size=9, header_size=9)

add_rect(s, 0.3, 5.2, 12.7, 1.8, fill=LIGHTBLUE)
add_text(s, 0.5, 5.3, 12.4, 0.4, "Read", size=13, bold=True)
add_bullets(s, 0.5, 5.7, 12.4, 1.3, [
    ("QNT's edge is logical qubits demoed (12 with Microsoft — Sept 2024) — closest to commercial fault-tolerance demonstrations.", True),
    "IONQ has slightly higher peak fidelity in lab but lower commercial track record on logical qubits.",
    "Google Willow proved error-suppression-at-scale (Dec 2024) but is not commercial; QNT-Microsoft is the commercial analog.",
    "QNT roadmap to Apollo (2029) is the most aggressive timeline to fault-tolerant universal QC of any pure-play public name.",
], size=11)

add_footer(s, 10)

# =====================================================================
# SLIDE 11 — RIKEN CONCENTRATION
# =====================================================================
s = blank()
add_header(s, "RIKEN concentration cliff — the underwritten risk", num="11",
           subtitle="Single Japanese government customer was 60-90% of recent revenue; collapsed to 7% in Q1'26")
s.shapes.add_picture(riken_chart, Inches(0.3), Inches(1.1),
                     width=Inches(8.5), height=Inches(5.0))

add_rect(s, 9.0, 1.1, 4.0, 5.6, fill=LIGHTBLUE)
add_text(s, 9.15, 1.2, 3.8, 0.4, "What happened", size=13, bold=True)
add_bullets(s, 9.15, 1.6, 3.8, 5.0, [
    ("RIKEN bought a QNT on-prem system", True),
    "  Lumpy hardware delivery revenue",
    "  +ongoing maintenance/service",
    "",
    ("Q1'26: install complete → 7% of rev", True),
    "  Helios was supposed to fill the gap",
    "  Helios commercial launch Nov 2025",
    "  Launch customers ramping in 2H'26",
    "",
    ("Risk:", True),
    "  Q2/Q3 26 prints could show revenue",
    "  decline YoY before Helios ramp",
    "  Market may react harshly to first",
    "  reported quarter as public co.",
    "",
    ("Mitigant:", True),
    "  Quantum Origin SaaS (recurring)",
    "  Helios cloud + on-prem expanding",
    "  $100M CHIPS LOI (govt)",
], size=10)

add_footer(s, 11)

# =====================================================================
# SLIDE 12 — REVENUE BUILD
# =====================================================================
s = blank()
add_header(s, "Revenue build — 4 streams, 10-year horizon", num="12",
           subtitle="Helios cloud + on-prem becomes the primary growth engine; RIKEN runs off")
s.shapes.add_picture(revenue_chart, Inches(0.3), Inches(1.1),
                     width=Inches(8.5), height=Inches(5.0))

add_rect(s, 9.0, 1.1, 4.0, 5.6, fill=LIGHTBLUE)
add_text(s, 9.15, 1.2, 3.8, 0.4, "Base case assumptions", size=13, bold=True)
add_bullets(s, 9.15, 1.6, 3.8, 5.0, [
    ("RIKEN: runs off to $0 by FY2030E", True),
    "",
    ("Helios cloud + on-prem:", True),
    "  $22M (FY26E) → $580M (FY30E)",
    "  → $2,480M (FY35E)",
    "  ~60% CAGR base case",
    "",
    ("Quantum Origin + InQuanto SaaS:", True),
    "  $8M (FY26E) → $600M (FY35E)",
    "  Higher gross margin (~70%)",
    "",
    ("Government / DARPA / CHIPS:", True),
    "  $6M (FY26E) → $200M (FY35E)",
    "  DARPA QBI Stage B + CHIPS LOI",
    "",
    ("Total FY35E revenue: $3,280M", True),
    "  Implied 2035 QC market share: ~6.6%",
    "  (of $50B mid-case TAM)",
], size=10)

add_footer(s, 12)

# =====================================================================
# SLIDE 13 — FOOTBALL FIELD
# =====================================================================
s = blank()
add_header(s, "Football field — implied price ranges", num="13",
           subtitle="Traditional DCF anchors low; QBTS-comp anchors high; current $68 sits in lower-mid")
s.shapes.add_picture(ff_chart, Inches(0.5), Inches(1.0),
                     width=Inches(8.5), height=Inches(5.6))

add_rect(s, 9.3, 1.0, 3.8, 5.7, fill=LIGHTBLUE)
add_text(s, 9.45, 1.1, 3.6, 0.4, "Method ranges", size=13, bold=True)
add_bullets(s, 9.45, 1.5, 3.6, 5.0, [
    ("Traditional DCF: $18-38", True),
    "  honest intrinsic",
    "  mid $25",
    "",
    ("IONQ-anchored: $45-55", True),
    "  applies IONQ ~172× FY25",
    "",
    ("Peer median: $50-62", True),
    "  median of 3 quantum peers",
    "",
    ("Reverse-engineered: $55-85", True),
    "  what $50B TAM + 12% share gets",
    "",
    ("Bull case: $75-135", True),
    "  McKinsey high + Apollo execution",
    "",
    ("QBTS-anchored: $110-165", True),
    "  applies QBTS ~1,023× FY25",
    "  distorted by tiny denominator",
    "",
    ("→ $68 sits below peer median", True),
    ("→ but well above intrinsic DCF", True),
], size=10)

add_footer(s, 13)

# =====================================================================
# SLIDE 14 — DUAL DCF DETAIL
# =====================================================================
s = blank()
add_header(s, "Dual DCF — traditional vs reverse-engineered", num="14",
           subtitle="Both frameworks land in different places; the gap is the optionality premium")

# Left: traditional DCF detail
add_rect(s, 0.3, 1.0, 6.3, 6.0, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.1, 6, 0.4, "Traditional DCF (intrinsic)", size=14, bold=True)
add_table(s, 0.5, 1.6, [
    ["Line", "Value"],
    ["WACC", "15.85%"],
    ["Terminal growth", "4.0%"],
    ["", ""],
    ["Sum PV(FCF) FY26E-FY35E", "$1,372M"],
    ["Terminal value (Gordon)", "$11,472M"],
    ["PV of TV @ 15.85%, 9.5 yrs", "$2,835M"],
    ["Enterprise value", "$4,207M"],
    ["+ Cash post-IPO", "$2,080M"],
    ["- Debt", "$0"],
    ["Equity value", "$6,287M"],
    ["÷ Shares (M)", "253.9"],
    [("Implied price ($)", NAVY), ("$24.76", RED)],
    ["Current price", "$68.00"],
    [("Premium to intrinsic", NAVY), ("+175%", RED)],
], col_widths_in=[3.3, 2.7], row_h=0.27, body_size=11, first_col_bold=True)

# Right: reverse-engineered detail
add_rect(s, 6.8, 1.0, 6.2, 6.0, fill=WHITE, line=ORANGE)
add_text(s, 7.0, 1.1, 6, 0.4, "Reverse-engineered (at $68 trading)", size=14, bold=True)
add_table(s, 7.0, 1.6, [
    ["Line", "Value"],
    ["Current price ($)", "$68.00"],
    ["× Shares (M)", "253.9"],
    ["= Market cap", "$17,265M"],
    ["- Cash", "$2,080M"],
    ["= EV (today)", "$15,185M"],
    ["", ""],
    ["Mature multiple (FY35E)", "8.0× EV/Rev"],
    ["Discount factor (9.5 yrs)", "4.16×"],
    [("Implied FY35E revenue", NAVY), ("$7,680M", ORANGE)],
    ["Model FY35E revenue", "$3,280M"],
    [("Gap (implied / model)", NAVY), ("+134%", ORANGE)],
    ["", ""],
    ["Implied 2035 QC TAM share", "15.4% of $50B"],
    [("Required 10-yr CAGR", NAVY), ("~75%", ORANGE)],
], col_widths_in=[3.2, 2.7], row_h=0.27, body_size=11, first_col_bold=True)

add_footer(s, 14)

# =====================================================================
# SLIDE 15 — BULL / BASE / BEAR
# =====================================================================
s = blank()
add_header(s, "Bull / base / bear scenarios", num="15",
           subtitle="What needs to be true at each price level — for the next 12-18 months")

add_table(s, 0.3, 1.1, [
    ["", "Bull ($110+)", "Base ($60-75)", "Bear ($25-45)"],
    [("FY30E revenue", NAVY), "$1.0B+", "$725M", "$300-450M"],
    [("FY35E revenue", NAVY), "$5.5B+", "$3.3B", "$1.5-2.0B"],
    [("2035 QC market share", NAVY), "~11%", "~6.6%", "~3-4%"],
    [("FCF crossover year", NAVY), ("2028", GREEN), "2029", ("2031+", RED)],
    [("Helios ramp", NAVY), ("Beats Hololens-style launch", GREEN), "Hits launch-customer guidance", ("Slips by 6-12 months", RED)],
    [("Sol (2027) delivery", NAVY), ("On time, 192 fully connected qubits", GREEN), "On time, mostly", ("Slips to 2028", RED)],
    [("Apollo (2029) FTQC", NAVY), ("On time, demonstrated fault-tolerance", GREEN), "Slips ~6 months", ("Slips >12 months or partial", RED)],
    [("Competitor displacement", NAVY), "IONQ Tempo underperforms", "No change", ("IONQ scales faster on AWS", RED)],
    [("Macro for pre-profit tech", NAVY), "Rate cuts; ERP compresses", "Stable rates", ("Recession / risk-off", RED)],
    [("HON sells down", NAVY), "Orderly secondary, no overhang", "Held flat through 2027", ("Forced sale, overhang", RED)],
    [("Implied EV / FY25", NAVY), "~800x", "~490x (current)", "~200x"],
], col_widths_in=[3.3, 3.4, 3.0, 3.6], row_h=0.32, first_col_bold=True, body_size=10, header_size=11)

add_footer(s, 15)

# =====================================================================
# SLIDE 16 — DECISION + HONEYWELL SOTP
# =====================================================================
s = blank()
add_header(s, "Decision & Honeywell SOTP angle", num="16",
           subtitle="QNT: HOLD with downside skew. HON: marginal positive, ~5% NAV contribution")

# Left: QNT decision matrix
add_rect(s, 0.3, 1.0, 6.3, 6.0, fill=WHITE, line=NAVY)
add_text(s, 0.5, 1.1, 6, 0.4, "QNT (NASDAQ) — decision matrix", size=14, bold=True)
add_table(s, 0.5, 1.6, [
    ["Price", "Read", "Action"],
    [("$25-35", NAVY), "Trades to intrinsic DCF", ("BUY", GREEN)],
    [("$35-50", NAVY), "Below peer median; IONQ-comparable", ("BUY ↓", GREEN)],
    [("$50-65", NAVY), "Peer-median footing", ("BUY ↓", GREEN)],
    [("$65-80", NAVY), ("CURRENT — full optionality priced", ORANGE), ("HOLD", ORANGE)],
    [("$80-100", NAVY), "Premium vs peers + DCF", ("TRIM", RED)],
    [(">$100", NAVY), "Bull case priced in", ("SELL", RED)],
], col_widths_in=[1.5, 3.3, 1.5], row_h=0.32, body_size=11, first_col_bold=True)

add_text(s, 0.5, 4.0, 6, 0.4, "Key catalysts (next 12 months)", size=13, bold=True)
add_bullets(s, 0.5, 4.4, 6, 2.5, [
    ("Q2'26 print (Aug 2026) — first as public co.", True),
    "  RIKEN cliff visibility, Helios ramp narrative",
    ("Sol roadmap update (timing TBC)", True),
    ("Honeywell secondary offering (likely 6-12 mo)", True),
    "  Lockup expiry creates supply overhang",
    ("CHIPS Act $100M LOI conversion", True),
    ("Hyperscaler distribution deal expansions", True),
], size=10)

# Right: HON SOTP angle
add_rect(s, 6.8, 1.0, 6.2, 6.0, fill=WHITE, line=NAVY)
add_text(s, 7.0, 1.1, 6, 0.4, "HON (NYSE) — sum-of-parts angle", size=14, bold=True)
add_table(s, 7.0, 1.6, [
    ["Segment", "2025 rev", "Multiple", "EV"],
    ["Aerospace Tech", "$15.8B", "18× EBITDA", "~$80B"],
    ["Industrial Auto", "$9.2B", "15× EBITDA", "~$35B"],
    ["Building Auto", "$6.1B", "16× EBITDA", "~$22B"],
    ["Energy & Sustain.", "$7.4B", "12× EBITDA", "~$18B"],
    [("QNT stake (~48%)", NAVY), "—", "—", ("~$7.5B", GREEN)],
    [("Net cash / other", NAVY), "—", "—", "$(15B)"],
    [("Total EV", NAVY), "—", "—", "~$147B"],
    [("Implied HON share", NAVY), "—", "—", "~$226"],
    [("Current HON price", NAVY), "—", "—", "~$213 [E]"],
], col_widths_in=[1.9, 1.3, 1.5, 1.5], row_h=0.30, body_size=10, header_size=11, first_col_bold=True)

add_text(s, 7.0, 5.4, 6, 0.4, "HON read", size=13, bold=True)
add_bullets(s, 7.0, 5.75, 6, 1.4, [
    ("QNT IPO adds ~$5/share to HON SOTP (~2.4%)", True),
    "Real upside if HON markets a clean SOTP story",
    "Industrial-conglomerate discount may compress",
    ("Verdict for HON shareholders: modest positive", True),
], size=10)

add_footer(s, 16)

# =====================================================================
prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Slides: {len(prs.slides)}")
