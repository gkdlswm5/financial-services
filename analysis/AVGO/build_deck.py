"""
Build AVGO-Competitive-Analysis.pptx — 20-slide investment-decision deck for
Broadcom (AVGO), per the `competitive-analysis` skill. Mirrors analysis/MP and
analysis/WMT, extended with two valuation slides (comps football field + DCF
sensitivity) since this run includes the full model + DCF.

Run: python3 build_deck.py   (after build_model.py / build_comps.py)
Output: AVGO-Competitive-Analysis.pptx + charts/*.png
"""
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
CHARTS = HERE / "charts"; CHARTS.mkdir(exist_ok=True)
OUT = HERE / "AVGO-Competitive-Analysis.pptx"

NAVY = RGBColor(0x17, 0x36, 0x5D)
RED_AC = RGBColor(0xCC, 0x00, 0x00)     # Broadcom red accent
CHARCOAL = RGBColor(0x40, 0x40, 0x40)
GREY = RGBColor(0xBF, 0xBF, 0xBF)
LGREY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
FONT = "Times New Roman"

plt.rcParams.update({
    "font.family": "serif", "font.serif": ["DejaVu Serif"],
    "axes.edgecolor": "#404040", "axes.labelcolor": "#404040",
    "xtick.color": "#404040", "ytick.color": "#404040",
    "axes.grid": True, "grid.color": "#E0E0E0", "grid.linewidth": 0.5,
})


def chart_ai_ramp():
    fig, ax = plt.subplots(figsize=(9, 4.6), dpi=150)
    yrs = ["FY24A", "FY25A", "FY26E", "FY27E", "FY28E", "FY29E", "FY30E"]
    base = [12.2, 20, 48, 76, 98, 116, 131]
    colors = ["#888", "#888", "#CC0000", "#CC0000", "#CC0000", "#CC0000", "#CC0000"]
    bars = ax.bar(yrs, base, color=colors, edgecolor="#202020", linewidth=0.5)
    for b, v in zip(bars, base):
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+2, f"${v:.0f}B", ha="center", fontsize=9, fontweight="bold")
    # mgmt line-of-sight marker for FY27 ($100B+)
    ax.scatter([3], [100], marker="_", s=2000, color="#17365D", linewidths=3)
    ax.annotate("Mgmt 'line of sight\\nto $100B+ AI in FY27'", (3, 100), xytext=(3.3, 112),
                fontsize=8.5, color="#17365D", fontweight="bold")
    ax.set_ylabel("AI semiconductor revenue ($B)")
    ax.set_title("Broadcom AI revenue ramp — base case (red) haircuts mgmt FY27 guidance ~24%",
                 fontsize=12.5, color="#17365D", pad=10)
    ax.set_ylim(0, 145)
    fig.tight_layout(); out = CHARTS / "ai_ramp.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


def chart_2x2():
    fig, ax = plt.subplots(figsize=(9.5, 6), dpi=150)
    # (rev_$B TTM, growth %, name, color, size~mktcap)
    data = [
        (68.3, 29, "Broadcom",  "#CC0000", 1200),
        (253.5, 85, "NVIDIA",   "#17365D", 1500),
        (37.5, 38, "AMD",       "#7F8C8D", 700),
        (8.6, 28, "Marvell",    "#7F8C8D", 400),
        (44.5, -3, "Qualcomm",  "#95A5A6", 600),
        (18.4, 19, "Texas Inst.","#95A5A6", 650),
        (64.1, 22, "Oracle",    "#2E7D32", 800),
    ]
    for rev, gr, name, col, size in data:
        ax.scatter(rev, gr, s=size, color=col, alpha=0.7, edgecolors="#202020", linewidths=1.1)
        ax.annotate(name, (rev, gr), xytext=(8, 7), textcoords="offset points", fontsize=10, fontweight="bold")
    ax.set_xscale("log")
    ax.set_xlabel("Revenue scale, $B TTM (log)")
    ax.set_ylabel("Revenue growth, most recent qtr YoY (%)")
    ax.axhline(30, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.axvline(60, color="#BFBFBF", linewidth=0.7, linestyle=":")
    ax.set_xlim(6, 400); ax.set_ylim(-15, 100)
    ax.set_title("Positioning: scale x growth  (bubble = market cap)", fontsize=12.5, color="#17365D", pad=12)
    ax.text(7, 92, "Hyper-growth", fontsize=9, color="#7F8C8D")
    ax.text(80, 92, "Scaled + fast", fontsize=9, color="#7F8C8D")
    ax.text(7, -12, "Niche / mature", fontsize=9, color="#7F8C8D")
    ax.text(80, -12, "Scaled / mature", fontsize=9, color="#7F8C8D")
    fig.tight_layout(); out = CHARTS / "positioning.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


def chart_football():
    fig, ax = plt.subplots(figsize=(9.2, 4.6), dpi=150)
    # methods: (label, low, high)
    rows = [
        ("DCF — base (WACC/g grid $200-311)", 200, 311),
        ("DCF — bull (mgmt FY27 AI fully delivered)", 380, 540),
        ("Comps — EV/EBITDA (semis median)", 320, 490),
        ("Comps — forward P/E (semis median)", 350, 500),
    ]
    ypos = range(len(rows))
    for y, (lab, lo, hi) in zip(ypos, rows):
        ax.barh(y, hi-lo, left=lo, height=0.5, color="#17365D", alpha=0.75, edgecolor="#202020")
        ax.text(lo-6, y, f"${lo:.0f}", va="center", ha="right", fontsize=9)
        ax.text(hi+6, y, f"${hi:.0f}", va="center", ha="left", fontsize=9)
    ax.set_yticks(list(ypos)); ax.set_yticklabels([r[0] for r in rows], fontsize=9.5)
    ax.axvline(481.62, color="#CC0000", linewidth=2)
    ax.text(481.62, len(rows)-0.3, "  Current $482", color="#CC0000", fontsize=10, fontweight="bold")
    ax.axvline(264, color="#2E7D32", linewidth=1.5, linestyle="--")
    ax.text(264, -0.9, "DCF base $264", color="#2E7D32", fontsize=9, ha="center")
    ax.set_xlim(120, 600); ax.set_xlabel("Implied value per share ($)")
    ax.set_title("Valuation football field — market price sits at/above the bull case",
                 fontsize=12.5, color="#17365D", pad=10)
    ax.invert_yaxis()
    fig.tight_layout(); out = CHARTS / "football.png"; fig.savefig(out, bbox_inches="tight"); plt.close(fig)
    return out


AI_PNG = chart_ai_ramp(); MATRIX_PNG = chart_2x2(); FB_PNG = chart_football()

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = Emu(45720); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    runs = [text] if isinstance(text, str) else text
    for i, line in enumerate(runs):
        if i == 0:
            r = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align; r = p2.add_run()
        r.text = line; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def slide_title(slide, title, *, size=24):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.10))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED_AC; bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95), title, size=size, bold=True, color=NAVY)


def slide_footer(slide, txt):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3), txt, size=9, italic=True, color=GREY)


def add_table(slide, left, top, width, height, data, *, header=True, col_widths=None, first_col_bold=True, font_size=12):
    rows = len(data); cols = len(data[0])
    shp = slide.shapes.add_table(rows, cols, left, top, width, height); tbl = shp.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT; run.font.size = Pt(font_size)
                    if header and ri == 0:
                        run.font.bold = True; run.font.color.rgb = WHITE
                    elif first_col_bold and ci == 0:
                        run.font.bold = True; run.font.color.rgb = NAVY
                    else:
                        run.font.color.rgb = CHARCOAL
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Emu(45720); cell.margin_top = cell.margin_bottom = Emu(27432)
    return shp


def dots(n):
    return {3: "●●● ", 2: "●●○ ", 1: "●○○ "}.get(n, "○○○ ")


# ============ SLIDE 1 — COVER ============
s = prs.slides.add_slide(BLANK)
banner = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(4.5))
banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), prs.slide_width, Inches(0.12))
stripe.fill.solid(); stripe.fill.fore_color.rgb = RED_AC; stripe.line.fill.background()
add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5), "NASDAQ: AVGO — Investment Decision Frame", size=18, italic=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.4), "Broadcom Inc.", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.9), "Competitive Landscape, Comps & DCF Valuation", size=26, color=RED_AC)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5), "As of June 3, 2026 (post-Q2 FY26 print)", size=14, italic=True, color=CHARCOAL)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.6),
         ["Peer set: NVIDIA · AMD · Marvell · Qualcomm · Texas Instruments  +  Oracle (software anchor)",
          "Built with deep-research + competitive-analysis + comps-analysis + 3-statement-model + dcf-model skills",
          "Companion files: AVGO-Comps-Analysis.xlsx · AVGO-Model.xlsx (3-statement + DCF)"],
         size=13, color=CHARCOAL)

# ============ SLIDE 2 — THESIS ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Q2 print confirms the franchise — but at ~$482, even after the move, the price still embeds bull AI")
bullets = [
    ("1. The franchise", "AVGO is the #1 custom AI-silicon (XPU/ASIC) design house (~70% share) + leading AI networking (Tomahawk/Jericho). Q2 FY26: revenue $22.2B (+48%), AI revenue $10.8B (+143%), adj. EBITDA $15.2B (69%), FCF $10.3B (46% margin). Q3 guide $29.4B (+84%). $73B AI backlog; 6 committed hyperscaler customers."),
    ("2. The FCF engine", "VMware turned Infrastructure Software into a ~$27B/yr, ~78%-operating-margin annuity. Consolidated TTM FCF ~$33B+ (~44% margin) after Q2's $10.3B print. Net debt down to ~$49B from a ~$74B post-VMware peak; dividend $0.65/qtr (+10%)."),
    ("3. The valuation problem", "Hock Tan REITERATED — did NOT raise — the FY27 'AI $100B+' guide despite the Q2 beat. Our base-case DCF (now haircut ~15% to $85B FY27 AI) yields ~$264/share. The ENTIRE WACC×g grid ($200–$311) sits below the ~$482 market price; bull corner ~$395."),
    ("4. The read", "The Q2 beat closed some of the gap (stock +4.7% on the print) — but the structural finding holds: to justify $482 you still need the full mgmt FY27 AI delivered AND a low discount rate. High-quality compounder; no margin of safety. Long on pullbacks; pair vs MRVL; trim into strength."),
]
y = Inches(1.35)
for h, b in bullets:
    add_text(s, Inches(0.6), y, Inches(2.4), Inches(1.2), h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(3.1), y, Inches(9.6), Inches(1.2), b, size=11.5, color=CHARCOAL)
    y += Inches(1.36)
slide_footer(s, "Sources: AVGO Q2 FY26 IR (June 3, 2026); AVGO-Model.xlsx (base-case DCF); The Information (May 2026, OpenAI financing). Research, not investment advice.")

# ============ SLIDE 3 — THE QUESTION ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "The question this analysis answers")
add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(1.2),
         '"Broadcom Q2 FY26 (June 3) printed +48% revenue, AI +143%, with Q3 guided +84% — but the stock is at ~$482 (~$2.3T, ~29x fwd). '
         'How much of the AI bull case is already in the price post-print, and what has to be true to justify owning it here?"',
         size=20, italic=True, color=NAVY)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4), "We test that across five lenses:", size=14, bold=True, color=CHARCOAL)
for i, ln in enumerate([
    "(1) Market — is the hyperscaler AI-capex cycle still supportive?",
    "(2) Competitive — how durable is AVGO's custom-ASIC + networking position vs NVDA, AMD, MRVL?",
    "(3) Financial (comps) — how does AVGO screen vs the semis cohort + ORCL software anchor?",
    "(4) Intrinsic (DCF) — what is AVGO worth on a base case, and what does $460 imply?",
    "(5) Decision — bull/base/bear and the cleanest way to express a view.",
]):
    add_text(s, Inches(0.9), Inches(4.2)+i*Inches(0.46), Inches(12), Inches(0.4), ln, size=13, color=CHARCOAL)
slide_footer(s, "Framework per the `competitive-analysis` skill; valuation per `3-statement-model` + `dcf-model`.")

# ============ SLIDE 4 — MARKET ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Hyperscaler AI capex ~$630–725B in 2026, $1T+ by 2027 — mix shifting to custom silicon")
data = [
    ["Driver", "Data point", "Source / note"],
    ["Big-4 capex 2025A", "~$388–410B (record)", "MSFT, GOOGL, AMZN, META"],
    ["Big-4 capex 2026E", "~$630–725B (+62–77%)", "CNBC / Tom's Hardware (range)"],
    ["Total compute capex 2026E", "~$1.04T (first trillion-$ yr)", "incl. Oracle, neoclouds, sovereigns"],
    ["AI accelerator TAM", "~$200B → $500B+ by late-decade", "GPU + ASIC both growing"],
    ["Custom ASIC growth 2026E", "~+45% units vs ~+16% GPU", "ASIC units to surpass GPU by 2027 (Goldman)"],
    ["AVGO + MRVL share of custom ASIC", "~95% of co-design market", "Tom's Hardware / Dell'Oro"],
]
add_table(s, Inches(0.5), Inches(1.35), Inches(8.3), Inches(4.6), data, col_widths=[Inches(3.0), Inches(2.8), Inches(2.5)], font_size=11)
add_text(s, Inches(9.0), Inches(1.4), Inches(4.0), Inches(0.5), "The catch", size=13, bold=True, color=RED_AC)
add_text(s, Inches(9.0), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• Capex growing ~46% faster than revenue — wider than the 2001 telecom-bust divergence (Allianz).",
          "",
          "• ~$400B/yr implied AI depreciation > Big-4's combined profits.",
          "",
          "• Power is the binding constraint: Gartner sees ~40% of AI DCs power-limited by 2027.",
          "",
          "• Bull case intact near-term; digestion risk into 2027–28 is the key macro swing factor for AVGO."],
         size=11, color=CHARCOAL)
slide_footer(s, "Sources: CNBC, Tom's Hardware, Goldman Sachs, Dell'Oro, Allianz, Gartner (2026). Figures are third-party estimates; ranges where sources diverge.")

# ============ SLIDE 5 — SEGMENT SPLIT ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Two businesses under one roof: a ~60% AI-led semis engine and a ~40% software FCF annuity")
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6), [
    ["Segment (Q1 FY26)", "Revenue", "Mix", "Growth", "Op margin", "Role in thesis"],
    ["Semiconductor Solutions", "$12.5B", "~65%", "+52% YoY", "~high", "AI ASIC + networking = the growth"],
    ["  of which: AI semiconductors", "$8.4B", "~44% of total", "+106% YoY", "—", "The engine; Q2 guide $10.7B (+140%)"],
    ["  of which: non-AI semis", "$4.1B", "~21% of total", "~flat", "—", "Cyclical; Apple Wi-Fi insourcing a headwind"],
    ["Infrastructure Software", "$6.8B", "~35%", "+1% YoY", "78%", "VMware annuity = the FCF / margin"],
    ["Total", "$19.3B", "100%", "+29% YoY", "66% (non-GAAP)", "Adj. EBITDA $13.1B (68%)"],
], col_widths=[Inches(3.2), Inches(1.5), Inches(1.6), Inches(1.6), Inches(1.6), Inches(2.8)], font_size=11)
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3),
         ["• The semis engine is what re-rates the stock — AI revenue +106% with a $73B backlog and 'line of sight to $100B+ AI in FY27'.",
          "• The software annuity (78% operating margin, ~$27B/yr) is what de-risks it — recurring, high-margin cash funding debt paydown + buybacks.",
          "• This bifurcation is why the comps need a software anchor (ORCL) alongside the semis cohort — neither lens alone captures AVGO.",
          "• Caveat: AVGO reports ONE Semiconductor segment; the AI / non-AI split is per management commentary, partly estimated."],
         size=12.5, color=CHARCOAL)
slide_footer(s, "Source: Broadcom Q1 FY26 earnings release (Mar 2026). Op margins are segment-level; consolidated 66% is non-GAAP.")

# ============ SLIDE 6 — AVGO PROFILE ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Broadcom profile — AI inflection with a fortress FCF profile and de-levering balance sheet")
add_table(s, Inches(0.5), Inches(1.35), Inches(6.2), Inches(5.4), [
    ["Metric", "Value"],
    ["Q1 FY26 revenue / Q2 actual", "$19.3B (+29%) / $22.2B (+48%)"],
    ["Q1+Q2 FY26 AI revenue", "$8.4B + $10.8B = $19.2B (+143% Q2 YoY)"],
    ["Q3 FY26 guide (total)", "$29.4B (+84%) — beat consensus $28.5B"],
    ["FY26 revenue consensus [E]", "~$103-105B (raised post-Q2)"],
    ["Q2 adj. EBITDA / margin", "$15.2B / 69%"],
    ["TTM free cash flow [E]", "~$33B+ (~44% margin); Q2 alone $10.3B"],
    ["Total debt / cash [E]", "~$65B / ~$16B"],
    ["Net debt [E]", "~$49B"],
    ["Market cap / EV [E]", "~$2.28T / ~$2.33T"],
    ["P/E TTM / forward [E]", "~75x / ~29x"],
    ["Dividend (FY26)", "$0.65/qtr (~0.54% yield, +10%)"],
], col_widths=[Inches(3.3), Inches(2.9)], font_size=11)
s.shapes.add_picture(str(AI_PNG), Inches(6.95), Inches(1.5), width=Inches(6.0))
add_text(s, Inches(6.95), Inches(5.6), Inches(6.0), Inches(1.4),
         ["• AI revenue: $12.2B (FY24) → $20B (FY25) → base-case $48B/$76B (FY26/27E).",
          "• Red bars = our base case; navy marker = mgmt's $100B+ FY27 line-of-sight we haircut ~24%."],
         size=10.5, color=CHARCOAL)
slide_footer(s, "Sources: AVGO Q1 FY26 8-K/IR; StockAnalysis.com (price, shares, debt ~June 1 2026); AVGO-Model.xlsx for projected AI revenue.")

# ============ SLIDE 7 — INDUSTRY ECONOMICS ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Merchant GPU vs custom ASIC — the structural share shift is the heart of the AVGO bull case")
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
    ["Dimension", "Merchant GPU (NVDA, AMD)", "Custom ASIC / XPU (AVGO, MRVL)"],
    ["Who buys", "All AI builders; flexible, CUDA ecosystem", "Hyperscalers with scale to amortize design"],
    ["2026E unit growth", "~+16%", "~+45% (nearly 3x)"],
    ["Share of accelerators today", "~70% (NVDA) eroding at margin", "~15–25%, rising toward ~50% by 2027 (Goldman)"],
    ["Economics", "High merchant margin; NVDA ~75% GM", "Lower unit margin, but sticky multi-gen design lock-in"],
    ["AVGO position", "n/a (not merchant)", "~70% of custom AI accelerator DESIGN market"],
], col_widths=[Inches(2.6), Inches(4.6), Inches(5.1)], font_size=11)
add_text(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0),
         ["• The bull case: as hyperscalers internalize inference on cost-optimized custom silicon, the ASIC pool grows ~3x faster than GPUs — and AVGO + MRVL own ~95% of that design market.",
          "• The bear case: AVGO's AI revenue leans on a handful of programs (Google, Meta, ByteDance, OpenAI, Anthropic, Apple). Google adding MediaTek/Marvell as second sources dilutes per-program content; an OpenAI/program slip is a large hole.",
          "• Networking is the under-appreciated leg: Tomahawk 6 (102.4 Tb/s) + Jericho4 position open-Ethernet against NVIDIA NVLink/Spectrum/InfiniBand — Ethernet has now passed InfiniBand in AI back-end share (Dell'Oro)."],
         size=12, color=CHARCOAL)
slide_footer(s, "Sources: Goldman Sachs, Dell'Oro, Tom's Hardware, The Next Platform (2026). ASIC/GPU mix figures are estimates.")

# ============ SLIDE 8 — MOATS ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "AVGO moat — custom-silicon switching costs + IP + a $73B backlog + VMware lock-in")
add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.3), [
    ["Layer", "Detail / KPI"],
    ["Custom-silicon design lock-in", "Multi-year co-design with hyperscalers; switching cost is a full chip re-spin. ~70% of custom ASIC design market."],
    ["AI backlog visibility", "~$73B AI backlog; 'line of sight to $100B+ AI revenue in FY27' across 6 committed customers"],
    ["Confirmed anchor customers", "Google (TPU/Ironwood), Meta (MTIA), ByteDance; + Anthropic (~$21B), OpenAI (~10GW), Apple ('Baltra')"],
    ["AI networking franchise", "Tomahawk 6 (102.4 Tb/s) + Tomahawk Ultra + Jericho4 — open Ethernet vs NVIDIA NVLink/InfiniBand"],
    ["IP & SerDes leadership", "Industry-leading 200G+ SerDes; deep analog/mixed-signal + optical (CPO) IP portfolio"],
    ["VMware software annuity", "~$27B/yr, ~78% operating margin, recurring; deepens enterprise lock-in and funds capital return"],
    ["Capital-return discipline", "Hock Tan playbook: acquire, raise margins, pay down debt, grow dividend — ~$28.9B TTM FCF"],
], col_widths=[Inches(3.4), Inches(8.9)], font_size=11)
slide_footer(s, "Sources: Broadcom IR & earnings calls; CNBC (Anthropic, Dec 2025); company product releases (Tomahawk/Jericho). Some customer details reported/rumored.")

# ============ SLIDE 9 — PEER SET ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer set — semis cohort for the ~60%, an infrastructure-software anchor for the ~40%")
groups = [
    ("Closest competitor — custom silicon", ["Marvell (MRVL) — #2 custom-ASIC house; Amazon Trainium, Microsoft Maia", "The cleanest read-through to AVGO's AI ASIC franchise"], RED_AC),
    ("Merchant-GPU benchmarks", ["NVIDIA (NVDA) — AI accelerator leader; the multiple everyone anchors to", "AMD (AMD) — #2 data-center GPU"], NAVY),
    ("Mature-semi capital-return refs", ["Qualcomm (QCOM) — connectivity; value multiple, low AI-accelerator exposure", "Texas Instruments (TXN) — analog; dividend/quality reference"], GREY),
    ("Software anchor (VMware ~40%)", ["Oracle (ORCL) — large-cap infra software, now AI-capex-driven", "Frames the software segment (memo; excluded from semis stats)"], GREEN),
]
y = Inches(1.4)
for t, items, col in groups:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(5.0), Inches(0.5), t, size=14, bold=True, color=col)
    for i, it in enumerate(items):
        add_text(s, Inches(0.85), y+Inches(0.45)+i*Inches(0.32), Inches(11.5), Inches(0.3), "• "+it, size=12, color=CHARCOAL)
    y += Inches(1.35)
slide_footer(s, "Peer framing per `comps-analysis`. Statistics computed over the 5 semis peers; ORCL shown as a software-anchor memo (you don't median a software multiple into a semis set).")

# ============ SLIDE 10 — POSITIONING ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Positioning — NVDA is scaled-and-fast; AVGO is the scaled compounder; MRVL the high-beta proxy")
s.shapes.add_picture(str(MATRIX_PNG), Inches(0.4), Inches(1.3), width=Inches(8.8))
add_text(s, Inches(9.4), Inches(1.4), Inches(3.7), Inches(0.5), "What the matrix says", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.4), Inches(1.95), Inches(3.7), Inches(5.0),
         ["• NVIDIA: uniquely scaled AND +85% — the outlier; vast profitability justifies a 'low' P/E.",
          "",
          "• AVGO: $68B revenue at +29%, with the highest-margin/FCF profile — a scaled compounder, not a hyper-grower.",
          "",
          "• Marvell: small but +28% and re-rating — the high-beta way to play the same ASIC theme.",
          "",
          "• QCOM/TXN: scaled but mature/cyclical, limited AI-accelerator exposure.",
          "",
          "• Oracle: software anchor, +22% on AI-infra demand."],
         size=11, color=CHARCOAL)
slide_footer(s, "Bubble size = market cap. Growth = most recent reported quarter YoY. Source: companion AVGO-Comps-Analysis.xlsx.")

# ============ SLIDE 11 — NVDA DEEP DIVE ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dive: NVIDIA — the benchmark; scaled, hyper-profitable, and (optically) not expensive")
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.0), [
    ["Metric", "Value"],
    ["Market cap", "~$5.46T"],
    ["TTM revenue", "~$253.5B"],
    ["MRQ revenue growth", "+85% YoY (Q1 FY27)"],
    ["Gross margin", "74.9%"],
    ["TTM EBITDA / margin", "~$174.4B / ~69%"],
    ["TTM net income", "~$159.6B"],
    ["Net cash", "net cash (debt $8.5B < cash)"],
    ["P/E TTM / fwd", "~34x / ~26x"],
    ["EV/EBITDA / EV/Rev", "~31x / ~21x"],
], col_widths=[Inches(2.5), Inches(3.7)], font_size=11)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.0), [
    ["Category", "Assessment"],
    ["Business", "Merchant AI GPU + networking (NVLink, Spectrum-X, InfiniBand); CUDA software moat."],
    ["Strengths", "• ~70%+ accelerator share\n• Enormous profitability funds R&D lead\n• Full-stack (chip+system+software)"],
    ["Weaknesses", "• Custom ASIC share shift is the structural threat\n• Customer-concentration in a few hyperscalers\n• Geopolitical/export risk (China)"],
    ["Read-through to AVGO", "NVDA's 'low' P/E sets the ceiling the market pays for AI compute; AVGO trades RICHER on fwd P/E despite slower growth."],
], col_widths=[Inches(1.4), Inches(4.5)], font_size=10.5)
slide_footer(s, "Sources: NVIDIA Q1 FY27 newsroom; Macrotrends; GuruFocus; valueinvesting.io (~June 2026).")

# ============ SLIDE 12 — AMD + MRVL ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer deep dives: AMD (#2 GPU) and Marvell (closest custom-silicon comp)")
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.2), [
    ["AMD (AMD)", ""],
    ["Market cap", "~$839B"],
    ["TTM revenue / growth", "$37.5B / +38% YoY"],
    ["Gross margin", "~53%"],
    ["Net income (TTM)", "~$5.0B [E] (GAAP ramping)"],
    ["P/E TTM / fwd", "~150–167x / ~33x"],
    ["EV/Rev", "~22x"],
    ["Read", "#2 merchant GPU (MI-series); read on EV/Rev + forward, not GAAP TTM"],
], col_widths=[Inches(1.9), Inches(4.3)], font_size=10.5)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.2), [
    ["Marvell (MRVL)", ""],
    ["Market cap", "~$192B (+51% on NVDA endorsement)"],
    ["TTM revenue / growth", "~$8.6B [E] / +28% YoY"],
    ["Gross margin", "52.1% GAAP / 58.9% non-GAAP"],
    ["Net income (TTM)", "~$0.3B GAAP [E] (amortization-depressed)"],
    ["P/E TTM / fwd", "NM-high / ~48x"],
    ["Customers", "Amazon Trainium/Inferentia, Microsoft Maia"],
    ["Read", "Closest AVGO ASIC comp; smaller, higher-beta, single net-debt name"],
], col_widths=[Inches(1.9), Inches(4.0)], font_size=10.5)
slide_footer(s, "Sources: AMD Q1 2026 8-K; Marvell Q1 FY27 IR; StockAnalysis.com; GuruFocus (~June 2026).")

# ============ SLIDE 13 — QCOM + TXN ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Peer context: Qualcomm & Texas Instruments — mature-semi valuation anchors")
add_table(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(4.8), [
    ["Qualcomm (QCOM)", ""],
    ["Market cap / TTM rev", "~$215B / ~$44.5B"],
    ["MRQ growth", "~-3% YoY (handset soft)"],
    ["TTM EBITDA / margin", "~$14.0B / ~31%"],
    ["P/E TTM / fwd", "~20x / ~14x (cheapest fwd)"],
    ["Read", "Value/mature; limited AI-accelerator exposure; sets the cheap end"],
], col_widths=[Inches(2.2), Inches(4.0)], font_size=10.5)
add_table(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(4.8), [
    ["Texas Instruments (TXN)", ""],
    ["Market cap / TTM rev", "~$287B / ~$18.4B"],
    ["MRQ growth", "+19% YoY (industrial + DC)"],
    ["TTM EBITDA / margin", "~$8.8B / ~48%"],
    ["P/E TTM / fwd", "~50x / ~30x (trough earnings)"],
    ["Read", "Quality analog/dividend ref; multiples reflect capex-cycle trough"],
], col_widths=[Inches(2.4), Inches(3.5)], font_size=10.5)
add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
         "Both frame the mature-semi multiple band (fwd ~14–30x) — useful as the floor under AVGO's non-AI semis, but neither carries AVGO's AI growth.",
         size=12, italic=True, color=CHARCOAL)
slide_footer(s, "Sources: Qualcomm Q2 FY26; TI Q1 2026; Macrotrends; GuruFocus (~May–June 2026). QCOM TTM NI excludes one-time tax benefit.")

# ============ SLIDE 14 — VMWARE vs ORCL ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "The software side: VMware (within AVGO) vs Oracle — the ~40% the semis comps miss")
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
    ["Dimension", "AVGO Infrastructure Software (VMware)", "Oracle (ORCL) — anchor"],
    ["Revenue", "~$6.8B/qtr (~$27B/yr)", "~$64.1B TTM (+22% YoY)"],
    ["Profitability", "~78% segment operating margin", "~48% EBITDA margin"],
    ["Growth profile", "~+1% YoY (re-based post-acquisition; margin story)", "+22% (OCI +84%) — AI-capex growth story"],
    ["Balance sheet", "Funded within AVGO; group de-levering", "~$96B net debt, rising (+$30B raised Feb 2026)"],
    ["Multiple lens", "Annuity / FCF — high-margin, low-growth", "EV/Rev ~12x, EV/EBITDA ~18x (growth premium)"],
], col_widths=[Inches(2.4), Inches(5.4), Inches(4.5)], font_size=11)
add_text(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0),
         ["• ORCL is the best LARGE-CAP infra-software anchor, but it's an imperfect proxy: it's now an AI-datacenter capex story (high growth, ballooning debt), whereas VMware is a low-growth, ~78%-margin annuity.",
          "• Net: VMware likely warrants a HIGHER-margin / LOWER-growth multiple than ORCL's blended profile — value it as an FCF annuity, not on ORCL's growth multiple.",
          "• Practical takeaway for the comps: keep ORCL as a memo, not in the semis median — and treat AVGO's software EBITDA as the low-volatility ballast under the AI growth option."],
         size=12, color=CHARCOAL)
slide_footer(s, "Sources: Broadcom Q1 FY26 (software segment); Oracle Q3 FY26 (Mar 2026); Futurum; Macrotrends.")

# ============ SLIDE 15 — SCOREBOARD ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Comparative scoreboard — AVGO leads on margin/FCF; NVDA on growth; AVGO richest on fwd P/E")
add_table(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.4), [
    ["Dimension", "AVGO", "NVDA", "AMD", "MRVL", "QCOM", "TXN"],
    ["Revenue scale (TTM)", dots(2)+"$68B", dots(3)+"$254B", dots(1)+"$37B", dots(1)+"$9B", dots(2)+"$44B", dots(1)+"$18B"],
    ["Revenue growth (MRQ)", dots(2)+"+29%", dots(3)+"+85%", dots(2)+"+38%", dots(2)+"+28%", dots(1)+"-3%", dots(1)+"+19%"],
    ["AI exposure", dots(3)+"ASIC+net", dots(3)+"GPU leader", dots(2)+"#2 GPU", dots(3)+"ASIC #2", dots(1)+"edge", dots(1)+"content"],
    ["Margin / FCF quality", dots(3)+"68% EBITDA", dots(3)+"69%", dots(1)+"22%", dots(2)+"~31%[E]", dots(2)+"31%", dots(3)+"48%"],
    ["Balance sheet", dots(2)+"$52B net debt", dots(3)+"net cash", dots(3)+"net cash", dots(2)+"$1.4B nd", dots(2)+"$5.5B nd", dots(2)+"$8B nd"],
    ["Valuation (fwd P/E)", dots(1)+"~31x", dots(2)+"~26x", dots(1)+"~33x", dots(1)+"~48x", dots(3)+"~14x", dots(2)+"~30x"],
    ["Customer concentration", dots(1)+"few hyperscalers", dots(2)+"broad", dots(2)+"broad", dots(1)+"few", dots(3)+"diversified", dots(3)+"diversified"],
], col_widths=[Inches(2.5), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.6), Inches(1.6), Inches(1.6)], font_size=10)
slide_footer(s, "● = stronger / ○ = weaker on each axis (higher fwd P/E = weaker value score). Calibrated to AVGO-Comps-Analysis.xlsx.")

# ============ SLIDE 16 — VALUATION: COMPS FOOTBALL FIELD ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Valuation I — every method puts fair value at or below the $460 market price")
s.shapes.add_picture(str(FB_PNG), Inches(0.5), Inches(1.35), width=Inches(8.4))
add_text(s, Inches(9.1), Inches(1.4), Inches(4.0), Inches(0.5), "How to read it", size=13, bold=True, color=NAVY)
add_text(s, Inches(9.1), Inches(1.95), Inches(4.0), Inches(5.0),
         ["• DCF base ($264) and the full WACC×g grid ($200–311) sit well below $482.",
          "",
          "• Comps (semis-median EV/EBITDA and fwd P/E applied to AVGO) imply ~$320–500 — AVGO trades at/above the cohort post-Q2.",
          "",
          "• Only the DCF BULL case (mgmt's FY27 AI fully delivered) reaches the $482 zone.",
          "",
          "• Conclusion: even after Q2 +48% & Q3 guide +84%, the market still prices the bull case as the base case."],
         size=11, color=CHARCOAL)
slide_footer(s, "Football field combines DCF (AVGO-Model.xlsx) and comps (AVGO-Comps-Analysis.xlsx). Ranges are illustrative of method spread.")

# ============ SLIDE 17 — VALUATION: DCF + SENSITIVITY ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Valuation II — base-case DCF $264/share; the entire WACC × growth grid is below $482")
add_table(s, Inches(0.5), Inches(1.4), Inches(5.6), Inches(4.4), [
    ["DCF bridge (base, post-Q2 refresh)", "$M / $"],
    ["PV of explicit FCF (FY26–30)", "269,084"],
    ["PV of terminal value", "1,030,519"],
    ["Enterprise value", "1,299,603"],
    ["(–) Net debt", "(49,000)"],
    ["Equity value", "1,250,603"],
    ["÷ Diluted shares (M)", "4,730"],
    ["Implied value / share", "$264.40"],
    ["Current price", "$481.62"],
    ["Upside / (downside)", "(45.1%)"],
], col_widths=[Inches(3.4), Inches(2.2)], font_size=11)
add_text(s, Inches(6.4), Inches(1.4), Inches(6.4), Inches(0.4), "Sensitivity — implied $/share (WACC × terminal g)", size=12, bold=True, color=NAVY)
add_table(s, Inches(6.4), Inches(1.85), Inches(6.5), Inches(3.0), [
    ["WACC ↓ / g →", "2.5%", "3.0%", "3.5%", "4.0%", "4.5%"],
    ["8.5%", "275", "296", "323", "355", "395"],
    ["9.0%", "252", "270", "291", "317", "348"],
    ["9.5% (base)", "232", "247", "265", "286", "311"],
    ["10.0%", "215", "228", "243", "260", "281"],
    ["10.5%", "200", "211", "224", "239", "256"],
], col_widths=[Inches(1.7), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95), Inches(0.95)], font_size=10.5)
add_text(s, Inches(6.4), Inches(5.0), Inches(6.5), Inches(1.9),
         ["• Center cell ($265) = base case. Terminal value is ~79% of EV — valuation is very sensitive to WACC/g.",
          "• Even the most generous corner (WACC 8.5% / g 4.5%) = $395 on base FCF.",
          "• To reach $482 you ALSO need FCF above base — i.e., mgmt's full $100B+ FY27 AI delivered (the bull case)."],
         size=11, color=CHARCOAL)
slide_footer(s, "Source: AVGO-Model.xlsx (validated, 0 formula errors). Base case haircuts mgmt FY27 AI guidance ~24% for execution / OpenAI-financing risk.")

# ============ SLIDE 18 — SYNTHESIS / MOATS GRID ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Strategic synthesis — Moat: design lock-in + IP + FCF · Risk: concentration + valuation")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.6), [
    ["Moat type", "AVGO", "NVDA", "MRVL", "What it means"],
    ["Switching costs", "Strong", "Strong", "Mod", "Custom-silicon re-spin + VMware lock-in are sticky"],
    ["Intangibles (IP)", "Strong", "Strong", "Mod", "SerDes/optical/analog IP; CUDA is NVDA's equivalent"],
    ["Scale economies", "Strong", "Strong", "Weak", "R&D + supply scale; AVGO + NVDA dominate spend"],
    ["Network effects", "Weak", "Strong", "Weak", "NVDA's CUDA ecosystem is the real network effect"],
    ["Cost / capital discipline", "Strong", "Mod", "Mod", "Hock Tan FCF + de-levering playbook is a differentiator"],
], col_widths=[Inches(2.4), Inches(1.3), Inches(1.3), Inches(1.3), Inches(6.0)], font_size=11)
add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4), "Net assessment", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4),
         ["Durable advantages: (i) ~70% custom-ASIC design share with multi-gen lock-in, (ii) leading open-Ethernet AI networking, (iii) a ~$27B/yr ~78%-margin VMware annuity funding capital return.",
          "Structural vulnerabilities: (i) AI revenue concentrated in a few hyperscaler programs (OpenAI financing snag is live), (ii) Google multi-sourcing dilutes content, (iii) valuation leaves no margin of safety — $460 already prices the bull AI ramp."],
         size=12, color=CHARCOAL)
slide_footer(s, "Moat framework per `competitive-analysis` Step 9. Strong / Moderate (Mod) / Weak relative to peers.")

# ============ SLIDE 19 — BULL / BASE / BEAR ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Bull / base / bear — anchored to the DCF: the market price embeds the bull case")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.8), [
    ["Scenario", "Prob.", "Key drivers", "DCF value / read"],
    ["Bull", "30%", "Mgmt's $100B+ FY27 AI delivered & sustained; OpenAI deal funded; networking share gains; WACC ~8.5%, g ~4.5%",
     "~$470–520 — justifies (and exceeds) the $460 price; this is what the market currently embeds"],
    ["Base", "45%", "AI FY27 ~$76B (mgmt guidance haircut ~24%); software annuity steady; WACC 9.5%, g 3.5%",
     "~$264 — high-quality compounder, but ~45% below the current price"],
    ["Bear", "25%", "AI digestion into 2027–28; OpenAI program stalls; Google multi-sourcing erodes content; WACC ~10.5%, g ~2.5%",
     "~$185–210 — multiple compresses as growth normalizes"],
], col_widths=[Inches(1.3), Inches(0.9), Inches(5.6), Inches(4.5)], font_size=11)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4), "Signposts that move the weighting", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.9),
         "1) OpenAI $18B financing structure resolved or not   2) Q2/Q3 FY26 AI revenue vs the $10.7B→ramp guide   3) hyperscaler 2027 capex guidance (acceleration vs digestion)   4) any Google/Meta program re-sourcing",
         size=12, color=CHARCOAL)
slide_footer(s, "Probabilities are author's judgment, not consensus. DCF values from AVGO-Model.xlsx sensitivity + bull-FCF overlay.")

# ============ SLIDE 20 — DECISION + SOURCES ============
s = prs.slides.add_slide(BLANK)
slide_title(s, "Decision frame & sources")
add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4), "Four ways to express a view", size=14, bold=True, color=NAVY)
opts = [
    ("Own quality, but wait for a pullback", "Best-in-class franchise; no margin of safety at $460. Accumulate toward the base-case zone (high-$200s–low-$300s) or on AI-digestion scares.", NAVY),
    ("Pair trade: long AVGO / short a richer AI name", "Express the quality view while hedging the AI-multiple: long AVGO vs short a higher-beta/expensive AI name on relative valuation.", RED_AC),
    ("Play the theme via MRVL (higher beta)", "Same custom-ASIC thesis with more torque (and more risk); smaller, re-rating, single net-debt name.", CHARCOAL),
    ("Trim / avoid here", "If you believe 2027–28 AI digestion + OpenAI-financing risk, the price already discounts the bull case — trim into strength.", GREY),
]
y = Inches(1.7)
for t, b, col in opts:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.18), Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(0.85), y, Inches(11.8), Inches(0.35), t, size=13, bold=True, color=col)
    add_text(s, Inches(0.85), y+Inches(0.38), Inches(11.8), Inches(0.6), b, size=11, color=CHARCOAL)
    y += Inches(1.02)
add_text(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.35), "Sources & caveats", size=12, bold=True, color=NAVY)
add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.9),
         ["AVGO Q1 FY26 & FY23–25 SEC filings/IR; NVDA/AMD/MRVL/QCOM/TXN/ORCL latest 10-Q/IR; StockAnalysis/GuruFocus/Macrotrends (~June 2026); CNBC, The Information, Goldman, Dell'Oro, Allianz for AI-capex & customer detail.",
          "MCP terminal connectors (CapIQ/FactSet) NOT configured — public sources used, [E] flags in workbooks. AI/non-AI segment split partly estimated. Research & decision-framing only — NOT investment advice; cross-check a primary terminal before acting."],
         size=10, italic=True, color=CHARCOAL)
slide_footer(s, "Built per competitive-analysis + comps-analysis + 3-statement-model + dcf-model skills. Companions: AVGO-Comps-Analysis.xlsx, AVGO-Model.xlsx.")

prs.save(OUT)
print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
